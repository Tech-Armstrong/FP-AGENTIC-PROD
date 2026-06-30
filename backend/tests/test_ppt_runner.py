"""Tests for PPT generation helpers and workflow_state serialization."""

import sys
from datetime import date
from pathlib import Path
from unittest.mock import MagicMock, patch

import pytest

REPO_ROOT = Path(__file__).resolve().parent.parent.parent
sys.path.insert(0, str(REPO_ROOT / "backend"))
sys.path.insert(0, str(REPO_ROOT))

from financial_plan_runner import json_safe_workflow_state  # noqa: E402
from Financial_Planning.Utilities.ppt_runner import (  # noqa: E402
    PptTemplateError,
    default_template_path,
    generate_financial_plan_ppt,
    sanitize_client_filename,
)


def test_json_safe_workflow_state_round_trips_date():
    raw = {"as_of": date(2026, 6, 30), "nested": {"value": 1}}
    safe = json_safe_workflow_state(raw)
    assert safe["as_of"] == "2026-06-30"
    assert safe["nested"]["value"] == 1


def test_sanitize_client_filename_strips_invalid_chars():
    assert sanitize_client_filename('Aakash / "Test"') == "Aakash  Test"
    assert sanitize_client_filename("") == "Client"
    assert sanitize_client_filename("   ") == "Client"


def test_generate_financial_plan_ppt_rejects_empty_state():
    with pytest.raises(ValueError, match="non-empty dict"):
        generate_financial_plan_ppt({})


def test_generate_financial_plan_ppt_missing_template():
    with pytest.raises(PptTemplateError, match="template not found"):
        generate_financial_plan_ppt(
            {"client_data": {"client_data": {"name": "X"}}},
            template_path="/nonexistent/template.pptx",
        )


def test_generate_financial_plan_ppt_returns_pk_zip_with_mocks():
    """Smoke test: mocked populate + LLM produces a valid zip (pptx) buffer."""
    from pptx import Presentation

    template = default_template_path()
    if not Path(template).is_file():
        pytest.skip("PPT template not present on this machine")

    prs = Presentation(template)

    def fake_populate(presentation, final_state, llm_azure):
        return final_state["client_data"]["client_data"]["name"]

    with patch(
        "Financial_Planning.Utilities.ppt_runner._build_azure_llm",
        return_value=MagicMock(),
    ), patch(
        "Financial_Planning.Utilities.ppt_runner.populate_presentation",
        side_effect=lambda p, s, _llm: fake_populate(p, s, _llm),
    ), patch(
        "Financial_Planning.Utilities.ppt_runner.Presentation",
        return_value=prs,
    ):
        ppt_bytes, filename = generate_financial_plan_ppt(
            {"client_data": {"client_data": {"name": "Priya"}}},
            template_path=template,
        )

    assert ppt_bytes[:2] == b"PK"
    assert filename == "Priya_financial_plan.pptx"


def test_generate_financial_plan_ppt_missing_azure_env():
    template = default_template_path()
    if not Path(template).is_file():
        pytest.skip("PPT template not present on this machine")

    from Financial_Planning.Utilities.ppt_runner import PptDependencyError

    with patch.dict("os.environ", {}, clear=True), patch(
        "Financial_Planning.Utilities.ppt_runner.load_dotenv"
    ):
        with pytest.raises(PptDependencyError, match="AZURE_API"):
            generate_financial_plan_ppt(
                {"client_data": {"client_data": {"name": "X"}}},
                template_path=template,
            )


def test_build_term_cover_slide_adds_table():
    template = default_template_path()
    if not Path(template).is_file():
        pytest.skip("PPT template not present on this machine")

    from pptx import Presentation
    from Financial_Planning.Utilities.ppt_builder import PPTBuilder

    prs = Presentation(template)
    state = {
        "required_retirement_corpus": {"client_info": {"current_monthly_expenses": 100_000}},
        "client_data": {"education_planning_summary": [], "insurance_policies": []},
        "liabilities": [{"outstanding_balance": 500_000}],
        "liquid_pool": 200_000,
        "term_insurance_summary": {
            "pv_of_expenses": 20_000_000,
            "kids_education_cost": 1_000_000,
            "current_liabilities": 500_000,
            "existing_cover": 0,
            "liquidable_assets": 200_000,
            "total_term_required": 21_300_000,
        },
    }
    builder = PPTBuilder(prs, state)
    assert builder.build_term_cover_slide(22) is True
    tables = [s for s in prs.slides[22].shapes if s.has_table]
    assert tables
    assert "21,300,000" in tables[0].table.cell(6, 1).text or "2,13,00,000" in tables[0].table.cell(6, 1).text


def test_replace_pie_chart_without_template_chart_does_not_raise():
    """Wealth slide may lack a pie placeholder — chart is placed at default coords."""
    from pptx import Presentation
    from Financial_Planning.Utilities.ppt_utilities import replace_pie_chart_with_matplotlib

    prs = Presentation()
    layout = prs.slide_layouts[6]
    prs.slides.add_slide(layout)
    replace_pie_chart_with_matplotlib(
        prs,
        0,
        {"EPF": 0.5, "PPF": 0.5},
        "Wealth Created - 2045",
    )
    assert len(prs.slides[0].shapes) >= 1


def test_find_term_cover_slide_after_optional_deletes():
    template = default_template_path()
    if not Path(template).is_file():
        pytest.skip("PPT template not present on this machine")

    from pptx import Presentation
    from Financial_Planning.Utilities.ppt_runner import _resolve_term_cover_slide
    from Financial_Planning.Utilities.ppt_utilities import delete_slide, find_slide_index_by_text

    prs = Presentation(template)
    for idx in sorted([16, 17, 18, 19, 14, 15, 7], reverse=True):
        if idx < len(prs.slides):
            delete_slide(prs, idx)

    assert find_slide_index_by_text(prs, "{term cover heading}") is not None
    term_idx = _resolve_term_cover_slide(prs, slide_offset=0)
    assert "{term cover heading}" in (
        prs.slides[term_idx].shapes[2].text_frame.text.lower()
    )
