"""
PowerPoint Report Generator - Main Script

What this file does:
This script generates the complete financial planning PowerPoint presentation.
It populates a template with workflow results including client info, analysis, goals, and strategies.

What this file contains and processes:
This script populates 17 slides with workflow data:
- Slide 1: Cover page with client name
- Slide 2: Client profile table and financial goals summary
- Slide 3: Blank slide
- Slide 4: Asset distribution pie chart analysis
- Slide 5: Asset pool values (fixed/liquid/retirement)
- Slide 6: Goals roadmap timeline visualization
- Slide 7: Financial overview with savings/expense donut chart
- Slide 8: Natural loan repayment schedule table
- Slide 9: Optimal loan prepayment strategy with interest savings
- Slide 10: Retirement analysis with corpus gap and commentary
- Slides 11-13: EPF/NPS/PPF retirement scheme details
- Slide 14: Retirement funding strategy with allocations
- Slide 15: Education planning summary and strategy tables
- Slide 16: Other financial goals allocation strategy
- Output: Saves populated presentation to ppt_output/populated.pptx
"""
import os
import sys
current_dir = os.path.dirname(os.path.abspath(__file__))
parent_dir = os.path.abspath(os.path.join(current_dir, '../..'))
if parent_dir not in sys.path:
    sys.path.insert(0, parent_dir)

from langchain_openai import AzureChatOpenAI
from dotenv import load_dotenv
from pydantic import BaseModel, Field
from pptx import Presentation
from pptx.util import Inches, Pt
from datetime import datetime, date
from Financial_Planning.Utilities.ppt_utilities import (find_and_update_table, indx_text_boxes, update_text_of_textbox, create_goals_roadmap,
                                                       create_financial_health_donut, delete_text_box, delete_slide, replace_pie_chart_with_matplotlib, duplicate_slide, move_slide, delete_tables_from_slide, replace_tables_with_text,create_risk_band_visual,create_financial_goals_visual) #, replace_pie_chart_with_data)
from Financial_Planning.Utilities.utility_functions import (sip_required, convert_currency, analyze_asset_portfolio)
from Financial_Planning.Models.llm_schemas import (slide7_llm, slide8_llm, retirement_commentary, Finalblock, Edu_commentator)
from Financial_Planning.Workflow.workflow import workflow
from Financial_Planning.input_data_personas import client_data
#Finanical_Planning.input_data inport client_data
import copy
from Financial_Planning.Utilities.prompts import (final_block_summary, edu_system)
from Financial_Planning.Utilities.ppt_builder import PPTBuilder

load_dotenv() 

AZURE_API_KEY=os.getenv('AZURE_API_KEY')
AZURE_API_BASE=os.getenv('AZURE_API_BASE')
AZURE_API_VERSION=os.getenv('AZURE_API_VERSION')
AZURE_DEPLOYMENT_NAME=os.getenv('AZURE_DEPLOYMENT_NAME')

llm_azure = AzureChatOpenAI(
    api_key=AZURE_API_KEY,  # AZURE_API_KEY
    azure_endpoint=AZURE_API_BASE,  # AZURE_API_BASE
    api_version=AZURE_API_VERSION,  # AZURE_API_VERSION
    deployment_name=AZURE_DEPLOYMENT_NAME,  # AZURE_DEPLOYMENT_NAME
    temperature=0  # Optional
) 

# execute the workflow
initial_state={'client_data': client_data, 'EMI_allocated': False, 'loan_prepayed_times': 0, 'used_monthly_surplus':[0], 'optimal_selected': False, }
final_state=workflow.invoke(initial_state, config = {"recursion_limit": 50})

slide7_structured_llm=llm_azure.with_structured_output(slide7_llm)

slide8_structured_llm=llm_azure.with_structured_output(slide8_llm)

retirement_commentary_llm=llm_azure.with_structured_output(retirement_commentary)

final_block_summary_llm=llm_azure.with_structured_output(Finalblock)

education_commentator_llm=llm_azure.with_structured_output(Edu_commentator)
############################## retirement commentary llm ############################################
########## PPT GENERATOR FUNCTION ###################################################################

prs=Presentation(os.path.join(os.path.dirname(os.path.dirname(__file__)), 'PPT_io', 'PPT_template', 'Financial Plan Stencil v2.pptx'))
indx_text_boxes(prs, 0)
# new_s=duplicate_slide(prs,12)
# client_name="Aakash"
# prs.save(f'Financial_Planning/PPT_io/PPT_output/{client_name}_financial_plan.pptx')
# status=move_slide(prs, 19, 14)

client_name=final_state['client_data']['client_data']['name']                          
client_dob=final_state['client_data']['client_data']['date_of_birth']
client_age=final_state['client_data']['client_data']['client_age']
desired_retirement_age=final_state['client_data']['client_data']['retirement_age']
spouse_name=final_state['client_data']['client_data']['spouse_name']                   #optional
spouse_dob=final_state['client_data']['client_data']['spouse_dob']                     #optional
childrens=', '.join([x['child_name'] for x in final_state['client_data']['client_data']['children']]) #optional

fixed_assets_pool=final_state['fixed_asset_pool']    
liquid_asset_pool=final_state['liquid_pool']
retirement_asset_pool = sum(
    list(asset.values())[0].get('current_value', list(asset.values())[0].get('maturity_value', 0))
    for asset in final_state.get('retirement_assets', [])
    if list(asset.keys())[0] != 'ulip'
)
#asset_classification=final_state['financial_overview']['category_percentage']
liquidity_ratio=round(final_state['financial_overview']['liquidity_ratio'], 2)
liquidity_flag=final_state['financial_overview']['liquidity_flag']
flexibility=final_state['financial_overview']['flexibility']
flexibility = (
    "Medium High" if isinstance(flexibility, str) and "medium to high" in flexibility.lower()
    else "high"   if isinstance(flexibility, str) and "high" in flexibility.lower()
    else "medium" if isinstance(flexibility, str) and "medium" in flexibility.lower()
    else "low"    if isinstance(flexibility, str) and "low" in flexibility.lower()
    else flexibility
)
savings_ratio=round(final_state['financial_overview']['spending_behavior']['saving_ratio'], 2)
expense_ratio=round(final_state['financial_overview']['spending_behavior']['expense_ratio'], 2)
red_flag=final_state['financial_overview']['spending_behavior']['red_flag']
#esops = final_state['financial_overview']['']
current_year=date.today().year  
retirement_year=current_year+final_state.get('required_retirement_corpus',{}).get('client_info',{}).get('years_to_retirement',{})
years_to_retire=final_state.get('required_retirement_corpus',{}).get('client_info',{}).get('years_to_retirement',{})
inflation=0.06
current_monthly_expense=final_state.get('required_retirement_corpus', {}).get('client_info', {}).get('current_monthly_expenses', {})
current_annual_expense=final_state.get('required_retirement_corpus', {}).get('client_info', {}).get('current_annual_expenses', {})
future_annual_expense=final_state.get('required_retirement_corpus', {}).get('client_info', {}).get('future_annual_expenses_at_retirement', {})
expected_retirement_corpus_needed=final_state.get('required_retirement_corpus', {}).get('recommendation', {}).get('recommended_corpus', {})
fv_of_retirement_schemes=final_state['retirement_schemes_fv']['grand_total']
corpus_gap=expected_retirement_corpus_needed-fv_of_retirement_schemes

delete_list=[ ]

goals = final_state['sorted_goals']
goal_names = [goal['goal_name'] for goal in goals]
final_goal_names = []
for g in goal_names:
    if ('UG' in g or 'PG' in g):
        if "Kid's Education" not in final_goal_names:
            final_goal_names.append("Kid's Education")
    else:
        final_goal_names.append(g)
        
# Initialize PPTBuilder
builder = PPTBuilder(prs, final_state)

# slide 1, index = 0
builder.build_intro_slide(slide=0)

# slide 2, index = 1
builder.build_personal_details_slide(slide=1)

# slide 3 index - 2
slide=2
print(f"Blank: slide: {slide}")

# slide 4
output_portfolio = builder.build_asset_allocation_slide(slide=3)

# Extract asset names and percentages from portfolio
#def prepare_portfolio_chart_data(asset_portfolio):
#    pie_data = {}
    
    # Get the percentage distribution list
#    percentage_distribution = asset_portfolio.get('percentage_distribution', [])
    
#    # Extract non-zero assets
#    for item in percentage_distribution:
#        asset_name = item.get('asset_name', '')
#        percentage = item.get('percentage', 0)
        
        # Only include assets with non-zero percentage
#        if percentage > 0:
#            # Convert percentage to ratio (divide by 100)
#            pie_data[asset_name] = percentage / 100
  
builder.build_asset_classification_slide(slide=4) 
total_asset_value=output_portfolio['portfolio_summary']['total_asset_value']
indx_text_boxes(prs, slide)
update_text_of_textbox(prs, slide, 2, "Total Asset Value: "+f"{convert_currency(total_asset_value)}")
indx_text_boxes(prs, slide) 
slide = 4
update_text_of_textbox(prs, slide, 2, f"{convert_currency(fixed_assets_pool)}")
update_text_of_textbox(prs, slide, 3, f"{convert_currency(liquid_asset_pool)}")
update_text_of_textbox(prs, slide, 4, f"{convert_currency(retirement_asset_pool)}")
slide=5
builder.build_goals_roadmap_slide(slide=5)
slide=6
builder.build_financial_health_checkup_slide(slide=6)
slide=9
corpus_gap = builder.build_retirement_outlook_slide(retirement_commentary_llm, slide=9)
slide=10
builder.build_retirement_goal_planning_slide(llm_azure, slide=10)
# ULIP slide - optional, only if ULIPs exist
has_ulips = builder.build_ulips_slide(slide=16)
if not has_ulips:
    delete_list.append(16)

# SSY slide - optional, only if Sukanya Samriddhi scheme exists
has_ssy = builder.build_ssy_scheme_slide(slide=17)
if not has_ssy:
    delete_list.append(17)

# Marriage slide - optional, only if marriage goals exist
has_marriage, marriage_slide_offset = builder.build_marriage_slide(slide=18)
if not has_marriage:
    delete_list.append(18)

has_misc, misc_slide_offset = builder.build_misc_goals_slides(template_slide=19)
if not has_misc:
    delete_list.append(19)

has_kids, edu_slide_offset = builder.prepare_and_build_education_slides(
    summary_slide=14, details_slide=15)
s = edu_slide_offset

if not has_kids:
    delete_list.append(14)
    delete_list.append(15)

has_liabilities = builder.build_liabilities_slide(slide=7)

#Wealth created at reitrement slide
builder.build_wealth_created_ret_breakdown(slide=20)
delete_list.sort(reverse=True)
for slideno in delete_list:
    delete_slide(prs, slideno) 
 

builder.save(client_name=client_name)

"""
# slide 8 # liability slide must be optional(Liability slide must only be created only if loan exist, if it does not exist then that template must be deleted)
slide=8 
loan_exist=final_state['loans_exist'] 
if loan_exist: 
    print(f"Loan Details: slide: {slide}")
    liabilities = final_state['liability_allocation'][0]['per_loan'] 
    keys = ['accelerated_months', 'accelerated_close_date', 'interest_saved','monthly_extra_assigned', 'monthly_extra_applied', 'lump_sum_assigned','lump_sum_applied']

    loans = []
    for loan in liabilities: 
        # remove each unwanted key 
        for k in keys: 
            loan.pop(k, None) 
        loans.append(loan) 

    col1_loan_type=["Loan Type"] 
    col2_close_year=["Close Year"]
    col3_freed=["Freed Funds"]

    if loans!=[]: 
        for loan in loans:
            col1_loan_type.append(loan["type"])
            col2_close_year.append( loan["baseline_close_date"].split("-")[0] )
            col3_freed.append(convert_currency(list(loan["freed_by_year"].values())[0]) )

        natural_loan_data={"Type": col1_loan_type, "Close Year": col2_close_year, "Freed": col3_freed}

        #prs=Presentation('ppt_generator/ppt_template/Financial Plan Stencil.pptx') 
        #indx_text_boxes(prs, 8)
        find_and_update_table(prs, slide, natural_loan_data)
else: 
    delete_list.append(slide)


# slide 9 # this slide is optional as well, if optimal loan allocation does exist then only this slide must be created or else delete
for i, item in enumerate(final_state["optimal_loan_allocation"]["per_loan"]):
    print(i, item.get("interest_saved"), item)
total_interest_saved=sum([i.get("interest_saved",0) for i in final_state["optimal_loan_allocation"]["per_loan"]])
slide=9
if total_interest_saved>0: 
    print(f"Loan Pre-Payments: slide: {slide}")
    if final_state['optimal_loan_allocation']['per_loan'][0]['interest_saved']>0: # only those loans will be showed in PPT which result in interest saved.
        col1_type=["Debt Type"] 
        col2_baseclose=["Base Close Year"]
        col3_monthly_topup=["Monthly Top Up"]
        col4_lumpsum_topup=["Lumpsum"]
        col5_interest_saved=["Interest Saved"]
        col6_acc_close=["Close Year"]
        col7_free_fund=["Freed Fund"]
        for loan in final_state['optimal_loan_allocation']['per_loan']:

            col1_type.append(loan['type'])
            col2_baseclose.append(loan['baseline_close_date'].split("-")[0]) 
            col3_monthly_topup.append((convert_currency(loan['monthly_extra_applied'])))
            col4_lumpsum_topup.append(convert_currency(loan['lump_sum_applied']))
            col5_interest_saved.append(round(round(loan['interest_saved']), -2))
            col6_acc_close.append(loan['accelerated_close_date'].split("-")[0])
            col7_free_fund.append(convert_currency(list(loan['freed_by_year'].values())[0]))

        total_interest_saved=convert_currency(sum(col5_interest_saved[1:]))
        col1_type.append("Total")  
        col2_baseclose.append("") 
        col3_monthly_topup.append("")
        col4_lumpsum_topup.append("")
        col5_interest_saved.append(total_interest_saved) 
        col6_acc_close.append("")
        col7_free_fund.append("")

        loan_pre_data={"col1":col1_type, "col2_baseclose": col2_baseclose, "col3_monthly_topup": col3_monthly_topup, "col4_lumpsum_topup": col4_lumpsum_topup,
                "col6_acc_close": col6_acc_close, "col7_free_fund": col7_free_fund, "col5_interest_saved": col5_interest_saved }
    
        #indx_text_boxes(prs, 9)
        find_and_update_table(prs, slide, loan_pre_data) 
else: 
    delete_list.append(slide)

# slide 11 optional must be created only if epf scheme exist
slide=11 
epf_fv=final_state.get('retirement_schemes_fv', {}).get('category_totals',{}).get('epf',{})
epf_scheme=final_state.get('client_data',{}).get('investment_details',{}).get('retirement_investments',{}).get('epf', {})
if epf_fv!={}: 
    print(f"EPF: slide: {slide}")
    # indx_text_boxes(prs, 11) 
    epf_current_value=epf_scheme[0]['current_value']
    epf_monthly_contribution=epf_scheme[0]['employee_employer_contribution_monthly']
    epf_interest_rate=epf_scheme[0]['interest_rate'] 
    update_text_of_textbox(prs, slide, 9, f"{epf_interest_rate*100}%" )  
    update_text_of_textbox(prs, slide, 10, f"{convert_currency(epf_current_value)}" )
    update_text_of_textbox(prs, slide, 11, f"{convert_currency(epf_monthly_contribution)}" ) 
    update_text_of_textbox(prs, slide, 12, f"{convert_currency(epf_fv)}" )
    update_text_of_textbox(prs, slide, 13, str(years_to_retire) )
    update_text_of_textbox(prs, slide, 14, str(desired_retirement_age))
else: 
    delete_list.append(slide) 
 
# slide 12 optional, must be created only if nps exist
slide=12
nps_fv=final_state.get('retirement_schemes_fv', {}).get('category_totals',{}).get('nps',{})
nps_scheme=final_state.get('client_data',{}).get('investment_details',{}).get('retirement_investments',{}).get('nps', {})
if nps_fv!={}:  
    print(f"NPS: slide: {slide}")
    # indx_text_boxes(prs, 12) 
    nps_current_value=nps_scheme[0]['current_value']
    nps_maturity_year=nps_scheme[0]['maturity_year']
    nps_monthly_contribution=nps_scheme[0]['monthly_contribution']
    expected_corpus_growth_rate=nps_scheme[0]['expected_corpus_growth_rate']
    update_text_of_textbox(prs, slide, 8, f"{expected_corpus_growth_rate*100}%")
    update_text_of_textbox(prs, slide, 9, str(retirement_year))
    update_text_of_textbox(prs, slide, 10, str(final_state['client_data']['client_data']['retirement_age']))
    update_text_of_textbox(prs, slide, 11, str(convert_currency(nps_current_value)))
    update_text_of_textbox(prs, slide, 12, str(convert_currency(nps_monthly_contribution)))
    update_text_of_textbox(prs, slide, 13, str(convert_currency(nps_fv)))
else: 
    delete_list.append(slide) 

# slide 13 optional, must be created only if ppf exist.
slide=13
ppf_fv=final_state.get('retirement_schemes_fv', {}).get('category_totals',{}).get('ppf',{})
ppf_scheme=final_state.get('client_data',{}).get('investment_details',{}).get('retirement_investments',{}).get('ppf', {})
if ppf_fv!={}:
    print(f"PPF: slide: {slide}")
    # indx_text_boxes(prs, 13)
    ppf_current_value=ppf_scheme[0]['current_value']
    ppf_annual_contribution=ppf_scheme[0]['annual_contribution']
    ppf_interest_rate=ppf_scheme[0]['interest_rate']
    update_text_of_textbox(prs, slide, 9, f"{ppf_interest_rate*100}%")
    update_text_of_textbox(prs, slide, 10, str(retirement_year))
    update_text_of_textbox(prs, slide, 11, str(final_state['client_data']['client_data']['retirement_age']))
    update_text_of_textbox(prs, slide, 12, str(convert_currency(ppf_current_value)))
    update_text_of_textbox(prs, slide, 13, str(convert_currency(ppf_annual_contribution)))
    update_text_of_textbox(prs, slide, 6, str(convert_currency(ppf_fv)))
else: 
    delete_list.append(slide)
"""
# slide 14 compulsion, based on the retirement goal allocation it must populate the slide, 
# if the retirement goal is not met then it will state how much is required to achieve that goal
"""
slide=15
indx_text_boxes(prs, 15)
if final_state['client_data']['client_data']['if_any_kids']:
    edu_detail=[]
    for i in final_state['client_data']['education_planning_summary']:
        edu_detail.append({'child_name': i['name'], 'edu_type': i['type'], 'stream': i['stream'], 'destination': i['destination'], 
                           'target_year': i['target_year'], 'current_cost': i['current_cost'], 'future_cost': i['future_cost'],
                           'future_value_of_allocated_funds': i['total_future_corpus'], 'corpus_gap': i['final_gap'], 'sourced_from': i['funded_from']})
    
        for j in final_state['optimal_goal_allocation']['goals']:
            if i['name']+" "+i['type']==j['goal_name']:
                print(f"j: {j}")
                edu_detail[-1].update({'final_gap': j['corpus_gap'],'status': j['filter'][0]['type'], 'deprioritized': j['depriorized'], 'note':j['note'], 'funded_from':j['funded_from'], 'corpus_needed':j['corpus_needed']})
print("=" * 50)
update_text_of_textbox(prs,21, 1,f"Education Planning - {edu_detail[0]['child_name'] if len(edu_detail)>0 else 'N/A'}")
# Better approach - find UG and PG costs for a specific child
child_name = edu_detail[0]['child_name'] if len(edu_detail) > 0 else 'N/A'
ug_cost = "N/A"
pg_cost = "N/A"
for edu in edu_detail:
    if edu['child_name'] == child_name:
        if edu['edu_type'] == 'UG':
            ug_cost = convert_currency(edu['future_cost'])
            ug_future_cost = convert_currency(edu['future_cost'])
        elif edu['edu_type'] == 'PG':
            pg_cost = convert_currency(edu['future_cost'])
            pg_future_cost = convert_currency(edu['future_cost'])

kid_1_edu_data = {
    "Particulars": ["Particulars", "Est Educational Cost - UG", "Est Educational Cost - PG", "Inflation Rate"],
    "Values": ["Amount", ug_cost, pg_cost, "7%"]
}

print("\n" + "=" * 50)
print("SLIDE 20 TEXTBOXES:")
print("=" * 50)
# Build funded_from text for UG goal - each source on new line
ug_funded_text = ""
ug_corpus_created = 0

for edu in edu_detail:
    if edu['child_name'] == edu_detail[0]['child_name'] and edu['edu_type'] == 'UG':
        funded_from = edu.get('funded_from', [])
        if funded_from:
            parts = []
            for fund in funded_from:
                fund_type = fund.get("type", "")
                # Sum up fv_contribution from each funding source
                ug_corpus_created += fund.get('fv_contribution', 0)
                target_year = edu.get('target_year', 'N/A')
                
                if fund_type == "freed_sip":
                    parts.append(f"Freed EMI: {convert_currency(fund.get('monthly', 0))}/month")
                elif fund_type in ["sip_from_surplus", "sip_from_partial_surplus"]:
                    parts.append(f"SIP: {convert_currency(fund.get('monthly', 0))}/month")
                elif fund_type in ["lumpsum_from_liquid", "lumpsum_from_liquid_partial"]:
                    parts.append(f"Lumpsum: {convert_currency(fund.get('principal_used_today', 0))}")
            ug_funded_text = "\n".join(parts)
        else:
            ug_funded_text = "Fully Funded"
        break

    ug_corpus_created = convert_currency(ug_corpus_created) if ug_corpus_created > 0 else "N/A"
    ug_future_cost = convert_currency(edu_detail[0]['future_cost']) if len(edu_detail) > 0 else 'N/A'

    update_text_of_textbox(prs, 22, 1, f"Education Planning - {child_name}")
    update_text_of_textbox(prs, 22, 2, f"Undergraduate (UG) - {target_year}", font_size=28, bold=True,font_color=RGBColor(8,232, 222))
    formatted_funded_text = ug_funded_text.replace(", ", "\n  • ")

    update_text_of_textbox(prs, 22, 4, 
    f"Est. Future Cost - {ug_future_cost}\n"
    f"Funded From:\n"
    f"  • {formatted_funded_text}\n"
    f"Corpus Created - {ug_corpus_created}", 
    font_size=32, bold=True)
            
"""


# slide 15, optional must be created only if kids exist
# slide=15 
# slide+=1
# if final_state['client_data']['client_data']['if_any_kids']:

#     col1_child_name=["Child Name"]
#     col2_education=["Education"]
#     col3_target_year=["Target Year"]
#     col4_current_cost=["Current Cost"]
#     col5_allocated_funds=["Allocated Funds"]
#     col6_scheme=["Scheme"]
#     col7_total_available=["Future Value Of Investments"]
#     col8_future_cost=["Future Cost"]
#     col9_gap=["Corpus Gap"] 
#     for education in final_state['client_data']['education_planning_summary']:
#         # print(education["name"]) 
#         col1_child_name.append(education["name"])
#         # if education["type"]=='UG':
#         #     col2_education.append("UG")
#         # elif education["type"]=='PG':
#         #     col2_education.append("PG")
#         col2_education.append(education["type"])
#         col3_target_year.append(education["target_year"])
#         col4_current_cost.append(convert_currency(education["current_cost"])) 
#         col5_allocated_funds.append(convert_currency(education["allocated_funds"]))
#         col6_scheme.append(education["schemes"][0]['scheme_name']) 
#         col7_total_available.append(convert_currency(education["total_future_corpus"]))
#         col8_future_cost.append(convert_currency(education["future_cost"]))
#         col9_gap.append(convert_currency(education["final_gap"]))
    
#     edu_summary_data={"col1_child_name": col1_child_name, "col2_education": col2_education, "col3_target_year": col3_target_year, "col4_current_cost": col4_current_cost,
#                       "col8_future_cost": col8_future_cost, "col7_total_available": col7_total_available, "col9_gap": col9_gap }
    
#     #prs=Presentation('ppt_generator/ppt_template/Financial Plan Stencil.pptx') 
#     #indx_text_boxes(prs, 15) 
#     find_and_update_table(prs, slide, edu_summary_data)
#     #prs.save('ppt_generator/ppt_output/populated.pptx')

#     edu_allocations=[]
#     for allocations in sorted_goals:

#         #if (allocations["goal_name"].split("-")[-1]=="under_graduation" or allocations["goal_name"].split("-")[-1]=="post_graduation") and allocations['corpus_needed']!=0: 
#         if "UG" in allocations["goal_name"].split(" ") and allocations['corpus_needed']!=0: 
#             #allocations['goal_name']=allocations["goal_name"].split(" ")[0] + ' UG' 
#             edu_allocations.append(allocations) 
#         elif "PG" in allocations["goal_name"].split(" ") and allocations['corpus_needed']!=0: 
#             #allocations['goal_name']=allocations["goal_name"].split(" ")[0] + ' PG' 
#             edu_allocations.append(allocations)
#         else: 
#             continue

#     col1_goal_name=["Goal Name"]
#     col2_corpus_gap=["Corpus Gap"]
#     col3_investment_type=["Investment Type"]
#     col4_investment=["Investment"]
#     col5_from_year=["From Year"]
#     col6_to_year=["To Year"]
#     col7_rate=["Interest Rate"]
#     col8_fv=["Future Value"] #col8_fv
    
#     for edu in edu_allocations: 
#         col1_goal_name.append(edu["goal_name"])
#         col2_corpus_gap.append(convert_currency(edu["corpus_needed"]))
        
#         for fund in edu["funded_from"]: 
#             if fund.get("type", "None")=="freed_sip":
#                 col3_investment_type.append("Freed EMI")
#             elif fund.get("type", "None")=="sip_from_surplus" or fund.get("type", "None")=="sip_from_partial_surplus":
#                 col3_investment_type.append("SIP")
#             elif fund.get("type", "None")=="lumpsum_from_liquid" or fund.get("type", "None")=="lumpsum_from_liquid_partial":  
#                 col3_investment_type.append("Lumpsum")
#             else: 
#                 continue 
#             col4_investment.append(convert_currency(fund.get("monthly", fund.get('principal_used_today', 0) )))
#             col5_from_year.append(fund["from_year"]) 
#             col6_to_year.append(fund["to_year"]) 
#             col7_rate.append(fund['rate']) 
#             col8_fv.append(convert_currency(fund["fv_contribution"])) 
#         if len(edu["funded_from"])>1: 
#             for i in range(1,len(edu["funded_from"])):
#                 col1_goal_name.append("")
#                 col2_corpus_gap.append("")
    
#     education_strategy={"col1_goal_name": col1_goal_name, "col2_corpus_gap": col2_corpus_gap, "col3_investment_type": col3_investment_type, "col4_investment": col4_investment, "col5_from_year": col5_from_year, 
#                         "col6_to_year": col6_to_year, 'col7_rate':col7_rate , "col8_fv": col8_fv} 
    
#     #prs=Presentation('ppt_generator/ppt_template/Financial Plan Stencil.pptx') 
#     #indx_text_boxes(prs, 16)
#     find_and_update_table(prs, 15, education_strategy)
#     #prs.save('ppt_generator/ppt_output/populated.pptx')
# else: 
#     delete_list.append(15)
#     delete_list.append(16)

    
# slide 17
# slide+=1
#final_state['optimal_goal_allocation']["goals"]
# slide=17 
"""
other_goals=[]
for allocations in sorted_goals: 
    if "Retirement" in allocations["goal_name"].split(" "):
        continue 
    elif "UG" in allocations["goal_name"].split(" ") or "PG" in allocations["goal_name"].split(" "):
        continue
    else: 
        other_goals.append(allocations)
slide=15 
if other_goals!=[]: 
    
    indx_text_boxes(prs, slide+s)
    print(f"Other Goals: slide: {slide+s}")
    col1_goal_name=["Goal Name"]
    col2_corpus_gap=["Corpus Gap"]
    col3_investment_type=["Investment Type"]
    col4_investment=["Investment"]
    col5_from_year=["From Year"]
    col6_to_year=["To Year"]
    col7_rate=["Interest Rate"]
    col8_fv=["Future Value"] # col8_fv

    for fin_goal in other_goals: 
        col1_goal_name.append(fin_goal["goal_name"])
        col2_corpus_gap.append(convert_currency(fin_goal["corpus_needed"])) 

        for fund in fin_goal["funded_from"]:
            if fund.get("type", "None")=="freed_sip": 
                col3_investment_type.append("Freed EMI")
            elif fund.get("type", "None")=="sip_from_surplus" or fund.get("type", "None")=="sip_from_partial_surplus" :
                col3_investment_type.append("SIP")
            elif fund.get("type", "None")=="lumpsum_from_liquid" or fund.get("type", "None")=="lumpsum_from_liquid_partial":
                col3_investment_type.append("Lumpsum")
            else: 
                continue    
            col4_investment.append(convert_currency(fund.get("monthly", fund.get('principal_used_today', 0) )))
            col5_from_year.append(fund["from_year"]) 
            col6_to_year.append(fund["to_year"]) 
            col7_rate.append(fund['rate'])
            col8_fv.append(convert_currency(fund["fv_contribution"]))
        if len(fin_goal["funded_from"])>1: 
            for i in range(1,len(fin_goal["funded_from"])): 
                col1_goal_name.append("") 
                col2_corpus_gap.append("") 
    fin_goal_data={"col1_goal_name": col1_goal_name, "col2_corpus_gap": col2_corpus_gap, "col3_investment_type": col3_investment_type, "col4_investment": col4_investment, "col5_from_year": col5_from_year, 
                            "col6_to_year": col6_to_year, "col7_rate": col7_rate ,"col8_fv": col8_fv}     
    find_and_update_table(prs, slide+s, fin_goal_data) 
else: 
    delete_list.append(slide+s) 
"""    
# slide 18
# slide=18
#slide=16 
#notes=[]
#for goal in sorted_goals:
#    if goal['goal_name']=='Retirement':
#        if goal['corpus_gap']==0:
#            note1="Your retirement goal is fully funded"
#            
#        elif round(goal['corpus_gap']/goal['target_corpus'],2)<=0.40:
#            note1=f"Retirement goal is moderately feasible. Corpus gap of {round(goal['corpus_gap']/goal['target_corpus'],2)*100}% still remains. Invest X amount of SIP from year xxxx to year xxxx to achieve this corpus gap."
#        else:
#            note1=f"Retirement goal is not feasible. Corpus gap of {round(goal['corpus_gap']/goal['target_corpus'],2)*100}% still remains. That requires you to invest x amount from year xxxx to year xxxx to achieve your retirement." 
#      {"Values": ["Amount", ug_cost, pg_cost, "7%"]
#}

#update_text_of_textbox(prs, 21, 1, f"Education Planning - {child_name}")
#find_and_update_table(prs, 21, kid_1_edu_data)      notes.append(note1)

#    if goal['goal_name'].split(" ")[-1]=="UG":
#        if goal['corpus_gap']==0:
#            note2=f'UG goal of {goal['goal_name'].split(" ")[0]} is fully funded, which ensure your child career is secured and well planned.'
#        elif round(goal['corpus_gap']/goal['target_corpus'],2)<=0.40:
#            note2=f"UG goal of {goal['goal_name'].split(" ")[0]} is partially funded, you need to invest X amount from year xxxx to year xxxx to achieve this goal."
#        else: 
#            note2=f"UG goal of {goal['goal_name'].split(" ")[0]} is under funded, due to shortage of funds, you need to invest X from year xxxx to year xxxx to achieve this goal."
#        notes.append(note2)
#    if goal['goal_name'].split(" ")[-1]=="PG":
#        if goal['corpus_gap']==0:
#            note3=f"PG goal of {goal['goal_name'].split(" ")[0]} is fully funded, which ensure your child career is secured and well planned."
#        elif round(goal['corpus_gap']/goal['target_corpus'],2)<=0.40:
#            note3=f"PG goal of {goal['goal_name'].split(" ")[0]} is partially funded, you need to invest X amount from year xxxx to year xxxx to achieve this goal."
#        else: 
#            note3=f"PG goal of {goal['goal_name'].split(" ")[0]} is under funded, due to shortage of funds, you need to invest X from year xxxx to year xxxx to achieve this goal."
#        notes.append(note2)
#liquidity_ratio=final_state["financial_overview"]['liquidity_ratio']
#if liquidity_ratio<0.15:
#    diff=round(0.15-liquidity_ratio, 2)
#    grand_total=final_state['liquid_pool']+final_state['fixed_asset_pool']+final_state['retirement_schemes_fv']['grand_total']
#    increase=round(diff*(grand_total),2) 
#    note4=f"Your liquidity ratio is {liquidity_ratio} which below the recommended, we suggest you to rebalance for porfolio, i.e increase by {increase} in liquid instruments so that your liquidity ratio is atleast 0.15"
#    notes.append(note4)
#else:
#    note4=f"Your liquidity ratio is {liquidity_ratio} which indicates a healthy portfolio, this ensure that you have enough assets that can be liquidates in times of needful."
#    notes.append(note4)

#final_block_content=final_block_summary_llm.invoke([final_block_summary]+notes).final_block_summary


#indx_text_boxes(prs, slide+s)
#update_text_of_textbox(prs, slide+s, 2, final_block_content)


