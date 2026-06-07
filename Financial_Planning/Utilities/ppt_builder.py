"""
PPT Builder - PowerPoint Slide Builder for Financial Planning
"""

import copy
import os
import sys

from pydantic import BaseModel, Field

from datetime import date
from pptx.util import Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from lxml import etree
from pydantic import BaseModel, Field
from pptx.util import Inches
from Financial_Planning.Utilities.utility_functions import convert_currency, analyze_asset_portfolio, sip_required
from Financial_Planning.RSU.webscrapper import load_rsu_market_data
from Financial_Planning.RSU.constants import get_rsu_growth_rate

current_dir = os.path.dirname(os.path.abspath(__file__))
parent_dir = os.path.abspath(os.path.join(current_dir, '../..'))
if parent_dir not in sys.path:
    sys.path.insert(0, parent_dir)

from Financial_Planning.Utilities.ppt_utilities import (
    update_text_of_textbox,
    duplicate_slide,
    move_slide,
    find_and_update_table,
    indx_text_boxes,replace_pie_chart_with_matplotlib,
    create_goals_roadmap,
    delete_text_box,add_table_to_slide
)
from Financial_Planning.Utilities.utility_functions import convert_currency, analyze_asset_portfolio


class PPTBuilder:
    """
    PPT Builder for Financial Planning presentations.
    Handles education planning slides with support for multiple children.
    """

    @staticmethod
    def _add_rsu_esop_table(slide, data, left, top, width, height,
                            header_fill, header_font_color,
                            body_fill, body_font_color, alt_fill,
                            font_size=None):

        rows = len(data)
        cols = len(data[0])
        shape = slide.shapes.add_table(rows, cols, left, top, width, height)
        table = shape.table

        for r_idx, row in enumerate(data):
            for c_idx, value in enumerate(row):
                cell = table.cell(r_idx, c_idx)
                tf = cell.text_frame
                tf.word_wrap = True
                p = tf.paragraphs[0]
                p.alignment = PP_ALIGN.LEFT

                run = p.add_run() if not p.runs else p.runs[0]
                run.text = str(value)
                run.font.name = "Calibri"
                run.font.size = font_size if font_size is not None else Pt(11)

                if r_idx == 0:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = header_fill
                    run.font.bold = True
                    run.font.color.rgb = header_font_color
                else:
                    fill_color = alt_fill if r_idx % 2 == 0 else body_fill
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = fill_color
                    run.font.bold = False
                    run.font.color.rgb = body_font_color

                tc = cell._tc
                tcPr = tc.get_or_add_tcPr()
                for edge in ('lnL', 'lnR', 'lnT', 'lnB'):
                    ln = etree.SubElement(tcPr, '{http://schemas.openxmlformats.org/drawingml/2006/main}' + edge)
                    ln.set('w', '6350')
                    ln.set('cap', 'flat')
                    ln.set('cmpd', 'sng')
                    solidFill = etree.SubElement(ln, '{http://schemas.openxmlformats.org/drawingml/2006/main}solidFill')
                    srgbClr = etree.SubElement(solidFill, '{http://schemas.openxmlformats.org/drawingml/2006/main}srgbClr')
                    srgbClr.set('val', 'C0C0C0')

        return shape

    def __init__(self, prs, final_state):
        self.prs = prs
        self.final_state = final_state
        self.slide_offset = 0

    @staticmethod
    def _format_retirement_commentary(text):
          """Split commentary into 2-3 paragraphs for PPT readability."""
          sentences = text.split(". ")
          if len(sentences) <= 2:
            return text
          para1 = sentences[0] + "."
          para2 = " ".join(sentences[1:3]) + "."
          para3 = " ".join(sentences[3:]) if len(sentences) > 3 else ""
          return f"{para1}\n\n{para2}\n\n{para3}"

    def _write_formatted_commentary(self, slide, textbox_order, text):
        """
        Write formatted multi-paragraph commentary into a textbox.
        
        Special formatting rules:
        - Slide 10 (Retirement Outlook): All textboxes get formatted commentary
        - Slide 14, Textbox 2 (Retirement Funding): Single paragraph only
        - Other slides: Standard formatting
        """
        slide_obj = self.prs.slides[slide]
        text_shapes = [s for s in slide_obj.shapes if s.has_text_frame]

        if textbox_order >= len(text_shapes):
            raise IndexError(
                f"Textbox index {textbox_order} out of range. "
                f"Found only {len(text_shapes)} text boxes on slide {slide}."
            )

        shape = text_shapes[textbox_order]
        tf = shape.text_frame
        tf.clear()

        use_formatted_commentary = (
            slide == 9 or
            (slide == 10 and textbox_order == 2)
        )

        if use_formatted_commentary:
            paragraphs = text.split("\n\n")
            if slide == 10 and textbox_order == 2:
                paragraphs = [paragraphs[0]]

            for i, block in enumerate(paragraphs):
                p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
                p.text = block
                p.alignment = PP_ALIGN.LEFT
                p.font.name = "Calibri"
                p.font.color.rgb = RGBColor(44, 62, 80)
                if i == 0:
                    p.font.size = Pt(28)
                    p.font.bold = True
                else:
                    p.font.size = Pt(24)
                    p.font.bold = False
                p.space_before = Pt(10) if i > 0 else Pt(0)
                p.space_after = Pt(6)
                p.line_spacing = 1.3
        else:
            p = tf.paragraphs[0]
            p.text = text
            p.alignment = PP_ALIGN.LEFT
            p.font.name = "Calibri"
            p.font.size = Pt(24)
            p.font.bold = False
            p.font.color.rgb = RGBColor(44, 62, 80)

    
    def _get_child_data(self, child_name):
       ug_data = None
       pg_data = None
       for edu in self.edu_detail:
         if edu['child_name'] == child_name:
             if edu['edu_type'] == 'UG':
                ug_data = edu
             elif edu['edu_type'] == 'PG':
                pg_data = edu
       return ug_data, pg_data
    
    def _get_edu_funding_summary(self, edu_data):
        """
        Returns (sip_amount, lumpsum_amount) for an education goal,
        summing only sip_from_surplus / sip_from_partial_surplus and
        lumpsum_from_liquid / lumpsum_from_liquid_partial.
        freed_sip is intentionally excluded.
        """
        if edu_data is None:
            return None, None
        sip_total     = 0.0
        lumpsum_total = 0.0
        for fund in edu_data.get("funded_from", []):
            ft = fund.get("type", "")
            if ft in ("sip_from_surplus", "sip_from_partial_surplus"):
                sip_total += fund.get("monthly", 0)
            elif ft in ("lumpsum_from_liquid", "lumpsum_from_liquid_partial"):
                lumpsum_total += fund.get("principal_used_today", fund.get("amount_used", 0))
        return sip_total or None, lumpsum_total or None

    def _build_funding_text(self, edu_data):
      if edu_data is None:
          return "No goal planned"

      future_cost = convert_currency(edu_data.get('future_cost', 0))
      funded_from = edu_data.get('funded_from', [])
      funding_lines = []
      corpus_created = 0

      for fund in funded_from:
        fund_type = fund.get("type", "")
        if fund_type == "freed_sip":
            inv_label = "Freed EMI"
        elif fund_type in ["sip_from_surplus", "sip_from_partial_surplus"]:
            inv_label = "SIP"
        elif fund_type in ["lumpsum_from_liquid", "lumpsum_from_liquid_partial"]:
            inv_label = "Lumpsum"
        elif fund_type == "ssy_funds":
            inv_label = "SSY"
        elif fund_type == "esop_funds":
            inv_label = "ESOP"
        elif fund_type == "rsu_funds":
            inv_label = fund.get("source", "RSU")
        else:
            continue

        amount = fund.get("monthly", fund.get("principal_used_today", fund.get("amount_used", 0)))
        corpus_created += fund.get(
            "fv_contribution",
            fund.get("usable_esop_fv_at_goal", fund.get("amount_used", 0)),
        )

        # Added '•' as the default bullet point and bolded the label
        if fund_type in ["lumpsum_from_liquid", "lumpsum_from_liquid_partial"]:
            funding_lines.append(f"• {inv_label}: {convert_currency(amount)}")
        elif fund_type == "ssy_funds":
            funding_lines.append(f"• {inv_label}: {convert_currency(fund.get('amount_used', 0))}")
        elif fund_type == "esop_funds":
            funding_lines.append(
                f"• {inv_label}: {convert_currency(fund.get('fv_contribution', fund.get('amount_used', 0)))}"
            )
        elif fund_type == "rsu_funds":
            funding_lines.append(
                f"• {inv_label}: {convert_currency(fund.get('fv_contribution', fund.get('amount_used', 0)))}"
            )
        else:
            funding_lines.append(f"• {inv_label}: {convert_currency(amount)}/month")

      status = edu_data.get('status', '')
      notes = edu_data.get('note', [])

      # Bolding main headers for the final output
      if funding_lines:
        funded_text = "\n".join(funding_lines)
        corpus_text = convert_currency(corpus_created) if corpus_created > 0 else "N/A"
        result = f"Est. Future Cost: {future_cost}\n\nFunded From:\n{funded_text}\n\nCorpus Created: {corpus_text}"

        if status in ("partial_funded", "unfunded"):
          corpus_gap = convert_currency(edu_data.get('final_gap', edu_data.get('corpus_gap', 0)))
          result += f"\nCorpus Gap: {corpus_gap}"

      elif status in ("partial_funded", "unfunded"):
        corpus_gap = convert_currency(edu_data.get('final_gap', edu_data.get('corpus_gap', 0)))
        result = f"Est. Future Cost: {future_cost}\nCorpus Gap: {corpus_gap}\nNo existing funds allocated"
      else:
        result = f"Est. Future Cost: {future_cost}\nStatus: Fully Funded"

      if status in ("partial_funded", "unfunded") and notes:
        note_text = notes[0] if isinstance(notes, list) else str(notes)
        result += f"\n\n**Note: {note_text}"

      return result
   
    def _duplicate_two_slides(self, slide1, slide2, insert_position):
       # Clone summary template, move to target position
       duplicate_slide(self.prs, slide1)
       move_slide(self.prs, len(self.prs.slides) - 1, insert_position)

       # Clone details template, move right after summary
       duplicate_slide(self.prs, slide2)
       move_slide(self.prs, len(self.prs.slides) - 1, insert_position + 1)

       self.slide_offset += 2
       return insert_position, insert_position + 1
    
    def _get_unique_children(self):
       seen = []
       for edu in self.edu_detail:
         if edu['child_name'] not in seen:
            seen.append(edu['child_name'])
       return seen


    def build_education_slides(self, summary_slide=21, details_slide=22, inflation_rate="7%"):
        """
        Build education slides for all children dynamically.
        
        For each child:
        - Summary slide: title + table (UG cost, PG cost, inflation)
        - Details slide: funding breakdown per UG/PG
        
        First child uses template slides directly.
        Additional children get duplicated slides.
        
        Returns:
            slide_offset: number of extra slides added
        """
        children = self._get_unique_children()

        if not children:
            print("No children found.")
            return 0

        print(f"Building education slides for {len(children)} children: {children}")

        for idx, child_name in enumerate(children):
            print(f"\n[Child {idx + 1}] {child_name}")

            ug_data, pg_data = self._get_child_data(child_name)

            if idx == 0:
                # First child: use template slides directly
                current_summary = summary_slide
                current_details = details_slide
            else:
                # Additional children: duplicate and insert after previous pair
                insert_pos = summary_slide + (idx * 2)
                print(f"  Duplicating slides at position {insert_pos}...")
                current_summary, current_details = self._duplicate_two_slides(
                    summary_slide, details_slide, insert_pos
                )

            # --- Summary slide ---
            print(f"  Summary slide: {current_summary}")
            update_text_of_textbox(self.prs, current_summary, 1, f"Education Planning - {child_name}")

            ug_cost = convert_currency(ug_data['future_cost']) if ug_data else "N/A"
            pg_cost = convert_currency(pg_data['future_cost']) if pg_data else "N/A"

            summary_table = {
                "Particulars": ["Particulars", "Est Educational Cost - UG", "Est Educational Cost - PG", "Inflation Rate"],
                "Values": ["Amount", ug_cost, pg_cost, inflation_rate]
            }
            find_and_update_table(self.prs, current_summary, summary_table)

            # --- Funding summary table (SIP / Lumpsum per UG / PG) ---
            ug_sip, ug_lump = self._get_edu_funding_summary(ug_data)
            pg_sip, pg_lump = self._get_edu_funding_summary(pg_data)


            col_headers = [""]
            if ug_data:
                col_headers.append("UG")
            if pg_data:
                col_headers.append("PG")

            def _fmt_sip(v):
                return f"{convert_currency(v)}/mo" if v else "-"
            def _fmt_lump(v):
                return convert_currency(v) if v else "-"

            sip_row  = ["Monthly SIP"]
            lump_row = ["Lumpsum"]
            if ug_data:
                sip_row.append(_fmt_sip(ug_sip))
                lump_row.append(_fmt_lump(ug_lump))
            if pg_data:
                sip_row.append(_fmt_sip(pg_sip))
                lump_row.append(_fmt_lump(pg_lump))

            funding_table_data = [col_headers, sip_row, lump_row]

            slide_obj = self.prs.slides[current_summary]

            # Add heading above funding summary table
            heading_shape = slide_obj.shapes.add_textbox(Inches(0.9), Inches(6.5), Inches(5.0), Inches(0.4))
            heading_tf = heading_shape.text_frame
            heading_p = heading_tf.paragraphs[0]
            heading_p.text = "Saving Plan"
            heading_p.font.name = "Calibri"
            heading_p.font.size = Pt(28)
            heading_p.font.bold = True
            heading_p.font.color.rgb = RGBColor(0, 0, 0)
            heading_p.alignment = PP_ALIGN.LEFT
            
            funding_shape = add_table_to_slide(
                slide_obj,
                funding_table_data,
                left=Inches(0.9),
                top=Inches(7.5),
                width=Inches(5.0),
                height=Inches(2.0),
                font_size=Pt(24),
                font_color=RGBColor(0, 0, 0),
                header_fill=RGBColor(0, 176, 185),
                header_font_color=RGBColor(255, 255, 255),
                body_fill=RGBColor(255, 255, 255),
            )
            # Bold all body cells
            tbl = funding_shape.table
            for r_idx in range(1, len(funding_table_data)):
                for c_idx in range(len(funding_table_data[0])):
                    tbl.cell(r_idx, c_idx).text_frame.paragraphs[0].font.bold = True

            # --- Details slide ---
            print(f"  Details slide: {current_details}")
            update_text_of_textbox(self.prs, current_details, 1, f"Education Planning - {child_name}")

            ug_text = self._build_funding_text(ug_data)
            update_text_of_textbox(self.prs, current_details, 4, ug_text,bold=True,font_size=20)

            pg_text = self._build_funding_text(pg_data)
            update_text_of_textbox(self.prs, current_details, 5, pg_text,bold=True,font_size=20)

        print(f"\nEducation slides done. Children: {len(children)}, Extra slides added: {self.slide_offset}")
        return self.slide_offset
        

    def build_intro_slide(self, slide=0):
        """Build Slide 1 (Cover page) - populates client name."""
        client_name = self.final_state['client_data']['client_data']['name']
        print(f"Intro: slide: {slide}")
        indx_text_boxes(self.prs, slide)
        update_text_of_textbox(self.prs, slide, 2, client_name)

    def build_personal_details_slide(self, slide=1):
        """Build Slide 2 (Personal Details) - populates client profile table."""
        client_data = self.final_state['client_data']['client_data']
        client_name = client_data['name']
        client_dob = client_data['date_of_birth']
        client_age = client_data['client_age']
        desired_retirement_age = client_data['retirement_age']
        spouse_name = client_data['spouse_name']
        spouse_dob = client_data['spouse_dob']
        childrens = ', '.join([x['child_name'] for x in client_data['children']])

        print(f"Personal Details: slide: {slide}")
        slide2_table_data = {
            "Name": client_name,
            "Date of Birth": str(client_dob),
            "Current Age": str(client_age),
            "Desired Retirement Age": str(desired_retirement_age),
            "Spouse Name": str(spouse_name),
            "Date of Birth": str(spouse_dob),
            "Children": str(childrens)
        }

        col1_data = list(slide2_table_data.keys())
        col2_data = list(slide2_table_data.values())
        indx_text_boxes(self.prs, slide)

        table_data = {'Particulars': ['Particulars'] + col1_data, "Details": ["Details"] + col2_data}
        find_and_update_table(self.prs, slide, table_data, 0, debug=True)
    
    def build_asset_allocation_slide(self, slide=4):
        """Build Slide 4 (Asset Allocation Pie Chart) - analyzes portfolio and populates pie chart."""
        print(f"Pie Chart: slide: {slide}")
        output_portfolio = analyze_asset_portfolio(
            self.final_state['retirement_assets'],
            self.final_state['liquid_assets'],
            self.final_state['fixed_assets']
        )
        pie_data = {}
        for item in output_portfolio.get('percentage_distribution', []):
            asset_name = item.get('asset_name', '')
            percentage = item.get('percentage', 0)
            if percentage > 0:
                pie_data[asset_name] = percentage / 100

        replace_pie_chart_with_matplotlib(
            presentation=self.prs,
            slide_number=slide,
            pie_data=pie_data,
            title="Current Asset Allocation"
        )

        # Replace {CA_VALUE} placeholder in template textbox
        total_asset_value = output_portfolio['portfolio_summary']['total_asset_value']
        total_in_crores = round(total_asset_value / 1e7, 2)
        slide_obj = self.prs.slides[slide]
        for shape in slide_obj.shapes:
            if shape.has_text_frame and '{CA_VALUE}' in shape.text_frame.text:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if '{CA_VALUE}' in run.text:
                            run.text = run.text.replace('{CA_VALUE}', f'₹{total_in_crores} Cr')
                break

        return output_portfolio

    def build_asset_classification_slide(self, slide=5):
        print(f"Asset Classification: slide: {slide}")
        category_data = self.final_state['financial_overview']['category_percentage']
        replace_pie_chart_with_matplotlib(
            presentation=self.prs,
            slide_number=slide,
            pie_data=category_data,
            title="Asset Distribution Analysis"
        )

        # Populate fixed / liquid / retirement pool values on this slide
        fixed_assets_pool     = self.final_state['fixed_asset_pool']
        liquid_asset_pool     = self.final_state['liquid_pool']
        retirement_asset_pool = sum(
            list(asset.values())[0].get('current_value', list(asset.values())[0].get('maturity_value', 0))
            for asset in self.final_state.get('retirement_assets', [])
            if list(asset.keys())[0] != 'ulip'
        )
        update_text_of_textbox(self.prs, slide, 2, f"{convert_currency(fixed_assets_pool)}")
        update_text_of_textbox(self.prs, slide, 3, f"{convert_currency(liquid_asset_pool)}")
        update_text_of_textbox(self.prs, slide, 4, f"{convert_currency(retirement_asset_pool)}")

    def build_goals_roadmap_slide(self, slide=5):
        print(f"Roadmap: slide: {slide}")
        sorted_goals = copy.deepcopy(self.final_state['optimal_goal_allocation']['goals'])
        for goal in sorted_goals:
            if 'Under-Graduation' in goal['goal_name'].split(" "):
                goal['goal_name'] = goal['goal_name'].split(" ")[0] + " UG"
            elif 'Post-Graduation' in goal['goal_name'].split(" "):
                goal['goal_name'] = goal['goal_name'].split(" ")[0] + " PG"

        goals = [(x['goal_name'], x['target_year']) for x in sorted_goals]
        sorted_goals_plot = dict(sorted(goals, key=lambda x: x[1]))
        create_goals_roadmap(self.prs, slide, sorted_goals_plot)
    
    def build_financial_health_checkup_slide(self, slide=7):
        """Build Financial Health Checkup slide - savings vs expenses donut chart and ratios."""
        print(f"Financial Health Checkup: slide: {slide}")
        financial_overview = self.final_state['financial_overview']
        savings_ratio = round(financial_overview['spending_behavior']['saving_ratio'], 2)
        expense_ratio = round(financial_overview['spending_behavior']['expense_ratio'], 2)

        flexibility = financial_overview['flexibility']
        flexibility = (
            "Medium High" if isinstance(flexibility, str) and "medium to high" in flexibility.lower()
            else "high"   if isinstance(flexibility, str) and "high" in flexibility.lower()
            else "medium" if isinstance(flexibility, str) and "medium" in flexibility.lower()
            else "low"    if isinstance(flexibility, str) and "low" in flexibility.lower()
            else flexibility
        )

        indx_text_boxes(self.prs, slide_no=slide)
        update_text_of_textbox(self.prs, slide_no=slide, text_box_id=13, new_text=str(savings_ratio))
        update_text_of_textbox(self.prs, slide_no=slide, text_box_id=14, new_text=str(expense_ratio))
        update_text_of_textbox(self.prs, slide_no=slide, text_box_id=15, new_text=str(flexibility))

        pie_data = {"Savings": round(savings_ratio * 100, 2), "Expenses": round(expense_ratio * 100, 2)}
        replace_pie_chart_with_matplotlib(
            presentation=self.prs,
            slide_number=slide,
            pie_data=pie_data,
            title="Savings vs Expenses")
    
    def build_retirement_outlook_slide(self, retirement_commentary_llm, slide=10):
        """Build Slide 10 (Retirement Outlook) - retirement stats, corpus gap, and LLM commentary."""
        print(f"Retirement Outlook: slide: {slide}")

        # Extract retirement data from final_state
        current_year = date.today().year
        retirement_info = self.final_state.get('required_retirement_corpus', {})
        client_info = retirement_info.get('client_info', {})
        years_to_retire = client_info.get('years_to_retirement', 0)
        retirement_year = current_year + years_to_retire
        current_monthly_expense = client_info.get('current_monthly_expenses', 0)
        current_annual_expense = client_info.get('current_annual_expenses', 0)
        future_annual_expense = client_info.get('future_annual_expenses_at_retirement', 0)
        expected_retirement_corpus_needed = retirement_info.get('recommendation', {}).get('recommended_corpus', 0)
        fv_of_retirement_schemes = self.final_state['retirement_schemes_fv']['grand_total']
        corpus_gap = expected_retirement_corpus_needed - fv_of_retirement_schemes

        # LLM prompt
        prompt = (
            "You are a professional retirement specialist, YOU EXPLAIN LIKE RICHARD FEYNMAN, "
            "you craft your words such that the customer listens to you and understands what you say clearly, "
            "you also give correlations while explaining something. you are provided with a customer's basic details "
            "like retirement year, year left to retire, current monthly expense, current annual expense, "
            "expected required retirement corpus, future value of all the retirement schemes and the corpus gap. "
            "Your task is to provide an overview of the customer's retirement profile and provide a good analysis. "
            "EXPLAIN THE CUSTOMER HOW IMORTANT IT IS TO HAVE A RETIREMENT PLAN AND HOW CAN IT MAKE LIFE EASY. "
            "** USE TONE AS IF YOU ARE TALKING TO A CUSTOMER ITSELF, USE PRONOUNS SUCH AS YOU, YOUR. **"
            "** USE PROFFESSIONAL TONE **"
            "**AND DO NOT USE ANY PRONOUNS LIKE WE, I OR ANYTHING ELSE TO REPRESENT YOURSELF**"
            "**IF CORPUS GAP EXIST THEN EXPLAIN IT TO CLIENT THE IMPORTANCE OF HAVING RETIREMENT PLAN READY**"
            f"current_year: {current_year}, retirement_year: {retirement_year}, years_to_retire: {years_to_retire}, "
            f"current_monthly_expense: {current_monthly_expense}, current_annual_expense: {current_annual_expense}, "
            f"corpus_gap: {corpus_gap}."
            "This overview will be published in one of the slide of overall PPT, so craft this overview accordingly."
            "DO ENSURE THAT THE PROVIDED CONTENT CONTAINS NOT MORE THAN 100 WORDS ")

        ret_note = retirement_commentary_llm.invoke(prompt).commentary

        # Populate textboxes
        update_text_of_textbox(self.prs, slide, 2, str(retirement_year))
        update_text_of_textbox(self.prs, slide, 4, str(years_to_retire))
        update_text_of_textbox(self.prs, slide, 6, str(convert_currency(future_annual_expense)))
        update_text_of_textbox(self.prs, slide, 8, "6%")
        update_text_of_textbox(self.prs, slide, 10, str(convert_currency(expected_retirement_corpus_needed)))
        update_text_of_textbox(self.prs, slide, 12, str(convert_currency(fv_of_retirement_schemes)))

        if corpus_gap > 0:
            update_text_of_textbox(self.prs, slide, 14, f"{convert_currency(corpus_gap)}")
            formatted_note = self._format_retirement_commentary(ret_note)
            self._write_formatted_commentary(slide, 17, formatted_note)
        else:
            formatted_note = self._format_retirement_commentary(ret_note)
            self._write_formatted_commentary(slide, 17, formatted_note)
            delete_text_box(self.prs, slide, 14)
            delete_text_box(self.prs, slide, 15)

        return corpus_gap

    def build_retirement_goal_planning_slide(self, llm_azure, slide=14):
        """Build Slide 14 (Retirement Funding Strategy) - investment table and LLM commentary."""
        # Compute corpus gap
        retirement_info = self.final_state.get('required_retirement_corpus', {})
        expected_corpus = retirement_info.get('recommendation', {}).get('recommended_corpus', 0)
        fv_of_schemes = self.final_state['retirement_schemes_fv']['grand_total']
        ret_corpus_gap = expected_corpus - fv_of_schemes

        if ret_corpus_gap <= 0:
            return False

        print(f"Retirement Planning: slide: {slide}")

        current_year = date.today().year
        years_to_retire = retirement_info.get('client_info', {}).get('years_to_retirement', 0)
        retirement_year = current_year + years_to_retire

        # Find retirement goal
        retirement_plan = {}
        for x in self.final_state['optimal_goal_allocation']['goals']:
            if x['goal_name'] == "Retirement":
                retirement_plan = x
                break

        # Categorize funding sources
        freed_funds = []
        sip_funds = []
        lumpsum_funds = []
        retirement_assets = []
        for allocation in retirement_plan.get('funded_from', []):
            alloc_type = allocation.get('type', "None")
            if alloc_type == "freed_sip":
                freed_funds.append(allocation)
            elif alloc_type in ('sip_from_surplus', 'sip_from_partial_surplus'):
                sip_funds.append(allocation)
            elif alloc_type == 'lumpsum_from_liquid':
                lumpsum_funds.append(allocation)

        for source in retirement_plan.get('sourced_from', []):
            if source.get('source') == 'future_values_retirement_investments':
                retirement_assets.append(source)
                

        # Build table columns
        col1_type = ["Investment Type"]
        col2_investment = ["Investment"]
        col3_from = ["From Year"]
        col4_to = ["To Year"]
        col5_rate = ["ROI"]
        col7_fv = ["Future Value"]

        for instrument in freed_funds:
            col1_type.append("EMI freed monthly")
            col2_investment.append(convert_currency(instrument["monthly"]))
            col3_from.append(instrument['from_year'])
            col4_to.append(instrument['to_year'])
            col5_rate.append(instrument['rate'])
            col7_fv.append(convert_currency(instrument['fv_contribution']))

        for instrument in sip_funds:
            col1_type.append("Monthly SIP")
            col2_investment.append(convert_currency(instrument["monthly"]))
            col3_from.append(instrument['from_year'])
            col4_to.append(instrument['to_year'])
            col5_rate.append(instrument['rate'])
            col7_fv.append(convert_currency(instrument['fv_contribution']))

        for instrument in lumpsum_funds:
            col1_type.append("Lumpsum")
            col2_investment.append(convert_currency(instrument["principal_used_today"]))
            col3_from.append(instrument['from_year'])
            col4_to.append(instrument['to_year'])
            col5_rate.append(instrument['rate'])
            col7_fv.append(convert_currency(instrument['fv_contribution']))

        for instrument in retirement_assets:
            amount = instrument.get('amount', 0)
            col1_type.append("Retirement Assets")
            col2_investment.append(convert_currency(amount))
            col3_from.append("-")
            col4_to.append("-")
            col5_rate.append("-")
            col7_fv.append(convert_currency(amount))

        # Total future value
        total_fv = (
            sum(i['fv_contribution'] for i in freed_funds) +
            sum(i['fv_contribution'] for i in sip_funds) +
            sum(i['fv_contribution'] for i in lumpsum_funds) +
            sum(i.get('amount', 0) for i in retirement_assets)
        )

        # LLM schema and structured output
        class slide14_llm(BaseModel):
            retirement_funding_commentary: str = Field(
                description=(
                    "A concise commentary (max 80 words) explaining the retirement funding strategy. "
                    "Mention: (1) Freed EMI funds from specific loan closures and their contribution, "
                    "(2) Monthly SIP from surplus income if any, "
                    "(3) Lumpsum investments from liquid assets if any, "
                    "(4) Overall funding adequacy status. "
                    "Be specific about loan types and amounts where applicable."
                ))

        slide14_structured_llm = llm_azure.with_structured_output(slide14_llm)

        # Build LLM prompt
        prompt = (
            f"You are analyzing a retirement funding strategy for a client who needs a retirement corpus of {convert_currency(ret_corpus_gap)} "
            f"by year {retirement_year} (in {years_to_retire} years).\n\n"
            f"The funding strategy uses the following sources:\n\n"
        )

        if freed_funds:
            prompt += "1. FREED EMI FUNDS (from loan closures):\n"
            for instrument in freed_funds:
                loan_type = "loan"
                if 'optimal_loan_allocation' in self.final_state and 'per_loan' in self.final_state['optimal_loan_allocation']:
                    for loan in self.final_state['optimal_loan_allocation']['per_loan']:
                        if 'freed_by_year' in loan and instrument['from_year'] in loan.get('freed_by_year', {}):
                            loan_type = loan.get('type', 'loan')
                            break
                prompt += (
                    f"   - From {loan_type}: ₹{convert_currency(instrument['monthly'])}/month "
                    f"from {instrument['from_year']} to {instrument['to_year']} "
                    f"({instrument['months']} months at {instrument['rate']}) "
                    f"→ Future Value: ₹{convert_currency(instrument['fv_contribution'])}\n"
                )

        if sip_funds:
            prompt += "\n2. MONTHLY SIP FROM SURPLUS INCOME:\n"
            for instrument in sip_funds:
                prompt += (
                    f"   - ₹{convert_currency(instrument['monthly'])}/month "
                    f"from {instrument['from_year']} to {instrument['to_year']} "
                    f"({instrument['months']} months at {instrument['rate']}) "
                    f"→ Future Value: ₹{convert_currency(instrument['fv_contribution'])}\n"
                )

        if lumpsum_funds:
            prompt += "\n3. LUMPSUM FROM LIQUID ASSETS:\n"
            for instrument in lumpsum_funds:
                prompt += (
                    f"   - ₹{convert_currency(instrument['principal_used_today'])} invested today "
                    f"for {instrument['years']} years (until {instrument['to_year']}) "
                    f"at {instrument['rate']} "
                    f"→ Future Value: ₹{convert_currency(instrument['fv_contribution'])}\n"
                )

        if retirement_assets:
            prompt += "\n4. EXISTING RETIREMENT ASSETS:\n"
            for instrument in retirement_assets:
                amount = instrument.get('amount', 0)
                prompt += (
                    f"   - {instrument.get('source', 'retirement_assets')}: "
                    f"Rs. {convert_currency(amount)}\n"
                )

        prompt += f"\n--- SUMMARY ---\n"
        prompt += f"Total Future Value from all sources: ₹{convert_currency(total_fv)}\n"
        prompt += f"Required Retirement Corpus: ₹{convert_currency(ret_corpus_gap)}\n"

        if total_fv >= ret_corpus_gap * 0.90:
            funding_percentage = round((total_fv / ret_corpus_gap) * 100, 1)
            prompt += f"Status: Goal is ADEQUATELY FUNDED (achieves {funding_percentage}% of target)\n"
        else:
            gap_left = ret_corpus_gap - total_fv
            sip_required_ = sip_required(gap_left, 0.085, years_to_retire * 12)
            funding_percentage = round((total_fv / ret_corpus_gap) * 100, 1)
            prompt += (
                f"Status: PARTIAL FUNDING (achieves {funding_percentage}% of target)\n"
                f"Remaining Gap: ₹{convert_currency(gap_left)}\n"
                f"Additional monthly SIP needed: ₹{convert_currency(sip_required_)} at 8.5% return\n"
            )

        prompt += (
            "\nBased on the above information, write a concise, client-friendly commentary (MAXIMUM 80 WORDS) that:\n"
            "- Mentions the specific funding sources used (freed EMIs with loan types, SIPs, lumpsums, retirement assets)\n"
            "- Highlights key amounts and timelines\n"
            "- Clearly states whether the retirement goal is adequately funded or if a gap remains\n"
            "- Uses professional but accessible language suitable for a financial presentation\n"
            "- MUST be written as a SINGLE continuous paragraph (NO paragraph breaks, NO \\n\\n)\n"
            "- **DO NOT use bullet points or multiple paragraphs\n**"
        )

        # Get LLM commentary and populate slide
        slide14_commentary = slide14_structured_llm.invoke(prompt).retirement_funding_commentary
        print(f"slide14_commentary: {slide14_commentary}")

        #ret_table_data = {
        #    "col1": col1_type, "col2": col2_investment, "col3": col3_from,
        #    "col4": col4_to, "col5": col5_rate, "col6": col6_months, "col7": col7_fv
        #}
        #find_and_update_table(self.prs, slide, ret_table_data)
        #indx_text_boxes(self.prs, slide)

        data = list(zip(col1_type, col2_investment, col3_from,
                col4_to, col5_rate, col7_fv))

        slide_obj = self.prs.slides[slide]
        add_table_to_slide(
    slide=slide_obj,
    data=data,
    left=Inches(0.94),
    top=Inches(2.5),
    width=Inches(12.7),
    height=Inches(1.0 * len(data)),
    font_size=Pt(22),
    header_fill=RGBColor(47, 85, 151),
    font_color=RGBColor(0, 0, 0),
    header_font_color=RGBColor(255, 255, 255),)
        final_retirement_note = self._format_retirement_commentary(slide14_commentary)
        self._write_formatted_commentary(slide, 2, final_retirement_note)

        return True
    
    def build_liabilities_slide(self, slide=7):
        """Build Liabilities slide - dynamically populates liabilities table.

        Conditional: only runs if liabilities exist.
        Returns True if populated, False if slide should be deleted.
        """
        liabilities = self.final_state.get('liabilities', [])

        if not liabilities:
            return False

        print(f"Liabilities: slide: {slide}")

        data = [
            ["Loan Type", "Outstanding Balance", "Interest Rate", "EMI/Month", "Penalty Period"]
        ]

        for liability in liabilities:
            loan_type = liability.get('type', '')
            outstanding = convert_currency(liability.get('outstanding_balance', 0))
            interest_rate = f"{round(liability.get('interest_rate', 0) * 100, 2)}%"
            emi = convert_currency(liability.get('emi_amount', 0))

            is_penalty = liability.get('is_under_penalty_period', False)
            penalty_months = liability.get('time_left_to_come_out_of_penalty_period(months)', 0)
            if is_penalty and penalty_months > 0:
                penalty_text = f"Yes ({penalty_months} months)"
            elif is_penalty:
                penalty_text = "Yes"
            else:
                penalty_text = "No"

            data.append([loan_type, outstanding, interest_rate, emi, penalty_text])

        # Remove existing template table before adding dynamic one
        slide_obj = self.prs.slides[slide]
        for shape in slide_obj.shapes:
            if shape.has_table:
                sp = shape._element
                sp.getparent().remove(sp)
                break

        add_table_to_slide( slide=slide_obj,data=data,
            left=Inches(1.5), top=Inches(3),
            width=Inches(16.5), height=Inches(1.0 * len(data)),
            font_size=Pt(24),header_fill=RGBColor(47, 85, 151),
            header_font_color=RGBColor(255, 255, 255),font_color=RGBColor(0, 0, 0))

        return True

    def build_ulips_slide(self, slide=17):
        """Build insurance slide - ULIP table followed by LIC table, each rendered only if data exists."""
        client_data = self.final_state.get('client_data', {})

        investment_details = client_data.get('investment_details', {})
        retirement_investments = investment_details.get('retirement_investments', {})

        ulip_data = retirement_investments.get('ulip', [])
        lic_data  = investment_details.get('lic_policies', [])

        if not ulip_data and not lic_data:
            return False

        print(f"ULIPs/LIC: slide: {slide}")

        slide_obj = self.prs.slides[slide]

        TABLE_LEFT      = Inches(1.5)
        TABLE_WIDTH     = Inches(16.5)
        ROW_HEIGHT_IN   = 1.0          # inches per row
        HEADING_H_IN    = 0.55         # height of each section heading textbox
        HEADING_GAP_IN  = 0.1          # gap between heading and its table
        SECTION_GAP_IN  = 0.6          # gap between sections
        START_TOP_IN    = 3.0          # top of first heading

        next_top = START_TOP_IN        # tracks vertical cursor

        def _add_section_heading(text):
            nonlocal next_top
            txBox = slide_obj.shapes.add_textbox(TABLE_LEFT, Inches(next_top), TABLE_WIDTH, Inches(HEADING_H_IN))
            tf = txBox.text_frame
            tf.word_wrap = False
            para = tf.paragraphs[0]
            run = para.add_run()
            run.text = text
            run.font.bold = True
            run.font.size = Pt(32)
            run.font.color.rgb = RGBColor(0, 0, 0)
            next_top += HEADING_H_IN + HEADING_GAP_IN

        # ── ULIP table ──────────────────────────────────────────────────────
        if ulip_data:
            _add_section_heading("ULIP(S)")

            ulip_rows = [
                ["Policy Name", "Started Since", "Premium/Yr(₹)", "PPT", "Maturity Year", "Maturity Amount", "Comments"]
            ]
            for policy in ulip_data:
                ulip_rows.append([
                    policy.get('policy_name', ''),
                    policy.get('commencement_date', ''),
                    convert_currency(policy.get('annual_premium', 0)),
                    str(policy.get('premium_payment_term', '')),
                    str(policy.get('maturity_year', '')),
                    convert_currency(policy.get('maturity_value', 0)), ''])

            add_table_to_slide(
                slide=slide_obj,
                data=ulip_rows,
                left=TABLE_LEFT,
                top=Inches(next_top),
                width=TABLE_WIDTH,
                height=Inches(ROW_HEIGHT_IN * len(ulip_rows)),
                font_size=Pt(24),
                header_fill=RGBColor(128, 128, 128),
                header_font_color=RGBColor(255, 255, 255),
                font_color=RGBColor(0, 0, 0),
            )
            next_top += ROW_HEIGHT_IN * len(ulip_rows) + SECTION_GAP_IN

        # ── LIC / insurance policies table ───────────────────────────────────
        if lic_data:
            _add_section_heading("INSURANCE POLICIES")

            lic_rows = [["Policy Name", "Started Since", "Premium/Yr(₹)", "PPT", "Policy Period (Yrs)", "Maturity Amount", "Comments"]]
            for policy in lic_data:
                lic_rows.append([
                    policy.get('policy_name', ''),
                    policy.get('commencement_date', ''),
                    convert_currency(policy.get('premium', 0)),
                    str(policy.get('ppt', '')),
                    str(policy.get('policy_period', '')),
                    convert_currency(policy.get('maturity_value', 0)),
                    ''
                ])

            add_table_to_slide(
                slide=slide_obj,
                data=lic_rows,
                left=TABLE_LEFT,
                top=Inches(next_top),
                width=TABLE_WIDTH,
                height=Inches(ROW_HEIGHT_IN * len(lic_rows)),
                font_size=Pt(24),
                header_fill=RGBColor(70, 130, 180),
                header_font_color=RGBColor(255, 255, 255),
                font_color=RGBColor(0, 0, 0), )

        return True

        # Populate commentary textbox
        #opportunity_cost = self.final_state.get('ulip_opportunity_cost', [])
        #if opportunity_cost:     #    # Find the "*ULIPS COMMENTRY" textbox and update it
        #    indx_text_boxes(self.prs, slide)
        #    text_shapes = [s for s in slide_obj.shapes if s.has_text_frame]
        #    for idx, shape in enumerate(text_shapes):
        #        if "COMMENTRY" in shape.text_frame.text.upper() or "COMMENTARY" in shape.text_frame.text.upper():
        #            update_text_of_textbox(self.prs, slide, idx, str(commentary))
        #            break

        return True
        #    commentary = opportunity_cost[0] if isinstance(opportunity_cost, list) else str(opportunity_cost)
   

    def _get_marriage_goals(self):
        """
        Return a list of dicts for every marriage goal found in optimal_goal_allocation.
        Each dict has: goal_name, child_name, current_cost, target_corpus, target_year, funded_from.
        """
        raw_lookup = {
            g['goal_name']: g.get('capital_required_today', 0)
            for g in self.final_state.get('client_data', {}).get('financial_goals', [])
        }

        marriage_goals = []
        for goal in self.final_state.get('optimal_goal_allocation', {}).get('goals', []):
            if str(goal.get('goal_name', '')).strip().endswith('Marriage'):
                child_name = goal['goal_name'].rsplit(' ', 1)[0]
                marriage_goals.append({
                    'goal_name':    goal['goal_name'],
                    'child_name':   child_name,
                    'current_cost': raw_lookup.get(goal['goal_name'], 0),
                    'target_corpus': goal.get('target_corpus', 0),
                    'target_year':  goal.get('target_year', '-'),
                    'funded_from':  goal.get('funded_from', []),
                })
        return marriage_goals

    def build_marriage_slide(self, slide=18):
      """
      Build one marriage slide per child with hardcoded table positioning.
      """
      marriage_goals = self._get_marriage_goals()

      if not marriage_goals:
        print("No marriage goals found - skipping marriage slide.")
        return False, 0

      marriage_offset = 0
      for idx, mg in enumerate(marriage_goals):
        current_slide = slide + idx

        # Duplicate template for every child after the first
        if idx > 0:
            duplicate_slide(self.prs, slide)
            move_slide(self.prs, len(self.prs.slides) - 1, current_slide)
            self.slide_offset += 1
            marriage_offset += 1

        slide_obj = self.prs.slides[current_slide]


        # Slide heading: "{CHILD NAME} MARRIAGE - {MARRIAGE YEAR}"
        heading_text = f"{mg['child_name']} MARRIAGE - {mg['target_year']}"
        update_text_of_textbox(self.prs, current_slide, 1,
                               heading_text,
                               font_size=40, bold=True)

        # Current Estimated Cost
        update_text_of_textbox(self.prs, current_slide, 2,
                               convert_currency(mg['current_cost']),
                               font_size=28, bold=True)
        # Inflation (Hardcoded to 6%)
        update_text_of_textbox(self.prs, current_slide, 3,
                               "6%",
                               font_size=28, bold=True)

        # --- 2. Prepare Table Data ---
        FUND_MAP = {
            'sip_from_surplus': 'SIP',
            'sip_from_partial_surplus': 'SIP',
            'freed_sip': 'Freed EMI',
            'ssy_funds': 'SSY',
            'lumpsum_from_liquid': 'Lumpsum',
            'lumpsum_from_liquid_partial': 'Lumpsum',
            'esop_funds': 'ESOP',
        }

        funding_rows = []
        for fund in mg['funded_from']:
            ftype = fund.get('type', '')
            label = FUND_MAP.get(ftype)
            if not label: continue

            if ftype in ('sip_from_surplus', 'sip_from_partial_surplus', 'freed_sip'):
                amount_str = convert_currency(fund.get('monthly', 0)) + '/month'
            else:
                amount = fund.get('amount_used') or fund.get('principal_used_today') or 0
                amount_str = convert_currency(amount)

            from_year = str(fund.get('from_year', '-'))
            to_year   = str(fund.get('to_year', '-'))
            fv        = fund.get('fv_contribution', fund.get('amount_used', 0))
            fv_str    = convert_currency(fv) if fv else '-'

            funding_rows.append([label, amount_str, from_year, to_year, fv_str])

        table_data = [["Investment Type", "Amount", "From", "To", "Future Value"]]
        table_data.extend(funding_rows if funding_rows else [["-", "-", "-", "-", "-"]])

        add_table_to_slide(slide_obj, data=table_data, left=Inches(7.35), top=Inches(6), width=Inches(11.0),
            height=Inches(0.75 * len(table_data)), font_size=Pt(24), header_fill=RGBColor(91, 178, 192),
            body_fill=RGBColor(255, 255, 255), border_color="5BB2C0", border_width="12000",
            header_font_color=RGBColor(255, 255, 255))

      return True, marriage_offset

    def _get_misc_goals(self):
        """
        Return all goals that are NOT Retirement, NOT education (UG/PG), and NOT Marriage.
        Each dict has: goal_name, current_cost, target_corpus, target_year, funded_from.
        """
        raw_lookup = {
            g['goal_name']: g.get('capital_required_today', 0)
            for g in self.final_state.get('client_data', {}).get('financial_goals', [])
        }

        misc_goals = []
        for goal in self.final_state.get('optimal_goal_allocation', {}).get('goals', []):
            name = str(goal.get('goal_name', '')).strip()
            parts = name.split()
            # Exclude Retirement, education (UG/PG), and Marriage goals
            if name == 'Retirement':
                continue
            if parts and parts[-1] in ('UG', 'PG'):
                continue
            if parts and parts[-1] == 'Marriage':
                continue
            misc_goals.append({
                'goal_name':    name,
                'current_cost': raw_lookup.get(name, 0),
                'target_corpus': goal.get('target_corpus', 0),
                'target_year':  goal.get('target_year', '-'),
                'funded_from':  goal.get('funded_from', []),
                'note':         goal.get('note', []),
            })
        return misc_goals

    def build_misc_goals_slides(self, template_slide):
        """
        Build one slide per miscellaneous goal (vacation, flat, etc.)
        using the same styling as the marriage slide.

        The template_slide index points to the stencil slide with {GOAL_NAME} placeholder.
        For the first goal the template is used directly; additional goals get duplicated slides.

        Returns (has_misc, misc_offset) — True/int if slides were built, False/0 if none.
        """
        FUND_MAP = {
            'sip_from_surplus':         'SIP',
            'sip_from_partial_surplus':  'SIP',
            'freed_sip':                 'Freed EMI',
            'ssy_funds':                 'SSY',
            'lumpsum_from_liquid':       'Lumpsum',
            'lumpsum_from_liquid_partial': 'Lumpsum',
            'esop_funds':                'ESOP',
        }

        misc_goals = self._get_misc_goals()

        if not misc_goals:
            print("No miscellaneous goals found - skipping misc goals slides.")
            return False, 0

        print(f"Building misc goal slides for: {[g['goal_name'] for g in misc_goals]}")

        misc_offset = 0
        for idx, mg in enumerate(misc_goals):
            current_slide = template_slide + idx

            # Duplicate the template for every goal after the first
            if idx > 0:
                duplicate_slide(self.prs, template_slide)
                move_slide(self.prs, len(self.prs.slides) - 1, current_slide)
                self.slide_offset += 1
                misc_offset += 1

            slide_obj = self.prs.slides[current_slide]

            target_amount_str = convert_currency(mg['current_cost']) if mg['current_cost'] else convert_currency(mg['target_corpus'])
            target_year_str   = str(mg['target_year'])

            # --- 1. Update goal name title (textbox index 1) ---
            update_text_of_textbox(self.prs, current_slide, 1, mg['goal_name'], font_size=40, bold=True)

            # --- 2. Update Target Amount value (textbox index 3) and Target Year value (textbox index 5) ---
            # The stencil layout has: [1]=title, [2]="Target Amount" label, [3]=value,
            #                         [4]="Target Year" label, [5]=value
            update_text_of_textbox(self.prs, current_slide, 2, target_amount_str, font_size=28, bold=True)
            update_text_of_textbox(self.prs, current_slide, 3, target_year_str, font_size=28, bold=True)

            # --- 3. Build funding table or unfunded note ---
            funding_rows = []
            for fund in mg['funded_from']:
                ftype = fund.get('type', '')
                label = FUND_MAP.get(ftype)
                if not label:
                    continue
                if ftype in ('sip_from_surplus', 'sip_from_partial_surplus', 'freed_sip'):
                    amount_str = convert_currency(fund.get('monthly', 0)) + '/month'
                else:
                    amount = fund.get('amount_used') or fund.get('principal_used_today') or 0
                    amount_str = convert_currency(amount)
                from_year = str(fund.get('from_year', '-')) if fund.get('from_year') else '-'
                to_year   = str(fund.get('to_year', '-'))   if fund.get('to_year')   else '-'
                fv_str    = convert_currency(fund.get('fv_contribution', 0)) if fund.get('fv_contribution') else '-'
                funding_rows.append([label, amount_str, from_year, to_year, fv_str])

            if not funding_rows:
                # Extract the SIP recommendation from the note field
                note_list = mg.get('note', [])
                note_text = note_list[0] if note_list else ''
                # Pull out just the SIP sentence if present, else show full note
                import re
                sip_match = re.search(r'(You will have to start SIP.*)', note_text, re.IGNORECASE)
                display_note = sip_match.group(1).strip() if sip_match else note_text.strip()
                if not display_note:
                    display_note = "Insufficient funds — no allocation could be made for this goal."

                note_box = slide_obj.shapes.add_textbox(
                    Inches(6.6), Inches(7.5), Inches(10.5), Inches(1.2)
                )
                tf = note_box.text_frame
                tf.word_wrap = True
                p = tf.paragraphs[0]
                p.alignment = PP_ALIGN.LEFT
                run = p.add_run()
                run.text = f"Note: {display_note}"
                run.font.name = "Calibri"
                run.font.size = Pt(22)
                run.font.bold = False
                run.font.color.rgb = RGBColor(0x00, 0x00, 0x00)
            else:
                table_data = [["Investment Type", "Amount", "From Year", "To Year", "Future Value"]]
                table_data.extend(funding_rows)
                add_table_to_slide(
                    slide_obj,
                    data=table_data,
                    left=Inches(6.6),
                    top=Inches(7.5),
                    width=Inches(10.5),
                    height=Inches(0.75 * len(table_data)),
                    font_size=Pt(24),
                    header_fill=RGBColor(91, 178, 192),
                    body_fill=RGBColor(255, 255, 255),
                    border_color="5BB2C0",
                    border_width="12000",
                    header_font_color=RGBColor(255, 255, 255),
                )

            print(f"  Misc goal slide [{idx}]: '{mg['goal_name']}' -> slide {current_slide}")

        return True, misc_offset  
    
    def _get_first_ssy_investment(self):
        """Return first SSY record found in children investments, else None."""
        children = self.final_state.get('client_data', {}).get('client_data', {}).get('children', [])
        for child in children:
            for inv in child.get('investments', []) or []:
                if str(inv.get('type', '')).strip().upper() == 'SUKANYA SAMRIDDHI YOJANA':
                    return child, inv
        return None, None

    def build_ssy_scheme_slide(self, slide=17):
        """Populate SSY slide only when SSY data exists. If not found, leave slide unchanged."""
        child, ssy = self._get_first_ssy_investment()
        if not child or not ssy:
            return False

        try:
            commencement = date.fromisoformat(ssy.get('commencement_date'))
        except Exception:
            return False

        maturity_year = commencement.year + 21
        child_dob = child.get('child_dob')
        daughter_age_at_maturity = "-"
        if child_dob:
            try:
                dob = date.fromisoformat(child_dob)
                daughter_age_at_maturity = str(maturity_year - dob.year)
            except Exception:
                daughter_age_at_maturity = "-"

        current_value = ssy.get('current_value', 0)
        annual_contribution = ssy.get('annual_contribution', 0)
        roi = ssy.get('interest_rate', 0.082)

        slide_obj = self.prs.slides[slide]
        text_shapes = [s for s in slide_obj.shapes if s.has_text_frame]

        for idx, shape in enumerate(text_shapes):
            lower = shape.text_frame.text.strip().lower()
            if 'maturity year' in lower:
                shape.text_frame.text = f"Maturity Year - {maturity_year}"
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.size = Pt(24)
                        run.font.bold = True
            elif 'daughter age' in lower:
                shape.text_frame.text = f"Daughter Age - {daughter_age_at_maturity}"
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.size = Pt(24)
                        run.font.bold = True

        update_text_of_textbox(self.prs, slide, 9,
                               f"{convert_currency(annual_contribution)}",
                               font_size=32, bold=True)
        update_text_of_textbox(self.prs, slide, 10,
                               f"\{convert_currency(current_value)}",
                               font_size=32, bold=True)
        update_text_of_textbox(self.prs, slide, 11,
                               f"{round(roi * 100, 1)}%",
                               font_size=32, bold=True)

        return True

    def prepare_and_build_education_slides(self, summary_slide=14, details_slide=15, inflation_rate="7%"):
        """
        Prepare edu_detail from final_state and build education slides.
        
        Returns:
            (has_kids, slide_offset) - whether slides were built, and how many extra slides were added
        """
        if not self.final_state['client_data']['client_data']['if_any_kids']:
            return False, 0

        print(f"Education Planning: preparing data...")

        # Build edu_detail from final_state
        edu_detail = []
        for i in self.final_state['client_data']['education_planning_summary']:
            entry = {
                'child_name': i['name'], 'edu_type': i['type'],
                'stream': i['stream'], 'destination': i['destination'],
                'target_year': i['target_year'], 'current_cost': i['current_cost'],
                'future_cost': i['future_cost'],
                'future_value_of_allocated_funds': i['total_future_corpus'],
                'corpus_gap': i['final_gap'], 'sourced_from': i['funded_from']
            }
            # Merge goal allocation data
            for j in self.final_state['optimal_goal_allocation']['goals']:
                if i['name'] + " " + i['type'] == j['goal_name']:
                    print(f"j: {j}")
                    entry.update({
                        'final_gap': j['corpus_gap'],
                        'status': j['filter'][0]['type'],
                        'deprioritized': j['depriorized'],
                        'note': j['note'],
                        'funded_from': j['funded_from'],
                        'corpus_needed': j['corpus_needed']
                    })
            edu_detail.append(entry)

        # Set on instance so build_education_slides can use it
        self.edu_detail = edu_detail

        # Build the slides (existing method handles per-child duplication + offset tracking)
        slide_offset = self.build_education_slides(
            summary_slide=summary_slide,
            details_slide=details_slide,
            inflation_rate=inflation_rate )

        return True, slide_offset
    
    def build_epf_slide(self, slide=11):
        """Build EPF (Provident Fund) slide - optional, only if EPF data exists.

        Populates retirement year, retirement age, current value, monthly contribution,
        interest rate, and projected future value from retirement_schemes_fv.

        Returns True if populated, False if slide should be deleted.
        """
        epf_fv = self.final_state.get('retirement_schemes_fv', {}).get('category_totals', {}).get('epf', {})
        epf_schemes = self.final_state.get('client_data', {}).get('investment_details', {}).get('retirement_investments', {}).get('epf', [])

        if not epf_fv or not epf_schemes:
            return False

        print(f"EPF: slide: {slide}")

        retirement_info = self.final_state.get('required_retirement_corpus', {})
        client_info = retirement_info.get('client_info', {})
        years_to_retirement = client_info.get('years_to_retirement', 0)
        retirement_age = self.final_state['client_data']['client_data']['retirement_age']
        retirement_year = date.today().year + years_to_retirement

        sc = epf_schemes[0]
        epf_current_value = sc['current_value']
        epf_monthly_contribution = sc['employee_employer_contribution_monthly']
        epf_interest_rate = sc['interest_rate']

        update_text_of_textbox(self.prs, slide, 9, f"{epf_interest_rate * 100}%")
        update_text_of_textbox(self.prs, slide, 10, convert_currency(epf_current_value))
        update_text_of_textbox(self.prs, slide, 11, convert_currency(epf_monthly_contribution))
        update_text_of_textbox(self.prs, slide, 12, convert_currency(epf_fv))
        update_text_of_textbox(self.prs, slide, 13, str(retirement_year))
        update_text_of_textbox(self.prs, slide, 14, str(retirement_age))

        return True
     

    def build_wealth_created_ret_breakdown(self, slide):
        """
        Build the 'Wealth Created' slide showing a breakdown of all asset future values
        at retirement, plus a donut chart of the asset mix.

        Table rows (matching the slide design):
          Asset                    | Corpus Created | Assumed ROI
          Equity Investments       | <fv>           | 11%
          Fixed Deposit            | <fv or ->      | <rate or ->
          Provident Fund (EPF)     | <fv>           | 8.5%
          Public Provident Fund    | <fv>           | 7.5%
          Sukanya Samriddhi Yojana | <fv or ->      | 8.2%
          Cash                     | <fv>           | -
          Real Estate              | <fv>           | 3%
          ─────────────────────────────────────────
          Total Corpus             | <sum>          |
          Less: Real Estate        | <re_fv>        |
          Less: Discounting Future Goal | <goals_corpus> |
          Corpus For Annuity       | <annuity>      |

        Donut chart uses the same per-asset FVs (excluding zero values).
        """
        print(f"Wealth Created Breakdown: slide: {slide}")

        client_data   = self.final_state.get('client_data', {})
        invest_detail = client_data.get('investment_details', {})
        ret_info      = self.final_state.get('required_retirement_corpus', {})
        client_info   = ret_info.get('client_info', {})
        years_to_ret  = client_info.get('years_to_retirement', 0)
        retirement_year = date.today().year + years_to_ret

        # ── 1. Retirement schemes (EPF / PPF / NPS / other) ──────────────────
        schemes_fv      = self.final_state.get('retirement_schemes_fv', {})
        category_totals = schemes_fv.get('category_totals', {})

        epf_fv  = category_totals.get('epf',  0)
        ppf_fv  = category_totals.get('ppf',  0)
        nps_fv  = category_totals.get('nps',  0)

        ret_investments = invest_detail.get('retirement_investments', {})
        epf_rate = f"{ret_investments.get('epf', [{}])[0].get('interest_rate', 0.085) * 100:.1f}%" if ret_investments.get('epf') else "8.5%"
        ppf_rate = f"{ret_investments.get('ppf', [{}])[0].get('interest_rate', 0.075) * 100:.1f}%" if ret_investments.get('ppf') else "7.5%"
        nps_rate = f"{ret_investments.get('nps', [{}])[0].get('expected_corpus_growth_rate', 0.10) * 100:.1f}%" if ret_investments.get('nps') else "10%"
 
        # ── 3. Fixed Deposits ─────────────────────────────────────────────────
        fd_fv   = 0
        fd_rate = "-"
        for asset in self.final_state.get('liquid_assets', []):
            if 'fixed_deposits' in asset:
                fd = asset['fixed_deposits']
                principal  = fd.get('principal_amount', 0)
                rate       = fd.get('interest_rate', 0.065)
                fd_fv     += principal * ((1 + rate) ** years_to_ret)
                fd_rate    = f"{rate * 100:.1f}%"

        # ── 4. SSY ────────────────────────────────────────────────────────────
        ssy_fv   = 0
        ssy_rate = "8.2%"
        for child in client_data.get('client_data', {}).get('children', []):
            for inv in (child.get('investments') or []):
                if str(inv.get('type', '')).strip().upper() == 'SUKANYA SAMRIDDHI YOJANA':
                    ssy_fv += inv.get('future_value', inv.get('current_value', 0))


        # ── 6. Real Estate (3% p.a. growth to retirement) ────────────────────
        real_estate_fv  = 0
        for asset in self.final_state.get('fixed_assets', []):
            if 'real_estate_investment' in asset:
                current_val    = asset['real_estate_investment'].get('current_market_value', 0)
                real_estate_fv += current_val * ((1.03) ** years_to_ret)

        # ── 7. SIP / Lumpsum / Freed-SIP contributions for Retirement from goal allocation ──
        sip_fv_retirement      = 0   # sip_from_surplus / sip_from_partial_surplus only
        freed_sip_fv_retirement = 0  # freed_sip (released EMI redirected to retirement)
        lumpsum_fv_retirement   = 0

        for goal in self.final_state.get('optimal_goal_allocation', {}).get('goals', []):
            if goal.get('goal_name', '').strip().lower() == 'retirement':
                for fund in goal.get('funded_from', []):
                    ftype = fund.get('type', '')
                    fv    = fund.get('fv_contribution', 0)
                    if ftype in ('sip_from_surplus', 'sip_from_partial_surplus'):
                        sip_fv_retirement += fv
                    elif ftype == 'freed_sip':
                        freed_sip_fv_retirement += fv
                    elif ftype in ('lumpsum_from_liquid', 'lumpsum_from_liquid_partial'):
                        lumpsum_fv_retirement += fv

        # ── 7b. Leftover RSU (total usable RSU − consumed by goals, grown to retirement) ──
        RSU_GROWTH_RATE = get_rsu_growth_rate(invest_detail)
        RSU_USABLE_CAP  = 0.60
        rsu_fv = 0.0
        rsu_data = invest_detail.get('rsu', [])
        if rsu_data:
            try:
                market_df = load_rsu_market_data()
            except FileNotFoundError:
                market_df = None

            if market_df is not None:
                total_rsu_value_all = 0.0
                for rsu_entry in rsu_data:
                    ticker = rsu_entry.get('ticker', '').upper()
                    vesting_schedule = rsu_entry.get('vesting_schedule', [])
                    if not ticker or not vesting_schedule:
                        continue
                    ticker_row = market_df[market_df['ticker'] == ticker]
                    if ticker_row.empty:
                        continue
                    price_usd  = float(ticker_row.iloc[0]['price_usd'])
                    usd_to_inr = float(ticker_row.iloc[0]['usd_to_inr_rate'])
                    sorted_sched = sorted(vesting_schedule, key=lambda x: int(x['year']))
                    prev_price_inr = None
                    total_rsu_value = 0.0
                    for i, tranche in enumerate(sorted_sched):
                        vest_year = int(tranche['year'])
                        no_shares = tranche['no_shares']
                        if i == 0:
                            price_inr = price_usd * usd_to_inr
                        else:
                            years_gap = vest_year - int(sorted_sched[i - 1]['year'])
                            price_inr = prev_price_inr * ((1 + RSU_GROWTH_RATE) ** years_gap)
                        total_rsu_value += round(price_inr * no_shares, 2)
                        prev_price_inr = price_inr
                    total_rsu_value_all += total_rsu_value

                # Sum RSU already consumed by goals (funded_from entries with type='rsu_funds')
                rsu_consumed_total = 0.0
                for goal in self.final_state.get('optimal_goal_allocation', {}).get('goals', []):
                    for fund in goal.get('funded_from', []):
                        if fund.get('type') == 'rsu_funds':
                            rsu_consumed_total += fund.get('amount_used', 0)

                rsu_remaining = max(0.0, total_rsu_value_all * RSU_USABLE_CAP - rsu_consumed_total)
                if rsu_remaining > 0:
                    rsu_fv = round(rsu_remaining * ((1 + RSU_GROWTH_RATE) ** years_to_ret), 2)

        # ── 8. Aggregate totals ───────────────────────────────────────────────
        total_corpus = (fd_fv + epf_fv + ppf_fv + nps_fv + ssy_fv + real_estate_fv
                        + sip_fv_retirement + freed_sip_fv_retirement + lumpsum_fv_retirement
                        + rsu_fv)

        # ── 9. Build table rows ───────────────────────────────────────────────
        TEAL  = RGBColor(0, 172, 193)
        WHITE = RGBColor(255, 255, 255)
        BLACK = RGBColor(0, 0, 0)

        table_data = [["Asset", "Corpus Created", "Assumed ROI"]]
        if fd_fv:
            table_data.append(["Fixed Deposit",            convert_currency(fd_fv),                   fd_rate])
        if epf_fv:
            table_data.append(["Provident Fund",           convert_currency(epf_fv),                  epf_rate])
        if ppf_fv:
            table_data.append(["Public Provident Fund",    convert_currency(ppf_fv),                  ppf_rate])
        if nps_fv:
            table_data.append(["NPS",                      convert_currency(nps_fv),                  nps_rate])
        if ssy_fv:
            table_data.append(["Sukanya Samriddhi Yojana", convert_currency(ssy_fv),                  ssy_rate])
        if real_estate_fv:
            table_data.append(["Real Estate",              convert_currency(real_estate_fv),          "3%"])
        if sip_fv_retirement:
            table_data.append(["SIP (Retirement)",         convert_currency(sip_fv_retirement),       "12%"])
        if freed_sip_fv_retirement:
            table_data.append(["Freed EMI (Retirement)",   convert_currency(freed_sip_fv_retirement), "12%"])
        if lumpsum_fv_retirement:
            table_data.append(["Lumpsum (Retirement)",     convert_currency(lumpsum_fv_retirement),   "12%"])
        if rsu_fv:
            table_data.append(["RSU (Leftover)",           convert_currency(rsu_fv),                  "10%"])

        # ── 10. Add table to slide ────────────────────────────────────────────
        slide_obj = self.prs.slides[slide]

        # Remove any existing template table
        for shape in list(slide_obj.shapes):
            if shape.has_table:
                shape._element.getparent().remove(shape._element)
                break

        tbl_shape = add_table_to_slide(
            slide=slide_obj,
            data=table_data,
            left=Inches(1.0),
            top=Inches(2.4),
            width=Inches(9.5),
            height=Inches(0.55 * len(table_data)),
            font_size=Pt(24),
            header_fill=TEAL,
            header_font_color=WHITE,
            font_color=BLACK,
            body_fill=WHITE,
        )

        # ── 11. Donut chart ───────────────────────────────────────────────────
        donut_data = {}
        label_map = {
            "PF":                   epf_fv,
            "PPF":                  ppf_fv,
            "NPS":                  nps_fv,
            "SSY":                  ssy_fv,
            "Real Estate":          real_estate_fv,
            "Fixed Deposit":        fd_fv,
            "SIP (Retirement)":     sip_fv_retirement,
            "Freed EMI (Ret)":      freed_sip_fv_retirement,
            "Lumpsum (Ret)":        lumpsum_fv_retirement,
            "RSU (Leftover)":       rsu_fv,
        }
        for label, val in label_map.items():
            if val and val > 0:
                donut_data[label] = val / total_corpus if total_corpus else 0

        if donut_data:
            replace_pie_chart_with_matplotlib(
                presentation=self.prs,
                slide_number=slide,
                pie_data=donut_data,
                title=f"Wealth Created - {retirement_year}"
            )

        # ── 12. Update slide title and bottom corpus display ──────────────────
        # Title: find the largest text shape (the heading) and update it
        slide_obj = self.prs.slides[slide]
        text_shapes = [s for s in slide_obj.shapes if s.has_text_frame]
        for shape in text_shapes:
            txt = shape.text_frame.text.strip()
            if 'ret year' in txt.lower() or 'wealth created' in txt.lower():
                tf = shape.text_frame
                for para in tf.paragraphs:
                    for run in para.runs:
                        run.text = run.text.replace('{Ret year}', str(retirement_year)) \
                                           .replace('{ret year}', str(retirement_year))
                    if not para.runs:
                        para.text = f"Wealth Created – {retirement_year}"
                break

        # Bottom corpus display: find "{total corpus}" or "Crores" placeholder
        corpus_in_crores = round(total_corpus / 1e7, 2)
        for shape in text_shapes:
            txt = shape.text_frame.text.strip()
            if 'total corpus' in txt.lower() or 'crores' in txt.lower():
                tf = shape.text_frame
                for para in tf.paragraphs:
                    for run in para.runs:
                        run.text = run.text.replace('{total corpus}', str(corpus_in_crores)) \
                                           .replace('{Total corpus}', str(corpus_in_crores))
                    if not para.runs and ('total corpus' in para.text.lower() or 'crores' in para.text.lower()):
                        para.text = f"{corpus_in_crores} Crores"
                break

        return True 
    def build_ppf_slide(self, slide=13):
        """Build PPF (Public Provident Fund) slide - optional, only if PPF data exists.

        Populates retirement year, retirement age, current value, annual contribution,
        interest rate, and projected future value from retirement_schemes_fv.

        Returns True if populated, False if slide should be deleted.
        """
        ppf_fv = self.final_state.get('retirement_schemes_fv', {}).get('category_totals', {}).get('ppf', {})
        ppf_schemes = self.final_state.get('client_data', {}).get('investment_details', {}).get('retirement_investments', {}).get('ppf', [])

        if not ppf_fv or not ppf_schemes:
            return False

        print(f"PPF: slide: {slide}")

        retirement_info = self.final_state.get('required_retirement_corpus', {})
        client_info = retirement_info.get('client_info', {})
        years_to_retirement = client_info.get('years_to_retirement', 0)
        retirement_age = self.final_state['client_data']['client_data']['retirement_age']
        retirement_year = date.today().year + years_to_retirement

        sc = ppf_schemes[0]
        ppf_current_value = sc['current_value']
        ppf_annual_contribution = sc['annual_contribution']
        ppf_interest_rate = sc['interest_rate']

        update_text_of_textbox(self.prs, slide, 9, f"{ppf_interest_rate * 100}%")
        update_text_of_textbox(self.prs, slide, 10, str(retirement_year))
        update_text_of_textbox(self.prs, slide, 11, str(retirement_age))
        update_text_of_textbox(self.prs, slide, 12, convert_currency(ppf_current_value))
        update_text_of_textbox(self.prs, slide, 13, convert_currency(ppf_annual_contribution))
        update_text_of_textbox(self.prs, slide, 6, convert_currency(ppf_fv))

        return True
    
    
    def _populate_rsu_slide(self, slide_idx, rsu):
        """Paint one RSU entry onto the slide at slide_idx. Called once per RSU grant."""
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor

        HEADER_FILL = RGBColor(0x1F, 0x38, 0x64)
        HEADER_COLOR = RGBColor(0xFF, 0xFF, 0xFF)
        BODY_FILL   = RGBColor(0xFF, 0xFF, 0xFF)
        BODY_COLOR  = RGBColor(0x2C, 0x3E, 0x50)
        ALT_FILL    = RGBColor(0xF0, 0xF4, 0xF8)
        BOX_BG      = RGBColor(0xE8, 0xED, 0xF5)
        BOX_BORDER  = RGBColor(0x1F, 0x38, 0x64)

        slide_obj = self.prs.slides[slide_idx]

        # Update heading to "RSU"
        update_text_of_textbox(self.prs, slide_idx, 1, "RSU")

        # Clear any existing tables left over from the template
        for s in [s for s in slide_obj.shapes if s.shape_type == 19]:
            s._element.getparent().remove(s._element)

        ticker    = rsu.get('ticker', '')
        company   = rsu.get('company_name', ticker).upper()
        price_usd = rsu.get('price_usd_today', 0)
        fx_rate   = rsu.get('usd_to_inr_rate', 0)
        total_val = rsu.get('total_rsu_value_inr', 0)
        consumed  = rsu.get('rsu_total_consumed', 0)
        remaining = rsu.get('rsu_remaining', 0)
        tranches  = rsu.get('tranches', [])

        # ── Company info label ────────────────────────────────────────────────
        info_box = slide_obj.shapes.add_textbox(
            Inches(1.2), Inches(1.75), Inches(12.5), Inches(0.55)
        )
        info_tf = info_box.text_frame
        info_tf.word_wrap = False
        info_p = info_tf.paragraphs[0]
        info_p.alignment = PP_ALIGN.LEFT
        info_run = info_p.add_run()
        info_run.text = (
            f"{company} ({ticker})   |   "
            f"Current Price: ${price_usd:,.2f}   |   "
            f"USD / INR: {fx_rate:,.2f}"
        )
        info_run.font.name = "Calibri"
        info_run.font.size = Pt(24)
        info_run.font.bold = True
        info_run.font.color.rgb = BODY_COLOR

        # ── Left table: per-tranche vesting schedule ──────────────────────────
        tranche_table_data = [
            ["Year", "Shares Vesting", "Price / Share (₹)", "Tranche Value (₹)"]
        ]
        for t in tranches:
            tranche_table_data.append([
                str(t.get('year', '')),
                str(t.get('no_shares', '')),
                convert_currency(t.get('price_per_share_inr', 0)),
                convert_currency(t.get('tranche_value_inr', 0)),
            ])

        row_h = 0.78
        self._add_rsu_esop_table(
            slide_obj, tranche_table_data,
            left=Inches(1.2), top=Inches(2.55),
            width=Inches(7.5), height=Inches(row_h * len(tranche_table_data)),
            font_size=Pt(24),
            header_fill=HEADER_FILL, header_font_color=HEADER_COLOR,
            body_fill=BODY_FILL, body_font_color=BODY_COLOR,
            alt_fill=ALT_FILL,
        )

        # ── Right side: 3 stacked stat boxes ─────────────────────────────────
        summary_items = [
            ("Total Portfolio Value", convert_currency(total_val)),
            ("Allocated to Goals",    convert_currency(consumed)),
            ("Remaining (Usable)",    convert_currency(remaining)),
        ]
        box_left = Inches(9.9)
        box_w    = Inches(4.0)
        box_h    = Inches(1.15)
        box_gap  = Inches(0.2)
        box_top  = Inches(2.55)

        for label, value in summary_items:
            shape = slide_obj.shapes.add_shape(1, box_left, box_top, box_w, box_h)
            shape.fill.solid()
            shape.fill.fore_color.rgb = BOX_BG
            shape.line.color.rgb = BOX_BORDER
            shape.line.width = Pt(1.5)
            tf = shape.text_frame
            tf.word_wrap = False
            tf.margin_left   = Inches(0.15)
            tf.margin_right  = Inches(0.1)
            tf.margin_top    = Inches(0.1)
            tf.margin_bottom = Inches(0.05)
            p_label = tf.paragraphs[0]
            p_label.alignment = PP_ALIGN.LEFT
            r_label = p_label.add_run()
            r_label.text = label
            r_label.font.name = "Calibri"
            r_label.font.size = Pt(18)
            r_label.font.bold = False
            r_label.font.color.rgb = BODY_COLOR
            p_value = tf.add_paragraph()
            p_value.alignment = PP_ALIGN.LEFT
            r_value = p_value.add_run()
            r_value.text = value
            r_value.font.name = "Calibri"
            r_value.font.size = Pt(26)
            r_value.font.bold = True
            r_value.font.color.rgb = HEADER_FILL
            box_top += box_h + box_gap

        # ── Goal breakdown — heading outside, box contains only goal lines ──────────────────
        rsu_used_tracker = rsu.get('rsu_used_tracker', {})
        if rsu_used_tracker:
            goal_lines = [(gn, convert_currency(amt)) for gn, amt in rsu_used_tracker.items()]
            table_bottom = Inches(2.55) + Inches(row_h * len(tranche_table_data))
            box_top = table_bottom + Inches(0.5)

            # Heading — sits above the box, outside it
            hdr_box = slide_obj.shapes.add_textbox(Inches(1.2), box_top, Inches(7.5), Inches(0.45))
            hdr_tf = hdr_box.text_frame
            hdr_tf.word_wrap = False
            hdr_p = hdr_tf.paragraphs[0]
            hdr_p.alignment = PP_ALIGN.LEFT
            hdr_r = hdr_p.add_run()
            hdr_r.text = "Goal Breakdown"
            hdr_r.font.name = "Calibri"
            hdr_r.font.size = Pt(28)
            hdr_r.font.bold = True
            hdr_r.font.color.rgb = HEADER_FILL

            # Box — only goal lines, tighter height
            breakdown_h = Inches(0.25 + 0.5 * len(goal_lines))
            bb = slide_obj.shapes.add_shape(1, Inches(1.2), box_top + Inches(0.45), Inches(7.5), breakdown_h)
            bb.fill.solid()
            bb.fill.fore_color.rgb = RGBColor(0xFF, 0xF3, 0xE0)
            bb.line.color.rgb = BOX_BORDER
            bb.line.width = Pt(2.0)
            btf = bb.text_frame
            btf.word_wrap = True
            btf.margin_left = Inches(0.2)
            btf.margin_right = Inches(0.15)
            btf.margin_top = Inches(0.12)
            btf.margin_bottom = Inches(0.08)
            first = True
            for goal_name, amt_str in goal_lines:
                p_g = btf.paragraphs[0] if first else btf.add_paragraph()
                first = False
                p_g.alignment = PP_ALIGN.LEFT
                r_name = p_g.add_run()
                r_name.text = f"{goal_name}:  "
                r_name.font.name = "Calibri"
                r_name.font.size = Pt(22)
                r_name.font.bold = True
                r_name.font.color.rgb = HEADER_FILL
                r_amt = p_g.add_run()
                r_amt.text = amt_str
                r_amt.font.name = "Calibri"
                r_amt.font.size = Pt(22)
                r_amt.font.bold = True
                r_amt.font.color.rgb = RGBColor(0x00, 0xAC, 0xC1)

    def build_rsu_esop_slide(self, slide=21):
        """Build RSU / ESOP summary slide(s).

        - One slide per RSU grant: the template slide is used for the first entry;
          additional entries get a fresh duplicate inserted immediately after.
        - ESOP falls back to a single slide with a vested/unvested table.
        - Returns (True, extra_slides_added) or (False, 0).
        """
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor

        investment_details = self.final_state.get('client_data', {}).get('investment_details', {})

        # ── Collect RSU portfolio ─────────────────────────────────────────────
        rsu_portfolio = []
        for fund_entry in self.final_state.get('goal_funding', []):
            if fund_entry.get('rsu_portfolio'):
                rsu_portfolio = fund_entry['rsu_portfolio']
                break

        esops_data    = investment_details.get('esops', [])
        vested_esop   = sum(e.get('vested_esops_value', 0) for e in esops_data)
        unvested_esop = sum(e.get('unvested_esops_value', 0) for e in esops_data)
        has_esop = (vested_esop + unvested_esop) > 0
        has_rsu  = bool(rsu_portfolio)

        if not has_rsu and not has_esop:
            return False, 0

        extra_slides = 0

        if has_rsu:
            print(f"RSU slide: {len(rsu_portfolio)} grant(s), template at slide {slide}")

            # Pre-duplicate the blank template for every grant beyond the first,
            # BEFORE painting anything — so each copy is a clean blank.
            for i in range(1, len(rsu_portfolio)):
                duplicate_slide(self.prs, slide)
                insert_at = slide + i
                move_slide(self.prs, len(self.prs.slides) - 1, insert_at)
                extra_slides += 1
                self.slide_offset += 1

            # Now paint each grant onto its own slide
            for i, rsu in enumerate(rsu_portfolio):
                self._populate_rsu_slide(slide + i, rsu)

        else:
            # ESOP — single slide, simple table
            HEADER_FILL = RGBColor(0x1F, 0x38, 0x64)
            HEADER_COLOR = RGBColor(0xFF, 0xFF, 0xFF)
            BODY_FILL   = RGBColor(0xFF, 0xFF, 0xFF)
            BODY_COLOR  = RGBColor(0x2C, 0x3E, 0x50)
            ALT_FILL    = RGBColor(0xF0, 0xF4, 0xF8)
            BOX_BG      = RGBColor(0xFF, 0xF3, 0xE0)
            BOX_BORDER  = RGBColor(0x1F, 0x38, 0x64)

            update_text_of_textbox(self.prs, slide, 1, "ESOP")
            slide_obj = self.prs.slides[slide]
            table_data = [
                ["Component", "Value (₹)"],
                ["Vested ESOPs",   convert_currency(vested_esop)],
                ["Unvested ESOPs", convert_currency(unvested_esop)],
                ["Total ESOPs",    convert_currency(vested_esop + unvested_esop)],
            ]
            self._add_rsu_esop_table(
                slide_obj, table_data,
                left=Inches(0.5), top=Inches(1.75),
                width=Inches(6.0), height=Inches(0.65 * len(table_data)),
                font_size=Pt(20),
                header_fill=HEADER_FILL, header_font_color=HEADER_COLOR,
                body_fill=BODY_FILL, body_font_color=BODY_COLOR,
                alt_fill=ALT_FILL,
            )

            # ── Goal breakdown: only goals where ESOP funds were actually used ──
            esop_goal_lines = []
            for goal in self.final_state.get('optimal_goal_allocation', {}).get('goals', []):
                for fund in goal.get('funded_from', []):
                    if fund.get('type') == 'esop_funds':
                        esop_goal_lines.append(
                            (
                                goal['goal_name'],
                                convert_currency(
                                    fund.get('fv_contribution', fund.get('amount_used', 0))
                                ),
                            )
                        )
                        break

            if esop_goal_lines:
                breakdown_h = Inches(0.55 + 0.48 * len(esop_goal_lines))
                bb = slide_obj.shapes.add_shape(1, Inches(0.5), Inches(1.75 + 0.65 * len(table_data) + 0.3),
                                                Inches(6.0), breakdown_h)
                bb.fill.solid()
                bb.fill.fore_color.rgb = BOX_BG
                bb.line.color.rgb = BOX_BORDER
                bb.line.width = Pt(1.5)
                btf = bb.text_frame
                btf.word_wrap = True
                btf.margin_left = Inches(0.15)
                btf.margin_right = Inches(0.1)
                btf.margin_top = Inches(0.1)
                btf.margin_bottom = Inches(0.05)
                p_hdr = btf.paragraphs[0]
                p_hdr.alignment = PP_ALIGN.LEFT
                r_hdr = p_hdr.add_run()
                r_hdr.text = "Goal Breakdown"
                r_hdr.font.name = "Calibri"
                r_hdr.font.size = Pt(20)
                r_hdr.font.bold = True
                r_hdr.font.color.rgb = HEADER_FILL
                for goal_name, amt_str in esop_goal_lines:
                    p_g = btf.add_paragraph()
                    p_g.alignment = PP_ALIGN.LEFT
                    r_g = p_g.add_run()
                    r_g.text = f"  {goal_name}:  {amt_str}"
                    r_g.font.name = "Calibri"
                    r_g.font.size = Pt(20)
                    r_g.font.bold = False
                    r_g.font.color.rgb = BODY_COLOR

        return True, extra_slides   

    def save(self, output_path=None, client_name="Client"):
        """Save the presentation."""
        if output_path is None:
            output_path = os.path.join(
                os.path.dirname(os.path.dirname(__file__)),
                'PPT_io',
                'PPT_output',
                f'{client_name}_financial_plan.pptx'
            )
        self.prs.save(output_path)
        print(f"Saved to: {output_path}")
        return output_path

 

