"""
PowerPoint generation from LangGraph financial-plan workflow state.

Used by the dashboard API (POST /financial-plan/ppt) and the standalone CLI in Main/main.py.
"""

from __future__ import annotations

import io
import os
import re
from datetime import date
from typing import Tuple

from dotenv import load_dotenv
from langchain_openai import AzureChatOpenAI
from pptx import Presentation

from Financial_Planning.Models.llm_schemas import retirement_commentary
from Financial_Planning.Utilities.ppt_builder import PPTBuilder
from Financial_Planning.Utilities.ppt_utilities import (
    delete_slide,
    find_slide_index_by_text,
    indx_text_boxes,
    update_text_of_textbox,
)
from Financial_Planning.Utilities.utility_functions import convert_currency

TEMPLATE_FILENAME = "Financial Plan Stencil v2.pptx"


class PptTemplateError(FileNotFoundError):
    """Raised when the PowerPoint template file is missing."""


class PptDependencyError(RuntimeError):
    """Raised when Azure OpenAI env vars required for PPT LLM slides are missing."""


def _fp_root() -> str:
    return os.path.dirname(os.path.dirname(os.path.abspath(__file__)))


def default_template_path() -> str:
    return os.path.join(_fp_root(), "PPT_io", "PPT_template", TEMPLATE_FILENAME)


def sanitize_client_filename(name: str) -> str:
    """Strip characters invalid in Windows filenames."""
    cleaned = re.sub(r'[\\/:*?"<>|\r\n]+', "", str(name or "Client").strip())
    return cleaned or "Client"


def _build_azure_llm() -> AzureChatOpenAI:
    load_dotenv()
    api_key = os.getenv("AZURE_API_KEY")
    azure_endpoint = os.getenv("AZURE_API_BASE")
    api_version = os.getenv("AZURE_API_VERSION")
    deployment_name = os.getenv("AZURE_DEPLOYMENT_NAME")
    missing = [
        k
        for k, v in {
            "AZURE_API_KEY": api_key,
            "AZURE_API_BASE": azure_endpoint,
            "AZURE_API_VERSION": api_version,
            "AZURE_DEPLOYMENT_NAME": deployment_name,
        }.items()
        if not v
    ]
    if missing:
        raise PptDependencyError(
            f"Missing Azure OpenAI env for PPT generation: {', '.join(missing)}"
        )
    return AzureChatOpenAI(
        api_key=api_key,
        azure_endpoint=azure_endpoint,
        api_version=api_version,
        deployment_name=deployment_name,
        temperature=0,
    )


TERM_COVER_TEMPLATE_SLIDE = 22


def _resolve_term_cover_slide(prs: Presentation, slide_offset: int) -> int:
    """Locate the Term Cover stencil slide after optional inserts/deletions."""
    found = find_slide_index_by_text(
        prs,
        "{term cover heading}",
        "term cover required",
    )
    if found is not None:
        return found
    fallback = TERM_COVER_TEMPLATE_SLIDE + slide_offset
    if 0 <= fallback < len(prs.slides):
        return fallback
    return len(prs.slides) - 3


def populate_presentation(prs: Presentation, final_state: dict, llm_azure: AzureChatOpenAI) -> str:
    """
    Fill an open Presentation with slides from workflow state.
    Returns the client display name (unsanitized).
    """
    retirement_commentary_llm = llm_azure.with_structured_output(retirement_commentary)

    client_name = final_state["client_data"]["client_data"]["name"]
    fixed_assets_pool = final_state["fixed_asset_pool"]
    liquid_asset_pool = final_state["liquid_pool"]
    retirement_asset_pool = sum(
        list(asset.values())[0].get(
            "current_value", list(asset.values())[0].get("maturity_value", 0)
        )
        for asset in final_state.get("retirement_assets", [])
        if list(asset.keys())[0] != "ulip"
    )

    delete_list: list[int] = []

    builder = PPTBuilder(prs, final_state)

    builder.build_intro_slide(slide=0)
    builder.build_personal_details_slide(slide=1)

    slide = 2
    print(f"Blank: slide: {slide}")

    output_portfolio = builder.build_asset_allocation_slide(slide=3)
    builder.build_asset_classification_slide(slide=4)
    total_asset_value = output_portfolio["portfolio_summary"]["total_asset_value"]
    indx_text_boxes(prs, slide)
    update_text_of_textbox(
        prs, slide, 2, "Total Asset Value: " + f"{convert_currency(total_asset_value)}"
    )
    indx_text_boxes(prs, slide)
    slide = 4
    update_text_of_textbox(prs, slide, 2, f"{convert_currency(fixed_assets_pool)}")
    update_text_of_textbox(prs, slide, 3, f"{convert_currency(liquid_asset_pool)}")
    update_text_of_textbox(prs, slide, 4, f"{convert_currency(retirement_asset_pool)}")
    slide = 5
    builder.build_goals_roadmap_slide(slide=5)
    slide = 6
    builder.build_financial_health_checkup_slide(slide=6)
    slide = 9
    builder.build_retirement_outlook_slide(retirement_commentary_llm, slide=9)
    slide = 10
    builder.build_retirement_goal_planning_slide(llm_azure, slide=10)

    has_ulips = builder.build_ulips_slide(slide=16)
    if not has_ulips:
        delete_list.append(16)

    has_ssy = builder.build_ssy_scheme_slide(slide=17)
    if not has_ssy:
        delete_list.append(17)

    has_marriage, _marriage_slide_offset = builder.build_marriage_slide(slide=18)
    if not has_marriage:
        delete_list.append(18)

    has_misc, _misc_slide_offset = builder.build_misc_goals_slides(template_slide=19)
    if not has_misc:
        delete_list.append(19)

    has_kids, _edu_slide_offset = builder.prepare_and_build_education_slides(
        summary_slide=14, details_slide=15
    )
    if not has_kids:
        delete_list.append(14)
        delete_list.append(15)

    has_liabilities = builder.build_liabilities_slide(slide=7)
    if not has_liabilities:
        delete_list.append(7)

    wealth_slide = 20 + builder.slide_offset
    print(f"Wealth Created: slide: {wealth_slide} (template 20 + offset {builder.slide_offset})")
    builder.build_wealth_created_ret_breakdown(slide=wealth_slide)

    slide_offset_before_delete = builder.slide_offset
    delete_list.sort(reverse=True)
    for slideno in delete_list:
        delete_slide(prs, slideno)

    term_cover_slide = _resolve_term_cover_slide(prs, slide_offset_before_delete)
    print(
        f"Term Cover: slide: {term_cover_slide} "
        f"(template {TERM_COVER_TEMPLATE_SLIDE} + offset {slide_offset_before_delete})"
    )
    builder.build_term_cover_slide(slide=term_cover_slide)

    return client_name


def generate_financial_plan_ppt(
    final_state: dict,
    *,
    template_path: str | None = None,
) -> Tuple[bytes, str]:
    """
    Load template, populate slides via PPTBuilder, return (pptx_bytes, safe_filename).

    safe_filename does not include the .pptx extension.
    """
    if not final_state or not isinstance(final_state, dict):
        raise ValueError("workflow_state must be a non-empty dict")

    path = template_path or default_template_path()
    if not os.path.isfile(path):
        raise PptTemplateError(
            f"PowerPoint template not found: {path}. "
            "Ensure Financial_Planning/PPT_io/PPT_template/ is present."
        )

    llm_azure = _build_azure_llm()
    prs = Presentation(path)
    client_name = populate_presentation(prs, final_state, llm_azure)

    buffer = io.BytesIO()
    prs.save(buffer)
    safe_name = sanitize_client_filename(client_name)
    return buffer.getvalue(), f"{safe_name}_financial_plan.pptx"
