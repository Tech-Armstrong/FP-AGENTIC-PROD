"""
PowerPoint Manipulation Utilities - Presentation Generation

What this file does:
This script provides utilities for creating and manipulating PowerPoint presentations.
It handles charts, tables, text boxes, and visualizations for financial plan reports.

What this file contains:
- delete_text_box: Removes text box from slide by index
- replace_donut_chart_with_data: Updates donut chart with new data while preserving styling
- replace_pie_chart_with_data: Updates pie chart with category names and percentages inside slices
- create_goals_roadmap: Creates modern horizontal timeline visualization with goals and target years
- move_slide: Moves slide from one position to another in presentation
- duplicate_slide: Creates exact duplicate of slide including all formatting and images
- find_and_update_table: Finds table in slide and updates it with new data
- find_all_tables_in_slide: Recursively searches for all tables including those in groups
- extract_table_style / apply_cell_style / set_cell_margins / set_cell_border: Table styling utilities
- apply_column_widths: Intelligently sets column widths based on content with min/max constraints
- find_and_update_all_tables / create_table_if_not_exists: Batch table operations
- show_slide_data / indx_text_boxes: Text box inspection utilities
- update_text_of_textbox: Updates text box content while preserving formatting
- find_tables_in_slide: Returns list of all tables in slide
"""

from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_DATA_LABEL_POSITION
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.xmlchemy import OxmlElement
from copy import deepcopy
from pptx.enum.shapes import MSO_SHAPE
import matplotlib
matplotlib.use('Agg')  # non-interactive backend — required when running inside a web server (no main thread GUI)
import matplotlib.pyplot as plt
import matplotlib.patches as patches
from matplotlib.patches import FancyBboxPatch, Circle, FancyArrowPatch
import io
from PIL import Image
import numpy as np
import io
from pptx import Presentation
from lxml import etree
from copy import deepcopy
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

########################################################## delete a slide #######################################################

def delete_slide(prs, slide_number):
    """
    Delete a slide from a PowerPoint presentation.
    
    Parameters:
    -----------
    prs : Presentation
        The presentation object from python-pptx
    slide_number : int
        The slide number to delete (1-based index, i.e., first slide is 1)
    
    Returns:
    --------
    Presentation
        The modified presentation object with the slide deleted
    
    Raises:
    -------
    ValueError
        If slide_number is out of range
    """
    # Convert to 0-based index
    slide_index = slide_number
    
    # Validate slide number
    if slide_index < 0 or slide_index >= len(prs.slides):
        raise ValueError(f"Slide number {slide_number} is out of range. "
                        f"Presentation has {len(prs.slides)} slides.")
    
    # Get the slide ID and relationship ID
    slide = prs.slides[slide_index]
    
    # Access the XML elements
    slide_id = prs.slides._sldIdLst[slide_index]
    
    # Remove the slide from the presentation
    prs.part.drop_rel(slide_id.rId)
    prs.slides._sldIdLst.remove(slide_id)

#################################################### delete a slide #############################################################

#################################### delete a textbox ######################################################################
def delete_text_box(ppt, slide_no, textbox_id):
    """
    Deletes a text box from a slide.
    
    Parameters:
    - ppt: The PowerPoint presentation object
    - slide_no: The slide number (1-indexed)
    - textbox_id: The text box index (1-indexed, as shown by indx_text_boxes)
    
    Returns:
    - True if deletion was successful, False otherwise
    """
    slide = ppt.slides[slide_no]
    
    # Counter for text boxes with text
    text_box_counter = 0
    
    for shape in slide.shapes:
        if shape.has_text_frame and shape.text:
            text_box_counter += 1
            
            # If this is the text box we want to delete
            if text_box_counter == textbox_id:
                # Get the shape element and remove it
                sp = shape.element
                sp.getparent().remove(sp)
                return True
    
    # If we didn't find the text box
    print(f"Text box {textbox_id} not found on slide {slide_no}")
    return False
################################### delete a textbox ######################################################################

############################################# donut chart #####################################################################

from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_DATA_LABEL_POSITION
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor


def create_financial_health_donut(
    presentation,
    slide_index,
    savings_ratio,
    left=Inches(1.0),
    top=Inches(2.0),
    width=Inches(4.5),
    height=Inches(4.5),
):
    """
    Creates a SAFE donut chart for Slide 7 (Financial Health).

    savings_ratio:
        float between 0 and 1 (e.g. 0.35)
    """

    # ---------------- Safety checks ----------------
    if savings_ratio is None:
        raise ValueError("savings_ratio cannot be None")

    savings_ratio = float(savings_ratio)
    if savings_ratio < 0:
        savings_ratio = 0
    if savings_ratio > 1:
        savings_ratio = 1

    expense_ratio = 1 - savings_ratio

    slide = presentation.slides[slide_index]

    # ---------------- Chart data ----------------
    chart_data = CategoryChartData()
    chart_data.categories = ["Savings", "Expenses"]
    chart_data.add_series(
        "Financial Health",
        (savings_ratio, expense_ratio)
    )

    # ---------------- Create chart ----------------
    chart_shape = slide.shapes.add_chart(
        XL_CHART_TYPE.DOUGHNUT,
        left,
        top,
        width,
        height,
        chart_data
    )

    chart = chart_shape.chart

    # ---------------- Title ----------------
    chart.has_title = True
    title_p = chart.chart_title.text_frame.paragraphs[0]
    title_p.text = "Savings vs Expenses"
    title_p.font.size = Pt(18)
    title_p.font.bold = True

    # ---------------- Donut hole ----------------
    chart.plots[0].donut_hole_size = 65

    # ---------------- Data labels ----------------
    plot = chart.plots[0]
    plot.has_data_labels = True
    labels = plot.data_labels
    labels.show_percentage = True
    labels.show_category_name = False
    labels.show_value = False
    labels.position = XL_DATA_LABEL_POSITION.BEST_FIT
    labels.font.size = Pt(12)
    labels.font.bold = True

    # ---------------- Legend ----------------
    chart.has_legend = True
    chart.legend.font.size = Pt(11)

    # ---------------- Colors ----------------
    colors = [
        RGBColor(0, 176, 185),   # Savings (teal)
        RGBColor(231, 76, 60),   # Expenses (red)
    ]

    for i, point in enumerate(chart.series[0].points):
        fill = point.format.fill
        fill.solid()
        fill.fore_color.rgb = colors[i]

    return presentation


###################################### donut chart ##########################################################################

######################################## goals visualization ##########################################################

# def create_goals_roadmap(presentation, slide_number, sorted_goals):
#     """
#     Creates a modern horizontal timeline visualization with goals and their target years.
    
#     Parameters:
#     -----------
#     presentation : pptx.Presentation
#         The presentation object
#     slide_number : int
#         The slide number (0-indexed)
#     sorted_goals : dict
#         Dictionary with goal names as keys and target years as values
#         Example: {'Goal 1': 2028, 'Goal 2': 2032, 'Goal 3': 2034, 'Goal 4': 2037, 'Goal 5': 2038}
#         Goals should be in ascending order of target year
    
#     Returns:
#     --------
#     pptx.Presentation
#         Updated presentation object
#     """
    
#     # Get the slide
#     slide = presentation.slides[slide_number]
    
#     # Create the visualization with transparent background
#     fig, ax = plt.subplots(figsize=(16, 6), facecolor='none')
#     ax.set_facecolor('none')
    
#     # Remove axes
#     ax.set_xlim(0, 10)
#     ax.set_ylim(0, 5)
#     ax.axis('off')
    
#     # Define colors - gradient from teal to blue
#     colors = ['#1ABC9C', '#16A085', '#3498DB', '#2980B9', '#8E44AD']
    
#     # Calculate positions for goals
#     goals_list = list(sorted_goals.items())
#     num_goals = len(goals_list)
    
#     if num_goals == 0:
#         return presentation
    
#     # Calculate spacing
#     spacing = 8 / max(1, num_goals - 1) if num_goals > 1 else 0
#     start_x = 1
#     y_center = 2.5
    
#     # Draw the main timeline line with gradient effect
#     if num_goals > 1:
#         for i in range(num_goals - 1):
#             x1 = start_x + i * spacing
#             x2 = start_x + (i + 1) * spacing
            
#             # Use gradient colors for the line segments
#             line_color = colors[i % len(colors)]
            
#             ax.plot([x1, x2], [y_center, y_center], 
#                     color=line_color, linewidth=5, zorder=1, 
#                     solid_capstyle='round', alpha=0.6)
    
#     # Draw circles and milestones
#     for i, (goal_name, target_year) in enumerate(goals_list):
#         x = start_x + i * spacing
        
#         # Choose color
#         color = colors[i % len(colors)]
        
#         # Draw outer glow circle
#         glow = Circle((x, y_center), 0.52, color=color, alpha=0.2, zorder=2)
#         ax.add_patch(glow)
        
#         # Draw main circle
#         circle = Circle((x, y_center), 0.45, color=color, ec='white', linewidth=4, zorder=3)
#         ax.add_patch(circle)
        
#         # Add year inside circle
#         ax.text(x, y_center, str(target_year), 
#                 ha='center', va='center',
#                 fontsize=24, fontweight='bold', 
#                 color='white', fontfamily='sans-serif',
#                 zorder=4)
        
#         # Alternate goal names above and below the timeline
#         if i % 2 == 0:
#             # Place below
#             line_y_start = y_center - 0.45
#             line_y_end = y_center - 0.75
#             text_y = y_center - 1.2
#             va_align = 'center'
#         else:
#             # Place above
#             line_y_start = y_center + 0.45
#             line_y_end = y_center + 0.75
#             text_y = y_center + 1.2
#             va_align = 'center'
        
#         # Draw vertical line from circle to text with matching color
#         ax.plot([x, x], [line_y_start, line_y_end], 
#                 color=color, linewidth=3, zorder=2, alpha=0.8)
        
#         # Add goal name in a rounded box
#         bbox_props = dict(boxstyle="round,pad=0.4", 
#                          facecolor=color, 
#                          edgecolor='white',
#                          linewidth=2,
#                          alpha=0.9)
        
#         ax.text(x, text_y, goal_name,
#                 ha='center', va=va_align,
#                 fontsize=14, fontweight='bold',
#                 color='white', fontfamily='sans-serif',
#                 bbox=bbox_props,
#                 zorder=3)
    
#     # Add arrows between circles with matching colors
#     if num_goals > 1:
#         for i in range(num_goals - 1):
#             x1 = start_x + i * spacing + 0.45
#             x2 = start_x + (i + 1) * spacing - 0.45
            
#             # Use the color of the current circle for the arrow
#             arrow_color = colors[i % len(colors)]
            
#             arrow = FancyArrowPatch((x1, y_center), (x2, y_center),
#                                    arrowstyle='->', mutation_scale=30,
#                                    linewidth=2, 
#                                    color=arrow_color, alpha=0.7, zorder=2)
#             ax.add_patch(arrow)
    
#     # Add a subtle title
#     ax.text(5, 4.2, 'Financial Goals Timeline',
#             ha='center', va='center',
#             fontsize=22, fontweight='bold',
#             color='#2C3E50', fontfamily='sans-serif',
#             alpha=0.8)
    
#     # Save the figure to a bytes buffer with transparent background
#     buf = io.BytesIO()
#     plt.tight_layout()
#     plt.savefig(buf, format='png', dpi=150, bbox_inches='tight', 
#                 facecolor='none', edgecolor='none', transparent=True)
#     buf.seek(0)
#     plt.close()
    
#     # Load image from buffer
#     img = Image.open(buf)
    
#     # Save to temporary buffer for PowerPoint
#     img_buf = io.BytesIO()
#     img.save(img_buf, format='PNG')
#     img_buf.seek(0)
    
#     # Get slide dimensions
#     slide_width = presentation.slide_width
#     slide_height = presentation.slide_height
    
#     # Calculate dimensions to center the image
#     img_width, img_height = img.size
#     aspect_ratio = img_width / img_height
    
#     # Set the image to take up most of the slide width
#     pic_width = slide_width * 0.85
#     pic_height = pic_width / aspect_ratio
    
#     # If height is too large, scale based on height instead
#     if pic_height > slide_height * 0.7:
#         pic_height = slide_height * 0.7
#         pic_width = pic_height * aspect_ratio
    
#     # Calculate centered position
#     left = (slide_width - pic_width) / 2
#     top = (slide_height - pic_height) / 2
    
#     # Add picture to slide
#     slide.shapes.add_picture(img_buf, left, top, pic_width, pic_height)
    
#     buf.close()
#     img_buf.close()
    
#     return presentation

# def create_goals_roadmap(presentation, slide_number, sorted_goals):
#     """
#     Creates a minimal, elegant timeline matching presentation theme
#     """
    
#     # Get the slide
#     slide = presentation.slides[slide_number]
    
#     # Create the visualization with transparent background
#     fig, ax = plt.subplots(figsize=(16, 3.5), facecolor='none')
#     ax.set_facecolor('none')
    
#     # Remove axes
#     ax.set_xlim(0, 10)
#     ax.set_ylim(0, 2.5)
#     ax.axis('off')
    
#     # Theme color - teal matching your table
#     theme_color = '#00B0B9'
#     text_color = '#2C3E50'
    
#     # Calculate positions for goals
#     goals_list = list(sorted_goals.items())
#     num_goals = len(goals_list)
    
#     if num_goals == 0:
#         return presentation
    
#     # Calculate spacing
#     spacing = 8 / max(1, num_goals - 1) if num_goals > 1 else 0
#     start_x = 1
#     y_center = 1.25
    
#     # Draw the main timeline line
#     if num_goals > 1:
#         ax.plot([start_x, start_x + (num_goals - 1) * spacing], 
#                 [y_center, y_center], 
#                 color=theme_color, linewidth=3, zorder=1, 
#                 solid_capstyle='round', alpha=0.4)
    
#     # Draw milestones
#     for i, (goal_name, target_year) in enumerate(goals_list):
#         x = start_x + i * spacing
        
#         # Draw drop shadow
#         shadow = Circle((x, y_center - 0.015), 0.26, color='#000000', alpha=0.08, zorder=2)
#         ax.add_patch(shadow)
        
#         # Draw white circle background
#         bg_circle = Circle((x, y_center), 0.25, color='white', zorder=3)
#         ax.add_patch(bg_circle)
        
#         # Draw teal border circle
#         circle = Circle((x, y_center), 0.23, color='white', 
#                        ec=theme_color, linewidth=3, zorder=4)
#         ax.add_patch(circle)
        
#         # Add year
#         ax.text(x, y_center, str(target_year), 
#                 ha='center', va='center',
#                 fontsize=16, fontweight='bold', 
#                 color=theme_color, fontfamily='sans-serif',
#                 zorder=5)
        
#         # Alternate goal names
#         if i % 2 == 0:
#             text_y = y_center - 0.55
#             va_align = 'top'
#         else:
#             text_y = y_center + 0.55
#             va_align = 'bottom'
        
#         # Clean label design
#         bbox_props = dict(boxstyle="round,pad=0.25", 
#                          facecolor='white', 
#                          edgecolor=theme_color,
#                          linewidth=1.5,
#                          alpha=0.98)
        
#         ax.text(x, text_y, goal_name,
#                 ha='center', va=va_align,
#                 fontsize=10, fontweight='600',
#                 color=text_color, fontfamily='sans-serif',
#                 bbox=bbox_props,
#                 zorder=3)
    
#     # Save with high quality
#     buf = io.BytesIO()
#     plt.tight_layout()
#     plt.savefig(buf, format='png', dpi=200, bbox_inches='tight', 
#                 facecolor='none', edgecolor='none', transparent=True)
#     buf.seek(0)
#     plt.close()
    
#     # Load and add to slide
#     img = Image.open(buf)
#     img_buf = io.BytesIO()
#     img.save(img_buf, format='PNG')
#     img_buf.seek(0)
    
#     slide_width = presentation.slide_width
#     slide_height = presentation.slide_height
    
#     img_width, img_height = img.size
#     aspect_ratio = img_width / img_height
    
#     pic_width = slide_width * 0.88
#     pic_height = pic_width / aspect_ratio
    
#     left = (slide_width - pic_width) / 2
#     top = slide_height * 0.18
    
#     slide.shapes.add_picture(img_buf, left, top, pic_width, pic_height)
    
#     buf.close()
#     img_buf.close()
    
#     return presentation

# def create_goals_roadmap(presentation, slide_number, sorted_goals):
#     """
#     Creates a modern horizontal timeline visualization with goals and their target years.
#     Styled to match presentation theme with subtle, elegant design.
    
#     Parameters:
#     -----------
#     presentation : pptx.Presentation
#         The presentation object
#     slide_number : int
#         The slide number (0-indexed)
#     sorted_goals : dict
#         Dictionary with goal names as keys and target years as values
#         Example: {'Goal 1': 2028, 'Goal 2': 2032, 'Goal 3': 2034, 'Goal 4': 2037, 'Goal 5': 2038}
#         Goals should be in ascending order of target year
    
#     Returns:
#     --------
#     pptx.Presentation
#         Updated presentation object
#     """
    
#     # Get the slide
#     slide = presentation.slides[slide_number]
    
#     # Create the visualization with transparent background
#     fig, ax = plt.subplots(figsize=(16, 4), facecolor='none')
#     ax.set_facecolor('none')
    
#     # Remove axes
#     ax.set_xlim(0, 10)
#     ax.set_ylim(0, 3)
#     ax.axis('off')
    
#     # Define a more subtle, elegant color scheme that works with any background
#     # Using teal/cyan to match your table headers
#     primary_color = '#00B0B9'  # Teal/cyan matching your table header
#     accent_colors = ['#00B0B9', '#00C5CF', '#00939E', '#008C94', '#007B82']
    
#     # Calculate positions for goals
#     goals_list = list(sorted_goals.items())
#     num_goals = len(goals_list)
    
#     if num_goals == 0:
#         return presentation
    
#     # Calculate spacing
#     spacing = 8 / max(1, num_goals - 1) if num_goals > 1 else 0
#     start_x = 1
#     y_center = 1.5
    
#     # Draw the main timeline line - single elegant line
#     if num_goals > 1:
#         ax.plot([start_x, start_x + (num_goals - 1) * spacing], 
#                 [y_center, y_center], 
#                 color='#D3D3D3', linewidth=4, zorder=1, 
#                 solid_capstyle='round', alpha=0.6)
    
#     # Draw circles and milestones
#     for i, (goal_name, target_year) in enumerate(goals_list):
#         x = start_x + i * spacing
        
#         # Use consistent teal color for all circles
#         color = primary_color
        
#         # Draw subtle shadow circle
#         shadow = Circle((x, y_center - 0.02), 0.32, color='#000000', alpha=0.1, zorder=2)
#         ax.add_patch(shadow)
        
#         # Draw white background circle
#         bg_circle = Circle((x, y_center), 0.3, color='white', zorder=3)
#         ax.add_patch(bg_circle)
        
#         # Draw main circle with teal border
#         circle = Circle((x, y_center), 0.28, color='white', 
#                        ec=color, linewidth=3, zorder=4)
#         ax.add_patch(circle)
        
#         # Add year inside circle in teal
#         ax.text(x, y_center, str(target_year), 
#                 ha='center', va='center',
#                 fontsize=18, fontweight='bold', 
#                 color=color, fontfamily='sans-serif',
#                 zorder=5)
        
#         # Alternate goal names above and below the timeline
#         if i % 2 == 0:
#             # Place below
#             text_y = y_center - 0.65
#             va_align = 'top'
#         else:
#             # Place above
#             text_y = y_center + 0.65
#             va_align = 'bottom'
        
#         # Add goal name in a clean, minimal style with white background
#         bbox_props = dict(boxstyle="round,pad=0.3", 
#                          facecolor='white', 
#                          edgecolor=color,
#                          linewidth=2,
#                          alpha=0.95)
        
#         ax.text(x, text_y, goal_name,
#                 ha='center', va=va_align,
#                 fontsize=11, fontweight='600',
#                 color='#2C3E50', fontfamily='sans-serif',
#                 bbox=bbox_props,
#                 zorder=3)
    
#     # Save the figure to a bytes buffer with transparent background
#     buf = io.BytesIO()
#     plt.tight_layout()
#     plt.savefig(buf, format='png', dpi=200, bbox_inches='tight', 
#                 facecolor='none', edgecolor='none', transparent=True)
#     buf.seek(0)
#     plt.close()
    
#     # Load image from buffer
#     img = Image.open(buf)
    
#     # Save to temporary buffer for PowerPoint
#     img_buf = io.BytesIO()
#     img.save(img_buf, format='PNG')
#     img_buf.seek(0)
    
#     # Get slide dimensions
#     slide_width = presentation.slide_width
#     slide_height = presentation.slide_height
    
#     # Calculate dimensions to position in upper portion of slide
#     img_width, img_height = img.size
#     aspect_ratio = img_width / img_height
    
#     # Set the image to take up most of the slide width
#     pic_width = slide_width * 0.9
#     pic_height = pic_width / aspect_ratio
    
#     # Position in the upper half of the slide to avoid background imagery
#     left = (slide_width - pic_width) / 2
#     top = slide_height * 0.15  # Position in upper portion
    
#     # Add picture to slide
#     slide.shapes.add_picture(img_buf, left, top, pic_width, pic_height)
    
#     buf.close()
#     img_buf.close()
    
#     return presentation

# horizontal grid
# def create_goals_roadmap(presentation, slide_number, sorted_goals):
#     """
#     Creates a clean grid of goal cards - ultra minimalist
#     """
#     from pptx.util import Inches, Pt
#     from pptx.enum.shapes import MSO_SHAPE
#     from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
#     from pptx.dml.color import RGBColor
    
#     slide = presentation.slides[slide_number]
    
#     # Theme colors
#     teal_color = RGBColor(0, 176, 185)
#     dark_text = RGBColor(44, 62, 80)
    
#     goals_list = list(sorted_goals.items())
#     num_goals = len(goals_list)
    
#     if num_goals == 0:
#         return presentation
    
#     # Layout - 3 columns or adjust based on number of goals
#     cols = min(3, num_goals)
#     rows = (num_goals + cols - 1) // cols
    
#     card_width = Inches(2.8)
#     card_height = Inches(1.6)
#     start_left = Inches(1.2)
#     start_top = Inches(2.0)
#     h_spacing = Inches(3.0)
#     v_spacing = Inches(1.9)
    
#     for i, (goal_name, target_year) in enumerate(goals_list):
#         row = i // cols
#         col = i % cols
        
#         left = start_left + col * h_spacing
#         top = start_top + row * v_spacing
        
#         # Add card
#         card = slide.shapes.add_shape(
#             MSO_SHAPE.ROUNDED_RECTANGLE,
#             left, top,
#             card_width, card_height
#         )
#         card.fill.solid()
#         card.fill.fore_color.rgb = RGBColor(255, 255, 255)
#         card.line.color.rgb = teal_color
#         card.line.width = Pt(2)
        
#         # Add year badge
#         badge = slide.shapes.add_shape(
#             MSO_SHAPE.ROUNDED_RECTANGLE,
#             left + Inches(0.15), top + Inches(0.15),
#             Inches(0.7), Inches(0.35)
#         )
#         badge.fill.solid()
#         badge.fill.fore_color.rgb = teal_color
#         badge.line.fill.background()
        
#         badge_tf = badge.text_frame
#         badge_tf.text = str(target_year)
#         badge_tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        
#         badge_p = badge_tf.paragraphs[0]
#         badge_p.alignment = PP_ALIGN.CENTER
#         badge_p.font.size = Pt(12)
#         badge_p.font.bold = True
#         badge_p.font.color.rgb = RGBColor(255, 255, 255)
        
#         # Add goal name
#         tf = card.text_frame
#         tf.text = goal_name
#         tf.word_wrap = True
#         tf.margin_top = Inches(0.6)
#         tf.margin_left = Inches(0.2)
#         tf.margin_right = Inches(0.2)
#         tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        
#         p = tf.paragraphs[0]
#         p.alignment = PP_ALIGN.CENTER
#         p.font.size = Pt(13)
#         p.font.bold = True
#         p.font.color.rgb = dark_text
    
#     return presentation

# stepped grid 
# def create_goals_roadmap(presentation, slide_number, sorted_goals):
#     """
#     Creates a stepped/staircase timeline showing upward progression
#     """
#     from pptx.util import Inches, Pt
#     from pptx.enum.shapes import MSO_SHAPE
#     from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
#     from pptx.dml.color import RGBColor
    
#     slide = presentation.slides[slide_number]
    
#     # Theme colors
#     teal_color = RGBColor(0, 176, 185)
#     dark_text = RGBColor(44, 62, 80)
    
#     goals_list = list(sorted_goals.items())
#     num_goals = len(goals_list)
    
#     if num_goals == 0:
#         return presentation
    
#     # Layout
#     step_width = Inches(1.6)
#     step_height = Inches(1.0)
#     start_left = Inches(0.8)
#     start_top = Inches(4.5)
#     horizontal_spacing = Inches(1.7)
    
#     for i, (goal_name, target_year) in enumerate(goals_list):
#         left = start_left + i * horizontal_spacing
#         top = start_top - i * Inches(0.6)  # Steps go up
        
#         # Add step box
#         step = slide.shapes.add_shape(
#             MSO_SHAPE.ROUNDED_RECTANGLE,
#             left, top,
#             step_width, step_height
#         )
#         step.fill.solid()
#         step.fill.fore_color.rgb = RGBColor(255, 255, 255)
#         step.line.color.rgb = teal_color
#         step.line.width = Pt(3)
        
#         # Add year at top
#         year_label = slide.shapes.add_shape(
#             MSO_SHAPE.RECTANGLE,
#             left, top - Inches(0.4),
#             step_width, Inches(0.35)
#         )
#         year_label.fill.solid()
#         year_label.fill.fore_color.rgb = teal_color
#         year_label.line.fill.background()
        
#         year_tf = year_label.text_frame
#         year_tf.text = str(target_year)
#         year_tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        
#         year_p = year_tf.paragraphs[0]
#         year_p.alignment = PP_ALIGN.CENTER
#         year_p.font.size = Pt(14)
#         year_p.font.bold = True
#         year_p.font.color.rgb = RGBColor(255, 255, 255)
        
#         # Add goal name
#         tf = step.text_frame
#         tf.text = goal_name
#         tf.word_wrap = True
#         tf.margin_left = Inches(0.15)
#         tf.margin_right = Inches(0.15)
#         tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        
#         p = tf.paragraphs[0]
#         p.alignment = PP_ALIGN.CENTER
#         p.font.size = Pt(11)
#         p.font.bold = True
#         p.font.color.rgb = dark_text
    
#     return presentation

# vertical milestone card
# def create_goals_roadmap(presentation, slide_number, sorted_goals):
#     """
#     Creates a vertical milestone card timeline - modern and minimalist
#     """
#     from pptx.util import Inches, Pt
#     from pptx.enum.shapes import MSO_SHAPE
#     from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
#     from pptx.dml.color import RGBColor
    
#     # Get the slide
#     slide = presentation.slides[slide_number]
    
#     # Theme colors
#     teal_color = RGBColor(0, 176, 185)
#     dark_text = RGBColor(44, 62, 80)
#     light_gray = RGBColor(236, 240, 241)
    
#     # Calculate positions
#     goals_list = list(sorted_goals.items())
#     num_goals = len(goals_list)
    
#     if num_goals == 0:
#         return presentation
    
#     # Layout parameters
#     card_width = Inches(7)
#     card_height = Inches(0.9)
#     start_left = Inches(1.5)
#     start_top = Inches(1.8)
#     vertical_spacing = Inches(1.1)
    
#     # Add connecting line on the left
#     line_left = Inches(1.0)
#     line_top = start_top + Inches(0.45)
#     line_height = (num_goals - 1) * vertical_spacing
    
#     if num_goals > 1:
#         connector = slide.shapes.add_shape(
#             MSO_SHAPE.RECTANGLE,
#             line_left, line_top,
#             Inches(0.08), line_height
#         )
#         connector.fill.solid()
#         connector.fill.fore_color.rgb = teal_color
#         connector.line.fill.background()
    
#     # Create milestone cards
#     for i, (goal_name, target_year) in enumerate(goals_list):
#         top = start_top + i * vertical_spacing
        
#         # Add year circle on the left
#         circle = slide.shapes.add_shape(
#             MSO_SHAPE.OVAL,
#             line_left - Inches(0.15), top + Inches(0.25),
#             Inches(0.5), Inches(0.5)
#         )
#         circle.fill.solid()
#         circle.fill.fore_color.rgb = teal_color
#         circle.line.color.rgb = RGBColor(255, 255, 255)
#         circle.line.width = Pt(3)
        
#         # Add year text in circle
#         year_tf = circle.text_frame
#         year_tf.text = str(target_year)
#         year_p = year_tf.paragraphs[0]
#         year_p.alignment = PP_ALIGN.CENTER
#         year_p.font.size = Pt(14)
#         year_p.font.bold = True
#         year_p.font.color.rgb = RGBColor(255, 255, 255)
#         year_tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        
#         # Add goal card
#         card = slide.shapes.add_shape(
#             MSO_SHAPE.ROUNDED_RECTANGLE,
#             start_left, top,
#             card_width, card_height
#         )
#         card.fill.solid()
#         card.fill.fore_color.rgb = RGBColor(255, 255, 255)
#         card.line.color.rgb = teal_color
#         card.line.width = Pt(2)
        
#         # Add goal name
#         tf = card.text_frame
#         tf.text = goal_name
#         tf.margin_left = Inches(0.3)
#         tf.margin_right = Inches(0.3)
#         tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        
#         p = tf.paragraphs[0]
#         p.alignment = PP_ALIGN.LEFT
#         p.font.size = Pt(16)
#         p.font.bold = True
#         p.font.color.rgb = dark_text
    
#     return presentation


# enhanced vertical grid v1
# def create_goals_roadmap(presentation, slide_number, sorted_goals):
#     """
#     Creates a vertical milestone card timeline - modern and minimalist
    
#     Parameters:
#     -----------
#     presentation : pptx.Presentation
#         The presentation object
#     slide_number : int
#         The slide number (0-indexed)
#     sorted_goals : dict
#         Dictionary with goal names as keys and target years as values
#         Example: {'Goal 1': 2028, 'Goal 2': 2032, 'Goal 3': 2034}
#         Goals should be in ascending order of target year
    
#     Returns:
#     --------
#     pptx.Presentation
#         Updated presentation object
#     """
    
#     # Get the slide
#     slide = presentation.slides[slide_number]
    
#     # Theme colors
#     teal_color = RGBColor(0, 176, 185)
#     dark_text = RGBColor(44, 62, 80)
#     white_color = RGBColor(255, 255, 255)
    
#     # Calculate positions
#     goals_list = list(sorted_goals.items())
#     num_goals = len(goals_list)
    
#     if num_goals == 0:
#         return presentation
    
#     # Layout parameters
#     card_width = Inches(7)
#     card_height = Inches(0.9)
#     start_left = Inches(1.5)
#     start_top = Inches(1.8)
#     vertical_spacing = Inches(1.1)
    
#     # Circle parameters - larger for better visibility
#     circle_size = Inches(0.65)
#     circle_left = Inches(0.85)
    
#     # Add connecting line on the left
#     line_left = circle_left + circle_size / 2 - Inches(0.04)
#     line_top = start_top + card_height / 2
#     line_height = (num_goals - 1) * vertical_spacing
    
#     if num_goals > 1:
#         connector = slide.shapes.add_shape(
#             MSO_SHAPE.RECTANGLE,
#             line_left, line_top,
#             Inches(0.08), line_height
#         )
#         connector.fill.solid()
#         connector.fill.fore_color.rgb = teal_color
#         connector.line.fill.background()
    
#     # Create milestone cards
#     for i, (goal_name, target_year) in enumerate(goals_list):
#         top = start_top + i * vertical_spacing
#         circle_top = top + (card_height - circle_size) / 2
        
#         # Add year circle on the left
#         circle = slide.shapes.add_shape(
#             MSO_SHAPE.OVAL,
#             circle_left, circle_top,
#             circle_size, circle_size
#         )
#         circle.fill.solid()
#         circle.fill.fore_color.rgb = teal_color
#         circle.line.color.rgb = white_color
#         circle.line.width = Pt(4)
        
#         # Add year text in circle with better formatting
#         year_tf = circle.text_frame
#         year_tf.clear()
        
#         # Split year into two lines for better readability
#         year_str = str(target_year)
#         if len(year_str) == 4:
#             # Display as "20\n28" or similar
#             year_tf.text = f"{year_str[:2]}\n{year_str[2:]}"
#         else:
#             year_tf.text = year_str
        
#         year_tf.word_wrap = False
#         year_tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        
#         year_p = year_tf.paragraphs[0]
#         year_p.alignment = PP_ALIGN.CENTER
#         year_p.font.name = 'League Spartan'
#         year_p.font.size = Pt(16)
#         year_p.font.bold = True
#         year_p.font.color.rgb = white_color
#         year_p.line_spacing = 0.9
        
#         # Add goal card
#         card = slide.shapes.add_shape(
#             MSO_SHAPE.ROUNDED_RECTANGLE,
#             start_left, top,
#             card_width, card_height
#         )
#         card.fill.solid()
#         card.fill.fore_color.rgb = white_color
#         card.line.color.rgb = teal_color
#         card.line.width = Pt(2)
        
#         # Add goal name
#         tf = card.text_frame
#         tf.text = goal_name
#         tf.margin_left = Inches(0.3)
#         tf.margin_right = Inches(0.3)
#         tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        
#         p = tf.paragraphs[0]
#         p.alignment = PP_ALIGN.LEFT
#         p.font.name = 'League Spartan'
#         p.font.size = Pt(16)
#         p.font.bold = True
#         p.font.color.rgb = dark_text
    
#     return presentation

# enhanced vertical grid v2

def create_goals_roadmap(presentation, slide_number, sorted_goals):
    """
    Creates a vertical milestone card timeline - modern and minimalist
    
    Parameters:
    -----------
    presentation : pptx.Presentation
        The presentation object
    slide_number : int
        The slide number (0-indexed)
    sorted_goals : dict
        Dictionary with goal names as keys and target years as values
        Example: {'Goal 1': 2028, 'Goal 2': 2032, 'Goal 3': 2034}
        Goals should be in ascending order of target year
    
    Returns:
    --------
    pptx.Presentation
        Updated presentation object
    """
    
    # Get the slide
    slide = presentation.slides[slide_number]
    
    # Theme colors
    teal_color = RGBColor(0, 176, 185)
    dark_text = RGBColor(44, 62, 80)
    white_color = RGBColor(255, 255, 255)
    
    # Calculate positions
    goals_list = list(sorted_goals.items())
    num_goals = len(goals_list)
    
    if num_goals == 0:
        return presentation
    
    # Layout parameters
    card_width = Inches(7)
    card_height = Inches(0.9)
    start_left = Inches(2.4)
    start_top = Inches(2.5)
    vertical_spacing = Inches(1.1)
    
    # Circle parameters - larger for better visibility
    circle_size = Inches(1.0)
    circle_left = Inches(1.15)
    
    # Add connecting line on the left
    line_left = circle_left + circle_size / 2 - Inches(0.04)
    line_top = start_top + card_height / 2
    line_height = (num_goals - 1) * vertical_spacing
    
    if num_goals > 1:
        connector = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            line_left, line_top,
            Inches(0.08), line_height
        )
        connector.fill.solid()
        connector.fill.fore_color.rgb = teal_color
        connector.line.fill.background()
    
    # Create milestone cards
    for i, (goal_name, target_year) in enumerate(goals_list):
        top = start_top + i * vertical_spacing
        circle_top = top + (card_height - circle_size) / 2
        
        # Add year circle on the left
        circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            circle_left, circle_top,
            circle_size, circle_size
        )
        circle.fill.solid()
        circle.fill.fore_color.rgb = teal_color
        circle.line.color.rgb = white_color
        circle.line.width = Pt(4)
        
        # Add year text in circle with better formatting
        year_tf = circle.text_frame
        year_tf.clear()
        
    
        year_tf.text = str(target_year)      
        year_tf.word_wrap = False
        year_tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        year_tf.margin_top = Inches(0.05)
        year_tf.margin_bottom = Inches(0.05)
        
        year_p = year_tf.paragraphs[0]
        year_p.alignment = PP_ALIGN.CENTER
        year_p.font.name = 'Calibri'
        year_p.font.size = Pt(18)
        year_p.font.bold = True
        year_p.font.color.rgb = white_color

        
        # Add goal card
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            start_left, top,
            card_width, card_height
        )
        card.fill.solid()
        card.fill.fore_color.rgb = white_color
        card.line.color.rgb = teal_color
        card.line.width = Pt(2)
        
        # Add goal name
        tf = card.text_frame
        tf.text = goal_name
        tf.margin_left = Inches(0.3)
        tf.margin_right = Inches(0.3)
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.font.name = 'Calibri'
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = dark_text
    
    return presentation

######################################### goals visualization #########################################################

################################ move slide #########################################################

def move_slide(presentation, from_index, to_index):
    """
    Move a slide from one position to another
    
    Args:
        presentation: The presentation object
        from_index (int): Current index of the slide to move (0-based)
        to_index (int): Target index where to move the slide (0-based)
    
    Returns:
        bool: True if successful, False otherwise
    """
    slides = presentation.slides
    
    # Validate indices
    if from_index < 0 or from_index >= len(slides):
        raise ValueError(f"Invalid from_index: {from_index}. Must be between 0 and {len(slides)-1}")
    
    if to_index < 0 or to_index >= len(slides):
        raise ValueError(f"Invalid to_index: {to_index}. Must be between 0 and {len(slides)-1}")
    
    if from_index == to_index:
        print("Source and target positions are the same. No move needed.")
        return True
    
    try:
        # Get the XML slide list
        xml_slides = presentation.slides._sldIdLst
        slides_list = list(xml_slides)
        
        # Remove the slide from its current position
        slide_to_move = slides_list[from_index]
        xml_slides.remove(slide_to_move)
        
        # Insert the slide at the new position
        xml_slides.insert(to_index, slide_to_move)
        
        print(f"Successfully moved slide from position {from_index} to position {to_index}")
        return True
        
    except Exception as e:
        print(f"Error moving slide: {e}")
        return False


####################################### move slide #####################################################

###################### slide duplication ##########################################################

# def duplicate_slide(pres, index):
#     source = pres.slides[index]
    
#     # Use the exact same slide layout
#     source_layout = source.slide_layout 
#     new_slide = pres.slides.add_slide(source_layout)

#     # Copy slide-level properties first
#     source_part = source.part
#     new_part = new_slide.part
    
#     # Copy the entire slide element to preserve all properties
#     source_slide_elem = source.element
#     new_slide_elem = new_slide.element
    
#     # Clear new slide content but preserve structure
#     cSld_new = new_slide_elem.find('.//{http://schemas.openxmlformats.org/presentationml/2006/main}cSld')
#     spTree_new = cSld_new.find('.//{http://schemas.openxmlformats.org/presentationml/2006/main}spTree')
    
#     # Remove all shapes from new slide
#     for sp in spTree_new.findall('.//{http://schemas.openxmlformats.org/presentationml/2006/main}sp'):
#         spTree_new.remove(sp)
#     for pic in spTree_new.findall('.//{http://schemas.openxmlformats.org/presentationml/2006/main}pic'):
#         spTree_new.remove(pic)
    
#     # Copy relationships for images
#     rel_mapping = {}
#     for rel_id, rel in source_part.rels.items():
#         if rel.reltype == 'http://schemas.openxmlformats.org/officeDocument/2006/relationships/image':
#             image_part = rel.target_part
#             new_rel_id = new_part.relate_to(image_part, rel.reltype)
#             rel_mapping[rel_id] = new_rel_id

#     # Copy all elements from source slide
#     cSld_source = source_slide_elem.find('.//{http://schemas.openxmlformats.org/presentationml/2006/main}cSld')
    
#     # Copy background if it exists
#     bg_source = cSld_source.find('.//{http://schemas.openxmlformats.org/presentationml/2006/main}bg')
#     if bg_source is not None:
#         bg_new = deepcopy(bg_source)
#         # Update any image references in background
#         for blip in bg_new.findall('.//{http://schemas.openxmlformats.org/drawingml/2006/main}blip'):
#             embed_attr = '{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed'
#             if embed_attr in blip.attrib:
#                 old_id = blip.attrib[embed_attr]
#                 if old_id in rel_mapping:
#                     blip.attrib[embed_attr] = rel_mapping[old_id]
        
#         # Remove existing background and add new one
#         existing_bg = cSld_new.find('.//{http://schemas.openxmlformats.org/presentationml/2006/main}bg')
#         if existing_bg is not None:
#             cSld_new.remove(existing_bg)
#         cSld_new.insert(0, bg_new)
    
#     # Copy all shapes with their exact formatting
#     spTree_source = cSld_source.find('.//{http://schemas.openxmlformats.org/presentationml/2006/main}spTree')
    
#     # Copy each shape element
#     for child in spTree_source:
#         if child.tag.endswith('}sp') or child.tag.endswith('}pic') or child.tag.endswith('}grpSp'):
#             new_child = deepcopy(child)
            
#             # Update image references if this is a picture
#             for blip in new_child.findall('.//{http://schemas.openxmlformats.org/drawingml/2006/main}blip'):
#                 embed_attr = '{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed'
#                 if embed_attr in blip.attrib:
#                     old_id = blip.attrib[embed_attr]
#                     if old_id in rel_mapping:
#                         blip.attrib[embed_attr] = rel_mapping[old_id]
            
#             spTree_new.append(new_child)

#     return new_slide

def duplicate_slide(pres, index):
    source = pres.slides[index]
    
    # Use the exact same slide layout
    source_layout = source.slide_layout 
    new_slide = pres.slides.add_slide(source_layout)

    # Copy slide-level properties first
    source_part = source.part
    new_part = new_slide.part
    
    # Copy the entire slide element to preserve all properties
    source_slide_elem = source.element
    new_slide_elem = new_slide.element
    
    # Clear new slide content but preserve structure
    cSld_new = new_slide_elem.find('.//{http://schemas.openxmlformats.org/presentationml/2006/main}cSld')
    spTree_new = cSld_new.find('.//{http://schemas.openxmlformats.org/presentationml/2006/main}spTree')
    
    # Remove all shapes from new slide (including tables/graphicFrames)
    for sp in spTree_new.findall('.//{http://schemas.openxmlformats.org/presentationml/2006/main}sp'):
        spTree_new.remove(sp)
    for pic in spTree_new.findall('.//{http://schemas.openxmlformats.org/presentationml/2006/main}pic'):
        spTree_new.remove(pic)
    for graphicFrame in spTree_new.findall('.//{http://schemas.openxmlformats.org/presentationml/2006/main}graphicFrame'):
        spTree_new.remove(graphicFrame)
    for grpSp in spTree_new.findall('.//{http://schemas.openxmlformats.org/presentationml/2006/main}grpSp'):
        spTree_new.remove(grpSp)
    
    # Copy relationships for images
    rel_mapping = {}
    for rel_id, rel in source_part.rels.items():
        if rel.reltype == 'http://schemas.openxmlformats.org/officeDocument/2006/relationships/image':
            image_part = rel.target_part
            new_rel_id = new_part.relate_to(image_part, rel.reltype)
            rel_mapping[rel_id] = new_rel_id

    # Copy all elements from source slide
    cSld_source = source_slide_elem.find('.//{http://schemas.openxmlformats.org/presentationml/2006/main}cSld')
    
    # Copy background if it exists
    bg_source = cSld_source.find('.//{http://schemas.openxmlformats.org/presentationml/2006/main}bg')
    if bg_source is not None:
        bg_new = deepcopy(bg_source)
        # Update any image references in background
        for blip in bg_new.findall('.//{http://schemas.openxmlformats.org/drawingml/2006/main}blip'):
            embed_attr = '{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed'
            if embed_attr in blip.attrib:
                old_id = blip.attrib[embed_attr]
                if old_id in rel_mapping:
                    blip.attrib[embed_attr] = rel_mapping[old_id]
        
        # Remove existing background and add new one
        existing_bg = cSld_new.find('.//{http://schemas.openxmlformats.org/presentationml/2006/main}bg')
        if existing_bg is not None:
            cSld_new.remove(existing_bg)
        cSld_new.insert(0, bg_new)
    
    # Copy all shapes with their exact formatting
    spTree_source = cSld_source.find('.//{http://schemas.openxmlformats.org/presentationml/2006/main}spTree')
    
    # Copy each shape element (including tables via graphicFrame)
    for child in spTree_source:
        # Check if it's a shape, picture, group shape, or graphic frame (table/chart)
        if child.tag.endswith(('}sp', '}pic', '}grpSp', '}graphicFrame')):
            new_child = deepcopy(child)
            
            # Update image references if this is a picture
            for blip in new_child.findall('.//{http://schemas.openxmlformats.org/drawingml/2006/main}blip'):
                embed_attr = '{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed'
                if embed_attr in blip.attrib:
                    old_id = blip.attrib[embed_attr]
                    if old_id in rel_mapping:
                        blip.attrib[embed_attr] = rel_mapping[old_id]
            
            spTree_new.append(new_child)

    return new_slide

#################################### slide duplication #######################################################################################

def find_and_update_table(presentation, slide_index, table_data, table_index=0, debug=False):
    """
    Update text inside an existing table WITHOUT touching structure or borders.
    Header row -> 32 pt, bold
    Body rows   -> 24 pt, bold
    """

    # Validate slide index
    if slide_index < 0 or slide_index >= len(presentation.slides):
        print(f"Error: Invalid slide_index {slide_index}")
        return False

    slide = presentation.slides[slide_index]

    # Find tables
    tables = find_all_tables_in_slide(slide, debug=debug)

    if not tables:
        print(f"Error: No tables found in slide {slide_index}")
        return False

    if table_index >= len(tables):
        print(f"Error: table_index {table_index} out of range")
        return False

    table = tables[table_index].table

    if not table_data:
        print("Error: table_data is empty")
        return False

    num_cols = len(table_data)
    num_rows = max(len(v) for v in table_data.values())

    # Do not exceed existing table dimensions
    max_rows = min(num_rows, len(table.rows))
    max_cols = min(num_cols, len(table.columns))

    col_keys = list(table_data.keys())

    for col_idx in range(max_cols):
        values = table_data[col_keys[col_idx]]

        for row_idx in range(max_rows):
            cell = table.cell(row_idx, col_idx)
            tf = cell.text_frame
            value = values[row_idx] if row_idx < len(values) else ""

            # Safe text replacement (preserves borders)
            if tf.paragraphs:
                p = tf.paragraphs[0]
                p.text = str(value)

                # Font rules
                if row_idx == 0:       # Header
                    p.font.size = Pt(32)
                    p.font.bold = True
                else:                  # Body
                    p.font.size = Pt(24)
                    p.font.bold = True

                # Clear extra paragraphs safely
                for extra_p in tf.paragraphs[1:]:
                    extra_p.text = ""
            else:
                tf.text = str(value)
                p = tf.paragraphs[0]

                if row_idx == 0:
                    p.font.size = Pt(32)
                    p.font.bold = True
                else:
                    p.font.size = Pt(24)
                    p.font.bold = True

    print(f"Successfully updated table in slide {slide_index} (borders preserved)")
    return True

def find_all_tables_in_slide(slide, debug=False):
    """
    Find all tables in a slide, including those in groups and placeholders
    
    Args:
        slide: The slide object
        debug (bool): Print debug information
    
    Returns:
        list: List of SHAPE objects that contain tables (not Table objects)
    """
    table_shapes = []
    
    def search_for_tables(shapes, level=0):
        """Recursively search for tables in shapes"""
        for shape in shapes:
            indent = "  " * level
            
            if debug:
                shape_type_name = get_shape_type_name(shape.shape_type)
                print(f"{indent}Checking shape: {shape_type_name} (type={shape.shape_type})")
            
            # Check if this shape is a table using has_table attribute (most reliable)
            if hasattr(shape, 'has_table'):
                try:
                    if shape.has_table:
                        if debug:
                            print(f"{indent}  -> Found table! (Appending SHAPE object)")
                        table_shapes.append(shape)  # Append the SHAPE, not shape.table
                        continue
                except:
                    pass
            
            # Check if shape type is explicitly TABLE (type 19)
            if hasattr(shape, 'shape_type') and shape.shape_type == 19:
                if debug:
                    print(f"{indent}  -> Found table by shape type! (Appending SHAPE object)")
                table_shapes.append(shape)  # Append the SHAPE
                continue
            
            # Check if this shape is a group shape (contains other shapes)
            if hasattr(shape, 'shape_type') and shape.shape_type == 6:  # MSO_SHAPE_TYPE.GROUP
                if debug:
                    print(f"{indent}  -> Group shape, searching inside...")
                if hasattr(shape, 'shapes'):
                    search_for_tables(shape.shapes, level + 1)
    
    search_for_tables(slide.shapes)
    
    if debug:
        print(f"\nTotal table shapes found: {len(table_shapes)}")
        for i, shape in enumerate(table_shapes):
            print(f"  Table {i}: {type(shape).__name__} object")
    
    return table_shapes

def get_shape_type_name(shape_type):
    """Get human-readable name for shape type"""
    shape_types = {
        1: "AUTO_SHAPE",
        2: "CALLOUT",
        3: "CHART",
        4: "COMMENT",
        5: "FREEFORM",
        6: "GROUP",
        7: "EMBEDDED_OLE_OBJECT",
        8: "FORM_CONTROL",
        9: "LINE",
        10: "LINKED_OLE_OBJECT",
        11: "LINKED_PICTURE",
        12: "MEDIA",
        13: "PICTURE",
        14: "PLACEHOLDER",
        15: "SCRIPT_ANCHOR",
        19: "TABLE",  # This is the correct table type
        17: "TEXT_BOX",
        18: "TEXT_EFFECT",
        20: "3D_MODEL",
        21: "GRAPHIC",
        22: "CANVAS",
        23: "DIAGRAM",
        24: "INK",
        25: "INK_COMMENT",
        26: "SMART_ART",
        28: "WEB_VIDEO"
    }
    return shape_types.get(shape_type, f"UNKNOWN({shape_type})")

def extract_table_style(table):
    """
    Extract style information from an existing table
    
    Args:
        table: The table object to extract style from
    
    Returns:
        dict: Dictionary containing style information
    """
    style_info = {
        'font_name': None,
        'font_size': None,
        'font_color': None,
        'bold': False,
        'fill_color': None,
        'border_color': None,
        'alignment': None,
        'vertical_anchor': None
    }
    
    try:
        # Sample the first cell for default styling
        if table.rows and len(table.rows) > 0 and len(table.rows[0].cells) > 0:
            first_cell = table.rows[0].cells[0]
            
            # Get text properties
            if first_cell.text_frame.paragraphs:
                para = first_cell.text_frame.paragraphs[0]
                if para.runs:
                    run = para.runs[0]
                    style_info['font_name'] = run.font.name
                    style_info['font_size'] = run.font.size
                    if run.font.color.type == 1:  # RGB color
                        style_info['font_color'] = run.font.color.rgb
                    style_info['bold'] = run.font.bold
                
                style_info['alignment'] = para.alignment
            
            # Get cell fill
            if first_cell.fill.type == 1:  # Solid fill
                style_info['fill_color'] = first_cell.fill.fore_color.rgb
            
            # Get vertical anchor
            style_info['vertical_anchor'] = first_cell.vertical_anchor
            
    except Exception as e:
        print(f"Warning: Could not extract all style properties: {e}")
    
    return style_info

def apply_cell_style(cell, style_info, row_idx, col_idx, num_rows, num_cols):
    """
    Apply style to a table cell
    
    Args:
        cell: The cell object to style
        style_info (dict): Style information dictionary
        row_idx (int): Row index
        col_idx (int): Column index
        num_rows (int): Total number of rows
        num_cols (int): Total number of columns
    """
    try:
        # Set text properties
        if cell.text_frame.paragraphs:
            para = cell.text_frame.paragraphs[0]
            
            # Set alignment - center for all cells
            para.alignment = PP_ALIGN.CENTER
            
            if para.runs:
                run = para.runs[0]
                
                # Set font to League Spartan
                #run.font.name = 'League Spartan'
                
                # Set font size to 20pt
                run.font.size = Pt(24)
                
                # Set font color based on row
                if row_idx == 0:
                    # First row (header) - white text
                    run.font.color.rgb = RGBColor(255, 255, 255)
                    run.font.bold = True
                else:
                    # Other rows - black text
                    run.font.color.rgb = RGBColor(0, 0, 0)
                    run.font.bold = True
        
        # Set cell fill
        if row_idx == 0:
            # Header row - keep the colored background (teal/cyan color from your image)
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(0, 176, 185)  # Teal/cyan color
        else:
            # All other rows - transparent background
            cell.fill.background()
        
        # Set vertical anchor to middle 
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        
        # Set cell margins for better spacing
        set_cell_margins(cell, left=Pt(5), right=Pt(5), top=Pt(3), bottom=Pt(3))
        
        # Set cell borders
        set_cell_border(cell, 
                       color=RGBColor(0, 0, 0),  # Black borders
                       width=Pt(4))  # 1pt width
        
    except Exception as e:
        print(f"Warning: Could not apply all styles to cell [{row_idx}, {col_idx}]: {e}")

def set_cell_margins(cell, left=None, right=None, top=None, bottom=None):
    """
    Set cell margins (padding)
    
    Args:
        cell: The cell object
        left, right, top, bottom: Margin values in Pt()
    """
    try:
        tc = cell._tc
        tcPr = tc.get_or_add_tcPr()
        
        # Create margin element if it doesn't exist
        for margin in tcPr.findall('.//{http://schemas.openxmlformats.org/drawingml/2006/main}lIns'):
            tcPr.remove(margin)
        for margin in tcPr.findall('.//{http://schemas.openxmlformats.org/drawingml/2006/main}rIns'):
            tcPr.remove(margin)
        for margin in tcPr.findall('.//{http://schemas.openxmlformats.org/drawingml/2006/main}tIns'):
            tcPr.remove(margin)
        for margin in tcPr.findall('.//{http://schemas.openxmlformats.org/drawingml/2006/main}bIns'):
            tcPr.remove(margin)
        
        if left is not None:
            tcPr.set('{http://schemas.openxmlformats.org/drawingml/2006/main}marL', str(int(left)))
        if right is not None:
            tcPr.set('{http://schemas.openxmlformats.org/drawingml/2006/main}marR', str(int(right)))
        if top is not None:
            tcPr.set('{http://schemas.openxmlformats.org/drawingml/2006/main}marT', str(int(top)))
        if bottom is not None:
            tcPr.set('{http://schemas.openxmlformats.org/drawingml/2006/main}marB', str(int(bottom)))
    except:
        pass

def set_cell_border(cell, color=RGBColor(0, 0, 0), width=Pt(1)):
    """
    Set borders for a table cell
    
    Args:
        cell: The cell object
        color: Border color (RGBColor object or tuple of (r, g, b))
        width: Border width (in Pt)
    """
    try:
        from lxml import etree
        
        tc = cell._tc
        tcPr = tc.get_or_add_tcPr()
        
        # Define namespace
        a_ns = 'http://schemas.openxmlformats.org/drawingml/2006/main'
        
        # Remove existing borders
        for border in tcPr.findall(f'.//{{{a_ns}}}lnL'):
            tcPr.remove(border)
        for border in tcPr.findall(f'.//{{{a_ns}}}lnR'):
            tcPr.remove(border)
        for border in tcPr.findall(f'.//{{{a_ns}}}lnT'):
            tcPr.remove(border)
        for border in tcPr.findall(f'.//{{{a_ns}}}lnB'):
            tcPr.remove(border)
        
        # Convert RGBColor to hex string
        # RGBColor has an internal _color attribute that's a 3-tuple (r, g, b)
        # if hasattr(color, '_color'):
        #     r, g, b = color._color
        # elif isinstance(color, tuple) and len(color) == 3:
        #     r, g, b = color
        # else:
        #     # Default to black if we can't parse
        #     r, g, b = 0, 0, 0

        r, g, b = 0, 0, 0
        
        color_hex = f'{r:02X}{g:02X}{b:02X}'
        
        # Create border elements for all sides
        sides = ['lnL', 'lnR', 'lnT', 'lnB']  # Left, Right, Top, Bottom
        
        for side in sides:
            ln = etree.SubElement(tcPr, f'{{{a_ns}}}{side}')
            width_emu = int(width) if width else 12700
            ln.set('w', str(width_emu))
            ln.set('cap', 'flat')
            ln.set('cmpd', 'sng')
            ln.set('algn', 'ctr')
            
            # Add solid fill
            solidFill = etree.SubElement(ln, f'{{{a_ns}}}solidFill')
            srgbClr = etree.SubElement(solidFill, f'{{{a_ns}}}srgbClr')
            srgbClr.set('val', color_hex)
            
    except Exception as e:
        print(f"Warning: Could not set cell borders: {e}")

# def apply_column_width(table, num_cols, total_width):
#     """
#     Apply proportional column widths
    
#     Args:
#         table: The table object
#         num_cols (int): Number of columns
#         total_width: Total width to distribute
#     """
#     try:
#         col_width = total_width / num_cols
#         for col_idx in range(num_cols):
#             table.columns[col_idx].width = int(col_width)
#     except Exception as e:
#         print(f"Warning: Could not set column widths: {e}")

# def apply_column_widths(table, num_cols, total_width):
#     """
#     Apply proportional column widths based on content length
    
#     Args:
#         table: The table object
#         num_cols (int): Number of columns
#         total_width: Total width to distribute
#     """
#     try:
#         # Calculate the maximum content length for each column
#         col_max_lengths = []
        
#         for col_idx in range(num_cols):
#             max_length = 0
#             for row in table.rows:
#                 cell_text = row.cells[col_idx].text
#                 # Use character count as proxy for width needed
#                 max_length = max(max_length, len(cell_text))
#             col_max_lengths.append(max_length)
        
#         # Calculate total length
#         total_length = sum(col_max_lengths)
        
#         # If all columns are empty, distribute equally
#         if total_length == 0:
#             col_width = total_width / num_cols
#             for col_idx in range(num_cols):
#                 table.columns[col_idx].width = int(col_width)
#         else:
#             # Distribute width proportionally based on content
#             for col_idx in range(num_cols):
#                 proportion = col_max_lengths[col_idx] / total_length
#                 table.columns[col_idx].width = int(total_width * proportion)
        
#     except Exception as e:
#         print(f"Warning: Could not set column widths: {e}")

def apply_column_widths(table, num_cols, total_width):
    """
    Apply intelligent column widths based on content with minimum width constraints
    
    Args:
        table: The table object
        num_cols (int): Number of columns
        total_width: Total width to distribute
    """
    try:
        from pptx.util import Inches
        
        # Calculate the maximum content length for each column
        col_max_lengths = []
        
        for col_idx in range(num_cols):
            max_length = 0
            for row in table.rows:
                cell_text = row.cells[col_idx].text
                max_length = max(max_length, len(cell_text))
            col_max_lengths.append(max_length)
        
        # Set minimum width (e.g., 0.8 inches) and maximum width (e.g., 3 inches)
        min_width = Inches(0.8)
        max_width = Inches(5)
        
        # Calculate proportional widths
        total_length = sum(col_max_lengths)
        
        if total_length == 0:
            # Equal distribution if no content
            col_width = total_width / num_cols
            for col_idx in range(num_cols):
                table.columns[col_idx].width = int(col_width)
        else:
            # First pass: calculate ideal widths
            ideal_widths = []
            for col_idx in range(num_cols):
                proportion = col_max_lengths[col_idx] / total_length
                ideal_width = total_width * proportion
                # Apply min/max constraints
                ideal_width = max(min_width, min(max_width, ideal_width))
                ideal_widths.append(ideal_width)
            
            # Second pass: normalize to fit total_width
            current_total = sum(ideal_widths)
            scale_factor = total_width / current_total
            
            for col_idx in range(num_cols):
                final_width = ideal_widths[col_idx] * (scale_factor)
                table.columns[col_idx].width = int(final_width)
        
    except Exception as e:
        print(f"Warning: Could not set column widths: {e}")

def find_and_update_all_tables(presentation, slide_index, tables_data):
    """
    Update multiple tables in a slide
    
    Args:
        presentation: The presentation object
        slide_index (int): Index of the slide
        tables_data (list): List of dictionaries, one for each table
    ,
    Returns:
        bool: True if all successful
    """
    success = True
    for table_idx, table_data in enumerate(tables_data):
        result = find_and_update_table(presentation, slide_index, table_data, table_idx)
        success = success and result
    
    return success

def create_table_if_not_exists(presentation, slide_index, table_data, position=None):
    """
    Create a table if none exists, or update existing table
    
    Args:
        presentation: The presentation object
        slide_index (int): Index of the slide
        table_data (dict): Table data dictionary
        position (dict): Optional position dict with keys: 'left', 'top', 'width', 'height' (in Inches)
    
    Returns:
        bool: True if successful
    """
    slide = presentation.slides[slide_index]
    
    # Check if table exists
    has_table = False
    for shape in slide.shapes:
        if shape.has_table:
            has_table = True
            break
    
    if has_table:
        # Update existing table
        return find_and_update_table(presentation, slide_index, table_data)
    else:
        # Create new table
        if position is None:
            position = {
                'left': Inches(1),
                'top': Inches(2),
                'width': Inches(8),
                'height': Inches(3)
            }
        
        num_cols = len(table_data)
        num_rows = max(len(values) for values in table_data.values())
        
        table_shape = slide.shapes.add_table(
            num_rows, num_cols,
            position['left'], position['top'],
            position['width'], position['height']
        )
        table = table_shape.table
        
        # Populate table
        col_keys = list(table_data.keys())
        for col_idx, col_key in enumerate(col_keys):
            values = table_data[col_key]
            for row_idx in range(num_rows):
                cell = table.cell(row_idx, col_idx)
                cell_value = values[row_idx] if row_idx < len(values) else ""
                cell.text = str(cell_value)
                
                # Apply default styling
                style_info = {
                    'font_name': 'Calibri',
                    'font_size': Pt(11),
                    'font_color': RGBColor(0, 0, 0),
                    'bold': False,
                    'fill_color': None,
                    'border_color': None,
                    'alignment': PP_ALIGN.CENTER,
                    'vertical_anchor': MSO_ANCHOR.MIDDLE
                }
                apply_cell_style(cell, style_info, row_idx, col_idx, num_rows, num_cols)
        
        print(f"Created new table in slide {slide_index}")
        return True

def show_slide_data(ppt, slide_no):
    slide = ppt.slides[slide_no]
    arr = []
    for shape in slide.shapes:
        if shape.has_text_frame and shape.text:
            # print(shape.text)
            arr.append(shape.text)
    return arr
# print(list(x for x in enumerate(ppt.slides)))

def indx_text_boxes(ppt, slide_no):
    for indx, text in enumerate(show_slide_data(ppt, slide_no),1):
        print(f"text box {indx}: {text}")

def update_text_of_textbox(ppt, slide_no, text_box_id, new_text, font_size=None, bold=None, italic=None, font_color=None):
    slide = ppt.slides[(slide_no)]
    count = 0
    for shape in slide.shapes:
        if shape.has_text_frame and shape.text:
            count += 1
            if count == text_box_id:
                text_frame = shape.text_frame
                first_paragraph = text_frame.paragraphs[0]
                first_run = first_paragraph.runs[0] if first_paragraph.runs else first_paragraph.add_run()
                # Preserve formatting of the first run
                font = first_run.font
                orig_font_name = font.name
                orig_font_size = font.size
                orig_font_bold = font.bold
                orig_font_italic = font.italic
                orig_font_underline = font.underline
                orig_font_color = font.color.rgb if hasattr(font.color, 'rgb') and font.color.rgb is not None else None

                # Preserve paragraph alignment
                paragraph_alignment = first_paragraph.alignment

                # Clear only the text content, not the entire text frame
                for paragraph in text_frame.paragraphs:
                    paragraph.clear()

                # Add new text with preserved formatting and alignment
                new_paragraph = text_frame.paragraphs[0]
                new_paragraph.alignment = paragraph_alignment  # Restore alignment
                new_run = new_paragraph.add_run()
                new_run.text = new_text

                # Reapply formatting (use custom values if provided, otherwise use original)
                new_run.font.name = orig_font_name
                new_run.font.size = Pt(font_size) if font_size is not None else orig_font_size
                new_run.font.bold = bold if bold is not None else orig_font_bold
                new_run.font.italic = italic if italic is not None else orig_font_italic
                new_run.font.underline = orig_font_underline
                if font_color is not None:
                    new_run.font.color.rgb = font_color
                elif orig_font_color is not None:
                    new_run.font.color.rgb = orig_font_color
                return

def find_tables_in_slide(ppt, slide_no):
    slide = ppt.slides[slide_no]
    tables = []
    for shape in slide.shapes:
        if shape.has_table:
            tables.append(shape.table)
    return tables

# def replace_pie_chart_with_data(presentation, slide_number, pie_data, title):
#     """
#     Identifies a pie chart in a slide and replaces it with new data while preserving styling.
#     Shows category names outside the pie chart and values inside each slice.
    
#     Parameters:
#     -----------
#     presentation : pptx.Presentation
#         The presentation object
#     slide_number : int
#         The slide number (0-indexed)
#     pie_data : dict
#         Dictionary with category names as keys and values as proportions
#         Example: {'Liquid Assets': 30, 'Fixed Assets': 50, 'Retirement Assets': 20}
#     title : str
#         The title text for the chart
    
#     Returns:
#     --------
#     pptx.Presentation
#         Updated presentation object
#     """
    
#     # Get the slide
#     slide = presentation.slides[slide_number]
    
#     # Find the pie chart in the slide
#     pie_chart_shape = None
#     for shape in slide.shapes:
#         if shape.has_chart:
#             chart = shape.chart
#             # Check if it's a pie chart
#             if chart.chart_type in [XL_CHART_TYPE.PIE, XL_CHART_TYPE.PIE_EXPLODED, 
#                                      XL_CHART_TYPE.THREE_D_PIE, XL_CHART_TYPE.THREE_D_PIE_EXPLODED]:
#                 pie_chart_shape = shape
#                 break
    
#     if pie_chart_shape is None:
#         raise ValueError(f"No pie chart found in slide {slide_number}")
    
#     # Extract styling information from the original chart
#     original_chart = pie_chart_shape.chart
    
#     # Store position and size
#     left = pie_chart_shape.left
#     top = pie_chart_shape.top
#     width = pie_chart_shape.width
#     height = pie_chart_shape.height
    
#     # Store chart type
#     chart_type = original_chart.chart_type
    
#     # Store title information for font styling
#     title_font_size = Pt(18)  # Default
#     title_font_name = None
#     title_font_bold = None
    
#     if original_chart.has_title:
#         title_obj = original_chart.chart_title
#         if title_obj.text_frame.paragraphs[0].font.size:
#             title_font_size = title_obj.text_frame.paragraphs[0].font.size
#         if title_obj.text_frame.paragraphs[0].font.name:
#             title_font_name = title_obj.text_frame.paragraphs[0].font.name
#         title_font_bold = title_obj.text_frame.paragraphs[0].font.bold
    
#     # Store colors from original chart
#     colors = []
#     try:
#         for i, series in enumerate(original_chart.series):
#             for j, point in enumerate(series.points):
#                 try:
#                     fill = point.format.fill
#                     if fill.type == 1:  # SOLID fill type
#                         color = fill.fore_color.rgb
#                         colors.append(color)
#                 except:
#                     colors.append(None)
#     except:
#         pass
    
#     # Delete the original chart
#     sp = pie_chart_shape._element
#     sp.getparent().remove(sp)
    
#     # Create new chart data
#     chart_data = CategoryChartData()
#     chart_data.categories = list(pie_data.keys())
#     chart_data.add_series('Series 1', tuple(pie_data.values()))
    
#     # Add new chart with the same position and size
#     new_chart_shape = slide.shapes.add_chart(
#         chart_type, left, top, width, height, chart_data
#     )
    
#     new_chart = new_chart_shape.chart
    
#     # Apply title
#     new_chart.has_title = True
#     new_title = new_chart.chart_title
#     new_title.text_frame.text = title
#     if title_font_size:
#         new_title.text_frame.paragraphs[0].font.size = title_font_size
#     if title_font_name:
#         new_title.text_frame.paragraphs[0].font.name = title_font_name
#     if title_font_bold is not None:
#         new_title.text_frame.paragraphs[0].font.bold = title_font_bold
    
#     # Remove legend
#     new_chart.has_legend = False
    
#     # Configure data labels - disable default labels first
#     new_plot = new_chart.plots[0]
#     new_plot.has_data_labels = True
#     new_data_labels = new_plot.data_labels
    
#     # Turn off all default label options
#     new_data_labels.show_category_name = False
#     new_data_labels.show_value = False
#     new_data_labels.show_percentage = False
#     new_data_labels.show_legend_key = False
    
#     # Calculate total for percentage calculation
#     total = sum(pie_data.values())
    
#     # Configure individual data points to show category name and percentage inside
#     series = new_chart.series[0]
#     categories = list(pie_data.keys())
#     values = list(pie_data.values())
    
#     for i, point in enumerate(series.points):
#         # Apply colors if available
#         if i < len(colors) and colors[i] is not None:
#             try:
#                 fill = point.format.fill
#                 fill.solid()
#                 fill.fore_color.rgb = colors[i]
#             except:
#                 pass
        
#         # Configure individual point data label to show category name and percentage inside
#         try:
#             point_label = point.data_label
#             point_label.position = XL_DATA_LABEL_POSITION.CENTER
            
#             # Calculate percentage
#             percentage = (values[i] / total) * 100
            
#             # Set text to show category name and percentage
#             point_label.text_frame.text = f"{categories[i]}\n{percentage:.1f}%"
            
#             # Apply League Spartan font, size 20, black color
#             for paragraph in point_label.text_frame.paragraphs:
#                 paragraph.font.name = 'League Spartan'
#                 paragraph.font.size = Pt(20)
#                 paragraph.font.color.rgb = RGBColor(0, 0, 0)
#         except:
#             pass
    
#     return presentation


# from pptx import Presentation
# from pptx.chart.data import CategoryChartData
# from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_DATA_LABEL_POSITION
# from pptx.util import Inches, Pt
# from pptx.dml.color import RGBColor

# def replace_pie_chart_with_data(presentation, slide_number, pie_data, title):
#     """
#     Identifies a pie chart in a slide and replaces it with new data while preserving styling.
#     Uses legends to display category names with League Spartan font.
    
#     Parameters:
#     -----------
#     presentation : pptx.Presentation
#         The presentation object
#     slide_number : int
#         The slide number (0-indexed)
#     pie_data : dict
#         Dictionary with category names as keys and values as proportions
#         Example: {'Liquid Assets': 30, 'Fixed Assets': 50, 'Retirement Assets': 20}
#     title : str
#         The title text for the chart
    
#     Returns:
#     --------
#     pptx.Presentation
#         Updated presentation object
#     """
    
#     # Get the slide
#     slide = presentation.slides[slide_number]
    
#     # Find the pie chart in the slide
#     pie_chart_shape = None
#     for shape in slide.shapes:
#         if shape.has_chart:
#             chart = shape.chart
#             # Check if it's a pie chart
#             if chart.chart_type in [XL_CHART_TYPE.PIE, XL_CHART_TYPE.PIE_EXPLODED, 
#                                      XL_CHART_TYPE.THREE_D_PIE, XL_CHART_TYPE.THREE_D_PIE_EXPLODED]:
#                 pie_chart_shape = shape
#                 break
    
#     if pie_chart_shape is None:
#         raise ValueError(f"No pie chart found in slide {slide_number}")
    
#     # Extract styling information from the original chart
#     original_chart = pie_chart_shape.chart
    
#     # Store position and size
#     left = pie_chart_shape.left
#     top = pie_chart_shape.top
#     width = pie_chart_shape.width
#     height = pie_chart_shape.height
    
#     # Store chart type
#     chart_type = original_chart.chart_type
    
#     # Store title information for font styling
#     title_font_size = Pt(18)  # Default
#     title_font_name = 'League Spartan'
#     title_font_bold = True
    
#     if original_chart.has_title:
#         title_obj = original_chart.chart_title
#         if title_obj.text_frame.paragraphs[0].font.size:
#             title_font_size = title_obj.text_frame.paragraphs[0].font.size
#         if title_obj.text_frame.paragraphs[0].font.bold is not None:
#             title_font_bold = title_obj.text_frame.paragraphs[0].font.bold
    
#     # Store colors from original chart
#     colors = []
#     try:
#         for i, series in enumerate(original_chart.series):
#             for j, point in enumerate(series.points):
#                 try:
#                     fill = point.format.fill
#                     if fill.type == 1:  # SOLID fill type
#                         color = fill.fore_color.rgb
#                         colors.append(color)
#                 except:
#                     colors.append(None)
#     except:
#         pass
    
#     # Delete the original chart
#     sp = pie_chart_shape._element
#     sp.getparent().remove(sp)
    
#     # Create new chart data
#     chart_data = CategoryChartData()
#     chart_data.categories = list(pie_data.keys())
#     chart_data.add_series('Series 1', tuple(pie_data.values()))
    
#     # Add new chart with the same position and size
#     new_chart_shape = slide.shapes.add_chart(
#         chart_type, left, top, width, height, chart_data
#     )
    
#     new_chart = new_chart_shape.chart
    
#     # Apply title with League Spartan font
#     new_chart.has_title = True
#     new_title = new_chart.chart_title
#     new_title.text_frame.text = title
#     new_title.text_frame.paragraphs[0].font.size = title_font_size
#     new_title.text_frame.paragraphs[0].font.name = title_font_name
#     new_title.text_frame.paragraphs[0].font.bold = title_font_bold
#     new_title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)
    
#     # Enable and configure legend with League Spartan font
#     new_chart.has_legend = True
#     new_legend = new_chart.legend
    
#     # Position legend on the right side
#     new_legend.position = XL_LEGEND_POSITION.RIGHT
    
#     # Configure legend font - League Spartan, size 16, black, bold
#     new_legend.font.name = 'League Spartan'
#     new_legend.font.size = Pt(16)
#     new_legend.font.bold = True
#     new_legend.font.color.rgb = RGBColor(0, 0, 0)
    
#     # Show only percentage values inside the pie slices
#     new_plot = new_chart.plots[0]
#     new_plot.has_data_labels = True
#     new_data_labels = new_plot.data_labels
    
#     # Configure to show only percentages
#     new_data_labels.show_category_name = False
#     new_data_labels.show_value = False
#     new_data_labels.show_percentage = True
#     new_data_labels.show_legend_key = False
#     new_data_labels.position = XL_DATA_LABEL_POSITION.CENTER
    
#     # Apply League Spartan font to percentage labels
#     new_data_labels.font.name = 'League Spartan'
#     new_data_labels.font.size = Pt(18)
#     new_data_labels.font.bold = True
#     new_data_labels.font.color.rgb = RGBColor(0, 0, 0)
    
#     # Apply colors to the new chart
#     series = new_chart.series[0]
#     for i, point in enumerate(series.points):
#         if i < len(colors) and colors[i] is not None:
#             try:
#                 fill = point.format.fill
#                 fill.solid()
#                 fill.fore_color.rgb = colors[i]
#             except:
#                 pass
    
#     return presentation

# from pptx import Presentation
# from pptx.chart.data import CategoryChartData
# from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_DATA_LABEL_POSITION
# from pptx.util import Inches, Pt
# from pptx.dml.color import RGBColor

# def replace_pie_chart_with_data(presentation, slide_number, pie_data, title):
#     """
#     Identifies a pie chart in a slide and replaces it with new data while preserving styling.
#     Uses legends to display category names with League Spartan font.
    
#     Parameters:
#     -----------
#     presentation : pptx.Presentation
#         The presentation object
#     slide_number : int
#         The slide number (0-indexed)
#     pie_data : dict
#         Dictionary with category names as keys and values as proportions
#         Example: {'Liquid Assets': 30, 'Fixed Assets': 50, 'Retirement Assets': 20}
#     title : str
#         The title text for the chart
    
#     Returns:
#     --------
#     pptx.Presentation
#         Updated presentation object
#     """
    
#     # Get the slide
#     slide = presentation.slides[slide_number]
    
#     # Find the pie chart in the slide
#     pie_chart_shape = None
#     for shape in slide.shapes:
#         if shape.has_chart:
#             chart = shape.chart
#             # Check if it's a pie chart
#             if chart.chart_type in [XL_CHART_TYPE.PIE, XL_CHART_TYPE.PIE_EXPLODED, 
#                                      XL_CHART_TYPE.THREE_D_PIE, XL_CHART_TYPE.THREE_D_PIE_EXPLODED]:
#                 pie_chart_shape = shape
#                 break
    
#     if pie_chart_shape is None:
#         raise ValueError(f"No pie chart found in slide {slide_number}")
    
#     # Extract styling information from the original chart
#     original_chart = pie_chart_shape.chart
    
#     # Store position and size
#     left = pie_chart_shape.left
#     top = pie_chart_shape.top
#     width = pie_chart_shape.width
#     height = pie_chart_shape.height
    
#     # Reduce chart width significantly to make room for legend
#     # Move it slightly left and reduce to 60% of original width
#     adjusted_width = int(width * 0.60)
#     adjusted_left = left
    
#     # Store chart type
#     chart_type = original_chart.chart_type
    
#     # Store title information for font styling
#     title_font_size = Pt(18)  # Default
#     title_font_name = 'League Spartan'
#     title_font_bold = True
    
#     if original_chart.has_title:
#         title_obj = original_chart.chart_title
#         if title_obj.text_frame.paragraphs[0].font.size:
#             title_font_size = title_obj.text_frame.paragraphs[0].font.size
#         if title_obj.text_frame.paragraphs[0].font.bold is not None:
#             title_font_bold = title_obj.text_frame.paragraphs[0].font.bold
    
#     # Store colors from original chart
#     colors = []
#     try:
#         for i, series in enumerate(original_chart.series):
#             for j, point in enumerate(series.points):
#                 try:
#                     fill = point.format.fill
#                     if fill.type == 1:  # SOLID fill type
#                         color = fill.fore_color.rgb
#                         colors.append(color)
#                 except:
#                     colors.append(None)
#     except:
#         pass
    
#     # Delete the original chart
#     sp = pie_chart_shape._element
#     sp.getparent().remove(sp)
    
#     # Create new chart data
#     chart_data = CategoryChartData()
#     chart_data.categories = list(pie_data.keys())
#     chart_data.add_series('Series 1', tuple(pie_data.values()))
    
#     # Add new chart with adjusted position and size to accommodate legend
#     new_chart_shape = slide.shapes.add_chart(
#         chart_type, adjusted_left, top, adjusted_width, height, chart_data
#     )
    
#     new_chart = new_chart_shape.chart
    
#     # Apply title with League Spartan font
#     new_chart.has_title = True
#     new_title = new_chart.chart_title
#     new_title.text_frame.text = title
#     new_title.text_frame.paragraphs[0].font.size = title_font_size
#     new_title.text_frame.paragraphs[0].font.name = title_font_name
#     new_title.text_frame.paragraphs[0].font.bold = title_font_bold
#     new_title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)
    
#     # Enable and configure legend with League Spartan font
#     new_chart.has_legend = True
#     new_legend = new_chart.legend
    
#     # Position legend on the right side
#     new_legend.position = XL_LEGEND_POSITION.RIGHT
    
#     # Set legend to not overlap with plot area
#     new_legend.include_in_layout = True
    
#     # Configure legend font - League Spartan, size 14, black, bold
#     new_legend.font.name = 'League Spartan'
#     new_legend.font.size = Pt(14)
#     new_legend.font.bold = True
#     new_legend.font.color.rgb = RGBColor(0, 0, 0)
    
#     # Show only percentage values inside the pie slices with larger font
#     new_plot = new_chart.plots[0]
#     new_plot.has_data_labels = True
#     new_data_labels = new_plot.data_labels
    
#     # Configure to show only percentages
#     new_data_labels.show_category_name = False
#     new_data_labels.show_value = False
#     new_data_labels.show_percentage = True
#     new_data_labels.show_legend_key = False
#     new_data_labels.position = XL_DATA_LABEL_POSITION.CENTER
    
#     # Apply League Spartan font to percentage labels with larger size
#     new_data_labels.font.name = 'League Spartan'
#     new_data_labels.font.size = Pt(22)
#     new_data_labels.font.bold = True
#     new_data_labels.font.color.rgb = RGBColor(0, 0, 0)
    
#     # Apply colors to the new chart
#     series = new_chart.series[0]
#     for i, point in enumerate(series.points):
#         if i < len(colors) and colors[i] is not None:
#             try:
#                 fill = point.format.fill
#                 fill.solid()
#                 fill.fore_color.rgb = colors[i]
#             except:
#                 pass
    
#     return presentation

# from pptx import Presentation
# from pptx.enum.chart import XL_CHART_TYPE
# import matplotlib.pyplot as plt
# import matplotlib.font_manager as fm
# import io
# from PIL import Image

# def replace_pie_chart_with_matplotlib(presentation, slide_number, pie_data, title):
#     """
#     Identifies a pie chart in a slide, captures its position, creates a matplotlib pie chart,
#     and replaces the original chart with the matplotlib image.
    
#     Parameters:
#     -----------
#     presentation : pptx.Presentation
#         The presentation object
#     slide_number : int
#         The slide number (0-indexed)
#     pie_data : dict
#         Dictionary with category names as keys and values as proportions
#         Example: {'Real Estate Investment': 69, 'Direct Equity': 6, 'PMS/AIF': 5, 'PPF': 5, 'others': 15}
#     title : str
#         The title text for the chart
    
#     Returns:
#     --------
#     pptx.Presentation
#         Updated presentation object
#     """
    
#     # Get the slide
#     slide = presentation.slides[slide_number]
    
#     # Find the pie chart in the slide
#     pie_chart_shape = None
#     for shape in slide.shapes:
#         if shape.has_chart:
#             chart = shape.chart
#             # Check if it's a pie chart
#             if chart.chart_type in [XL_CHART_TYPE.PIE, XL_CHART_TYPE.PIE_EXPLODED, 
#                                      XL_CHART_TYPE.THREE_D_PIE, XL_CHART_TYPE.THREE_D_PIE_EXPLODED]:
#                 pie_chart_shape = shape
#                 break
    
#     if pie_chart_shape is None:
#         raise ValueError(f"No pie chart found in slide {slide_number}")
    
#     # Capture position and size of the original chart
#     left = pie_chart_shape.left
#     top = pie_chart_shape.top
#     width = pie_chart_shape.width
#     height = pie_chart_shape.height
    
#     # Store colors from original chart if available
#     colors = []
#     try:
#         original_chart = pie_chart_shape.chart
#         for i, series in enumerate(original_chart.series):
#             for j, point in enumerate(series.points):
#                 try:
#                     fill = point.format.fill
#                     if fill.type == 1:  # SOLID fill type
#                         color = fill.fore_color.rgb
#                         # Convert to hex format for matplotlib
#                         hex_color = '#{:02x}{:02x}{:02x}'.format(color[0], color[1], color[2])
#                         colors.append(hex_color)
#                 except:
#                     colors.append(None)
#     except:
#         pass
    
#     # If no colors extracted, use default matplotlib colors
#     if not colors or all(c is None for c in colors):
#         colors = None
    
#     # Delete the original chart
#     sp = pie_chart_shape._element
#     sp.getparent().remove(sp)
    
#     # Create matplotlib pie chart with proper aspect ratio
#     fig, ax = plt.subplots(figsize=(14, 12), facecolor='none')
#     ax.set_facecolor('none')
    
#     # Prepare data
#     labels = list(pie_data.keys())
#     sizes = list(pie_data.values())
    
#     # Calculate percentages
#     total = sum(sizes)
#     percentages = [(size/total)*100 for size in sizes]
    
#     # Create the pie chart centered
#     wedges, texts, autotexts = ax.pie(
#         sizes,
#         labels=None,  # Don't show labels on pie
#         autopct='%1.0f%%',
#         startangle=90,
#         colors=colors,
#         textprops={'fontsize': 20, 'fontweight': 'bold', 'color': 'black'},
#         pctdistance=0.65,
#         wedgeprops={'linewidth': 3, 'edgecolor': 'white'},
#         radius=1.0
#     )
    
#     # Try to set League Spartan font for percentage labels
#     try:
#         for autotext in autotexts:
#             autotext.set_fontfamily('League Spartan')
#             autotext.set_fontsize(24)
#             autotext.set_fontweight('bold')
#             autotext.set_color('black')
#     except:
#         # If League Spartan not available, use default sans-serif
#         for autotext in autotexts:
#             autotext.set_fontfamily('sans-serif')
#             autotext.set_fontsize(24)
#             autotext.set_fontweight('bold')
#             autotext.set_color('black')
    
#     # Add legend at the bottom with much larger font
#     try:
#         legend = ax.legend(
#             labels,
#             loc='upper center',
#             bbox_to_anchor=(0.5, -0.08),
#             fontsize=20,
#             frameon=False,
#             ncol=3,  # Display in 3 columns for better horizontal layout
#             columnspacing=2.0,
#             handlelength=2.5,
#             handleheight=1.5,
#             prop={'family': 'League Spartan', 'weight': 'bold', 'size': 20}
#         )
#     except:
#         legend = ax.legend(
#             labels,
#             loc='upper center',
#             bbox_to_anchor=(0.5, -0.08),
#             fontsize=20,
#             frameon=False,
#             ncol=3,  # Display in 3 columns for better horizontal layout
#             columnspacing=2.0,
#             handlelength=2.5,
#             handleheight=1.5,
#             prop={'family': 'sans-serif', 'weight': 'bold', 'size': 20}
#         )
    
#     # Set legend text color to black
#     for text in legend.get_texts():
#         text.set_color('black')
    
#     # Equal aspect ratio ensures that pie is drawn as a circle
#     ax.axis('equal')
    
#     # Add title
#     try:
#         plt.title(title, fontsize=20, fontweight='bold', 
#                  fontfamily='League Spartan', color='black', pad=20)
#     except:
#         plt.title(title, fontsize=20, fontweight='bold', 
#                  fontfamily='sans-serif', color='black', pad=20)
    
#     # Save the figure to a bytes buffer with transparent background
#     buf = io.BytesIO()
#     plt.tight_layout()
#     plt.savefig(buf, format='png', dpi=150, bbox_inches='tight', 
#                 facecolor='none', edgecolor='none', transparent=True)
#     buf.seek(0)
#     plt.close()
    
#     # Load image from buffer
#     img = Image.open(buf)
    
#     # Save to temporary buffer for PowerPoint
#     img_buf = io.BytesIO()
#     img.save(img_buf, format='PNG')
#     img_buf.seek(0)
    
#     # Add picture to slide at the original chart's position
#     slide.shapes.add_picture(img_buf, left, top, width, height)
    
#     buf.close()
#     img_buf.close()
    
#     return presentation

from pptx import Presentation
from pptx.enum.chart import XL_CHART_TYPE
import matplotlib.pyplot as plt
import matplotlib.font_manager as fm
import io
from PIL import Image

def replace_pie_chart_with_matplotlib(presentation, slide_number, pie_data, title):
    """
    Identifies a pie chart in a slide, captures its position, creates a matplotlib pie chart,
    and replaces the original chart with the matplotlib image.
    
    Parameters:
    -----------
    presentation : pptx.Presentation
        The presentation object
    slide_number : int
        The slide number (0-indexed)
    pie_data : dict
        Dictionary with category names as keys and values as proportions
        Example: {'Real Estate Investment': 0.69, 'Direct Equity': 0.06, 'PMS/AIF': 0.05, 'PPF': 0.05, 'others': 0.15}
    title : str
        The title text for the chart
    
    Returns:
    --------
    pptx.Presentation
        Updated presentation object
    """
    
    # Get the slide
    slide = presentation.slides[slide_number]
    
    # Find the pie chart in the slide
    pie_chart_shape = None
    for shape in slide.shapes:
        if shape.has_chart:
            chart = shape.chart
            # Check if it's a pie chart
            if chart.chart_type in [XL_CHART_TYPE.PIE, XL_CHART_TYPE.PIE_EXPLODED, 
                                     XL_CHART_TYPE.THREE_D_PIE, XL_CHART_TYPE.THREE_D_PIE_EXPLODED]:
                pie_chart_shape = shape
                break
    
    if pie_chart_shape is None:
        raise ValueError(f"No pie chart found in slide {slide_number}")
    
    # Capture position and size of the original chart
    left = pie_chart_shape.left
    top = pie_chart_shape.top
    width = pie_chart_shape.width
    height = pie_chart_shape.height
    
    # Delete the original chart
    sp = pie_chart_shape._element
    sp.getparent().remove(sp)
    
    # Color palette matching the image - Cyan gradient to dark blue
    modern_colors = [
        '#00FFFF',  # Cyan (brightest)
        '#00D9E1',  # Bright Cyan
        '#00ADB5',  # Muted Teal
        '#0088A8',  # Medium Cyan-Blue
        '#0077B6',  # Star Blue
        '#0080FF',  # Azure
        '#005A8D',  # Medium Blue
        '#023E8A',  # Dark Blue
        '#03045E',  # Deep Night (darkest)
        '#1E1548',  # Purple-Dark (anchor)
    ]
    
    # Use chart intent instead of hardcoded slide index.
    # In current flow, "Current Asset Allocation" can be rendered on slide 3.
    is_current_asset_chart = title.strip().lower() == "current asset allocation"
    
    # Create matplotlib pie chart with proper aspect ratio
    fig, ax = plt.subplots(figsize=(16, 14), facecolor='none') 
    ax.set_facecolor('none')
    
    # Filter out zero or very small values
    filtered_data = {k: v for k, v in pie_data.items() if v > 0}
    
    # Prepare data
    labels = list(filtered_data.keys())
    sizes = list(filtered_data.values())
    
    # Calculate ratios
    total = sum(sizes)
    ratios = [(size/total) for size in sizes]
    
    # Create the donut chart with modern styling
    wedges, texts = ax.pie(
        sizes,
        labels=None,
        autopct=None,
        startangle=90,
        colors=modern_colors[:len(sizes)],
        wedgeprops={
            'linewidth': 2,
            'edgecolor': 'white',
            'width': 0.4,  # Larger hole
            'antialiased': True
        },
        radius=1.0
    )
    
    # Add center circle for donut effect
    centre_circle = plt.Circle((0, 0), 0.60, fc='white', linewidth=0)
    ax.add_artist(centre_circle)
    
    # NO TEXT IN CENTER
    
    # First, draw all connecting lines (so they appear behind the labels)
    for i, (wedge, ratio) in enumerate(zip(wedges, ratios)):
        # Get the angle of the wedge center
        angle = (wedge.theta2 - wedge.theta1) / 2. + wedge.theta1
        
        # Calculate position at the edge of the donut (outer edge)
        x_edge = 1.0 * wedge.r * np.cos(np.radians(angle))
        y_edge = 1.0 * wedge.r * np.sin(np.radians(angle))
        
        # Calculate position for the label (further outside for more spacing)
        if is_current_asset_chart:
            label_distance = 1.35  # Closer for larger labels on slide 4
        else:
            label_distance = 1.5
        
        x_label = label_distance * wedge.r * np.cos(np.radians(angle))
        y_label = label_distance * wedge.r * np.sin(np.radians(angle))
        
        # Draw connecting line from edge to label
        ax.plot([x_edge, x_label], [y_edge, y_label],
               color='black', linewidth=4, alpha=1.0, zorder=1)
    
    # Then, add labels on top of the lines
    for i, (wedge, ratio, label) in enumerate(zip(wedges, ratios, labels)):
        # Get the angle of the wedge center
        angle = (wedge.theta2 - wedge.theta1) / 2. + wedge.theta1
        
        # Calculate position for the label
        if is_current_asset_chart:
            label_distance = 1.35  # Closer for larger labels
        else:
            label_distance = 1.5
        
        x_label = label_distance * wedge.r * np.cos(np.radians(angle))
        y_label = label_distance * wedge.r * np.sin(np.radians(angle))
        
        # Determine horizontal alignment based on angle
        if x_label > 0:
            ha = 'left'
        else:
            ha = 'right'
        
        # Conditional label formatting based on slide
        if is_current_asset_chart:
            # Current asset chart: show only percentage in figure
            label_text = f'{ratio:.2%}'
            fontsize = 20  # Larger for slide 4
            pad = 0.6
            linewidth = 3.0
        else:
            # OTHER SLIDES: Show ratio as percentage, not decimal fraction
            label_text = f'{ratio:.2%}'
            fontsize = 28
            pad = 0.8
            linewidth = 3.5
        
        # Add the label
        try:
            ax.text(x_label, y_label, label_text,
                   ha=ha, va='center',
                   fontsize=fontsize, fontweight='bold', color='black',
                   fontfamily='Calibri',
                   bbox=dict(boxstyle='round,pad=' + str(pad), facecolor='white',
                            edgecolor='black', linewidth=linewidth, alpha=1.0),
                   zorder=2)
        except:
            ax.text(x_label, y_label, label_text,
                   ha=ha, va='center',
                   fontsize=fontsize, fontweight='bold', color='black',
                   fontfamily='sans-serif',
                   bbox=dict(boxstyle='round,pad=' + str(pad), facecolor='white',
                            edgecolor='black', linewidth=linewidth, alpha=1.0),
                   zorder=2)
    
    # Legend behavior:
    # - Current asset chart: color + asset names only
    # - Other charts: existing legend behavior
    if is_current_asset_chart:
        try:
            legend = ax.legend(
                wedges,
                labels,
                loc='upper center',
                bbox_to_anchor=(0.5, -0.05),
                fontsize=18,
                frameon=False,
                ncol=2,
                columnspacing=2.0,
                handlelength=1.8,
                handleheight=1.2,
                labelspacing=0.7,
                prop={'family': 'Calibri', 'weight': 'bold', 'size': 18}
            )
            for text in legend.get_texts():
                text.set_color('#2C3E50')
        except:
            pass
    else:
        try:
            legend = ax.legend(
                labels,
                loc='upper center',
                bbox_to_anchor=(0.5, -0.05),
                fontsize=24,
                frameon=False,
                ncol=2,
                columnspacing=3.0,
                handlelength=2.5,
                handleheight=1.5,
                labelspacing=0.8,
                prop={'family': 'Calibri', 'weight': 'bold', 'size': 24}
            )
            # Set legend text color
            for text in legend.get_texts():
                text.set_color('#2C3E50')
        except:
            pass
    
    # Equal aspect ratio ensures that pie is drawn as a circle
    ax.axis('equal')
    
    # Save the figure to a bytes buffer with transparent background
    buf = io.BytesIO()
    plt.tight_layout()
    plt.savefig(buf, format='png', dpi=150, bbox_inches='tight', 
                facecolor='none', edgecolor='none', transparent=True)
    buf.seek(0)
    plt.close()
    
    # Load image from buffer
    img = Image.open(buf)
    
    # Save to temporary buffer for PowerPoint
    img_buf = io.BytesIO()
    img.save(img_buf, format='PNG')
    img_buf.seek(0)
    
    # Add picture to slide at the original chart's position
    slide.shapes.add_picture(img_buf, left, top, width, height)
    
    buf.close()
    img_buf.close()
    
    return presentation


############################## delete a table from slide ##################################################

def delete_tables_from_slide(prs, slide_number):
    """
    Detects and deletes all tables from a specific slide.
    
    Parameters:
    -----------
    prs : Presentation object
        The presentation object containing the slides
    slide_number : int
        The slide number (0-indexed) from which to delete tables
    
    Returns:
    --------
    Presentation object
        The modified presentation object
    """
    # Get the slide
    slide = prs.slides[slide_number]
    
    # Get the slide element
    slide_elem = slide.element
    
    # Find the spTree (shape tree) element
    cSld = slide_elem.find('.//{http://schemas.openxmlformats.org/presentationml/2006/main}cSld')
    spTree = cSld.find('.//{http://schemas.openxmlformats.org/presentationml/2006/main}spTree')
    
    # Find all graphicFrame elements (tables are stored as graphicFrames)
    graphic_frames = spTree.findall('.//{http://schemas.openxmlformats.org/presentationml/2006/main}graphicFrame')
    
    tables_deleted = 0
    
    # Check each graphicFrame to see if it's a table
    for gf in graphic_frames:
        # Tables have a specific URI in their graphic element
        graphic = gf.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}graphic')
        if graphic is not None:
            graphicData = graphic.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}graphicData')
            if graphicData is not None:
                # Check if the URI indicates this is a table
                uri = graphicData.get('uri')
                if uri == 'http://schemas.openxmlformats.org/drawingml/2006/table':
                    # This is a table, remove it from spTree
                    spTree.remove(gf)
                    tables_deleted += 1
    
    print(f"Deleted {tables_deleted} table(s) from slide {slide_number}")
    
    return prs

############################### delete a table from slide ##################################################

################################## delete a table from a slide and add a text #########################################

def replace_tables_with_text(prs, slide_number, replacement_text):
    """
    Detects and deletes all tables from a specific slide and replaces them with a textbox.
    
    Parameters:
    -----------
    prs : Presentation object
        The presentation object containing the slides
    slide_number : int
        The slide number (0-indexed) from which to delete tables
    replacement_text : str
        The text to display in place of the deleted tables
    
    Returns:
    --------
    Presentation object
        The modified presentation object
    """
    # Get the slide
    slide = prs.slides[slide_number]
    
    # Get the slide element
    slide_elem = slide.element
    
    # Find the spTree (shape tree) element
    cSld = slide_elem.find('.//{http://schemas.openxmlformats.org/presentationml/2006/main}cSld')
    spTree = cSld.find('.//{http://schemas.openxmlformats.org/presentationml/2006/main}spTree')
    
    # Find all graphicFrame elements (tables are stored as graphicFrames)
    graphic_frames = spTree.findall('.//{http://schemas.openxmlformats.org/presentationml/2006/main}graphicFrame')
    
    tables_deleted = 0
    table_positions = []
    
    # Check each graphicFrame to see if it's a table
    for gf in graphic_frames:
        # Tables have a specific URI in their graphic element
        graphic = gf.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}graphic')
        if graphic is not None:
            graphicData = graphic.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}graphicData')
            if graphicData is not None:
                # Check if the URI indicates this is a table
                uri = graphicData.get('uri')
                if uri == 'http://schemas.openxmlformats.org/drawingml/2006/table':
                    # Get the position and size of the table before deleting
                    xfrm = gf.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}xfrm')
                    if xfrm is not None:
                        off = xfrm.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}off')
                        ext = xfrm.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}ext')
                        
                        if off is not None and ext is not None:
                            # Get position and size in EMUs (English Metric Units)
                            x = int(off.get('x'))
                            y = int(off.get('y'))
                            cx = int(ext.get('cx'))
                            cy = int(ext.get('cy'))
                            
                            # Store the position (convert EMUs to Inches)
                            table_positions.append({
                                'left': x / 914400,  # EMU to inches
                                'top': y / 914400,
                                'width': cx / 914400,
                                'height': cy / 914400
                            })
                    
                    # This is a table, remove it from spTree
                    spTree.remove(gf)
                    tables_deleted += 1
    
    # Add textbox(es) in place of deleted table(s)
    if tables_deleted > 0:
        # If we have position info, use it; otherwise use a default position
        if table_positions:
            for pos in table_positions:
                # Add textbox at the table's position
                textbox = slide.shapes.add_textbox(
                    Inches(pos['left']),
                    Inches(pos['top']),
                    Inches(pos['width']),
                    Inches(pos['height'])
                )
                
                # Add text to the textbox
                text_frame = textbox.text_frame
                text_frame.text = replacement_text
                
                # Format the text
                text_frame.word_wrap = True
                paragraph = text_frame.paragraphs[0]
                paragraph.alignment = PP_ALIGN.LEFT
                paragraph.font.size = Pt(14)
                paragraph.font.name = 'Arial'
        else:
            # Default position if we couldn't get table position
            left = Inches(1)
            top = Inches(3)
            width = Inches(8)
            height = Inches(2)
            
            textbox = slide.shapes.add_textbox(left, top, width, height)
            text_frame = textbox.text_frame
            text_frame.text = replacement_text
            text_frame.word_wrap = True
            paragraph = text_frame.paragraphs[0]
            paragraph.alignment = PP_ALIGN.LEFT
            paragraph.font.size = Pt(14)
            paragraph.font.name = 'Arial'
    
    print(f"Deleted {tables_deleted} table(s) from slide {slide_number} and added replacement text")
    
    return prs

#################################### delete a table from a slide and add a text ############################################

def create_risk_band_visual(presentation, slide_number, risk_level, position=None):
    """
    Creates a visual risk band/meter showing the risk assessment level.
    
    Parameters:
    -----------
    presentation : pptx.Presentation
        The presentation object
    slide_number : int
        The slide number (0-indexed)
    risk_level : str
        The risk level - one of: 'Low', 'Moderate', 'High', 'Very High', 'Conservative', 'Aggressive'
    position : dict, optional
        Dictionary with 'left', 'top', 'width', 'height' in Inches
        If None, uses default position
    
    Returns:
    --------
    pptx.Presentation
        Updated presentation object
    """
    
    # Get the slide
    slide = presentation.slides[slide_number]
    
    # Normalize risk level to standard categories
    risk_mapping = {
        'low': 1,
        'conservative': 1,
        'low to moderate': 2,
        'moderate': 2,
        'moderately conservative': 2,
        'moderate to high': 3,
        'moderately aggressive': 3,
        'high': 4,
        'aggressive': 4,
        'very high': 5,
        'very aggressive': 5
    }
    
    risk_level_normalized = risk_level.lower().strip()
    risk_value = risk_mapping.get(risk_level_normalized, 2)  # Default to moderate
    
    # Define risk categories and colors
    risk_categories = ['Low', 'Moderate', 'Moderately\nHigh', 'High', 'Very High']
    risk_colors = ['#2E7D32', '#558B2F', '#FF8F00', '#EF6C00', '#C62828']  # Green to Red gradient
    
    # Create figure with transparent background
    fig, ax = plt.subplots(figsize=(12, 3), facecolor='none')
    ax.set_facecolor('none')
    
    # Remove axes
    ax.set_xlim(0, 5)
    ax.set_ylim(0, 1)
    ax.axis('off')
    
    # Draw the risk band segments
    segment_width = 1.0
    for i in range(5):
        # Create rectangle for each risk level
        # Selected segment is fully opaque (1.0), others are more visible (0.65)
        rect = plt.Rectangle((i * segment_width, 0.3), segment_width, 0.4,
                             facecolor=risk_colors[i],
                             edgecolor='white',
                             linewidth=4,
                             alpha=1.0 if i == risk_value - 1 else 0.80,
                             zorder=1)
        ax.add_patch(rect)
        
        # Add label below each segment
        try:
            ax.text(i * segment_width + segment_width/2, 0.1,
                   risk_categories[i],
                   ha='center', va='center',
                   fontsize=24,
                   fontweight='bold' if i == risk_value - 1 else 'normal',
                   color='#2C3E50',
                   fontfamily='Calibri')
        except:
            ax.text(i * segment_width + segment_width/2, 0.1,
                   risk_categories[i],
                   ha='center', va='center',
                   fontsize=24,
                   fontweight='bold' if i == risk_value - 1 else 'normal',
                   color='#2C3E50',
                   fontfamily='sans-serif')
    
    # Add pointer/arrow to indicate current risk level
    pointer_x = (risk_value - 1) * segment_width + segment_width/2
    
    # Draw arrow pointing to the selected risk level
    arrow = plt.Polygon([[pointer_x - 0.15, 0.85],
                         [pointer_x + 0.15, 0.85],
                         [pointer_x, 0.75]],
                        facecolor=risk_colors[risk_value - 1],
                        edgecolor='white',
                        linewidth=3,
                        zorder=2)
    ax.add_patch(arrow)
    
    # Add "Your Risk Profile" label above the arrow
    try:
        ax.text(pointer_x, 0.92, 'Your Risk Profile',
               ha='center', va='bottom',
               fontsize=18,
               fontweight='bold',
               color=risk_colors[risk_value - 1],
               fontfamily='Calibri',
               bbox=dict(boxstyle='round,pad=0.6', facecolor='white', 
                        edgecolor=risk_colors[risk_value - 1], linewidth=3, alpha=1.0))
    except:
        ax.text(pointer_x, 0.92, 'Your Risk Profile',
               ha='center', va='bottom',
               fontsize=18,
               fontweight='bold',
               color=risk_colors[risk_value - 1],
               fontfamily='sans-serif',
               bbox=dict(boxstyle='round,pad=0.6', facecolor='white', 
                        edgecolor=risk_colors[risk_value - 1], linewidth=3, alpha=1.0))
    
    # Save the figure to a bytes buffer with TRANSPARENT background
    buf = io.BytesIO()
    plt.tight_layout()
    plt.savefig(buf, format='png', dpi=150, bbox_inches='tight',
                facecolor='none', edgecolor='none', transparent=True)  # Transparent
    buf.seek(0)
    plt.close()
    
    # Load image from buffer
    img = Image.open(buf)
    
    # Save to temporary buffer for PowerPoint
    img_buf = io.BytesIO()
    img.save(img_buf, format='PNG')
    img_buf.seek(0)
    
    # Set position - use provided position or default
    if position is None:
        left = Inches(1.5)
        top = Inches(3.5)
        width = Inches(7)
        height = Inches(1.5)
    else:
        left = position.get('left', Inches(1.5))
        top = position.get('top', Inches(3.5))
        width = position.get('width', Inches(7))
        height = position.get('height', Inches(1.5))
    
    # Add picture to slide
    slide.shapes.add_picture(img_buf, left, top, width, height)
    
    buf.close()
    img_buf.close()
    
    return presentation

def draw_card_gradient(ax, x, y, w, h, zorder=1):
    gradient = np.linspace(1.0, 0.85, 256).reshape(256, 1)

    img = ax.imshow(
        gradient,
        extent=[x, x + w, y, y + h],
        origin="lower",
        cmap="Greys",
        alpha=0.35,
        zorder=zorder
    )
    return img

def create_financial_goals_visual(presentation, slide_number, goal_names, position=None):
    """
    Clean executive-style goal cards:
    - White rounded cards
    - Left color strip
    - Large size
    """

    import io
    import matplotlib.pyplot as plt
    from matplotlib.patches import FancyBboxPatch
    from PIL import Image
    from pptx.util import Inches

    slide = presentation.slides[slide_number]

    # Professional strip colors (matches screenshot style)
    strip_colors = ['#C62828', '#1A237E', '#6A1B9A', '#2E7D32']

    # ---- SIZING ----
    num_goals = len(goal_names)
    card_height = 1.1
    vertical_stride = 1.4
    font_size = 26
    strip_width = 0.28
    corner_radius = 0.25

    fig_height = 1.4 + num_goals * vertical_stride
    fig, ax = plt.subplots(figsize=(10, fig_height), facecolor='none')
    ax.set_facecolor('none')

    ax.set_xlim(0, 10)
    ax.set_ylim(0, num_goals * vertical_stride + 0.5)
    ax.axis('off')

    card_x = 0.5
    card_width = 9.0

    for i, goal_name in enumerate(goal_names):
        y = (num_goals - i - 1) * vertical_stride + 0.4
        strip_color = strip_colors[i % len(strip_colors)]

        # ---- Shadow (very subtle) ----
        shadow = FancyBboxPatch(
            (card_x + 0.05, y - 0.04),
            card_width,
            card_height,
            boxstyle=f"round,pad=0.02,rounding_size={corner_radius}",
            linewidth=0,
            facecolor='black',
            alpha=0.08,
            zorder=0
        )
        ax.add_patch(shadow)

        # ---- White Card ----
        card = FancyBboxPatch(
            (card_x, y),
            card_width,
            card_height,
            boxstyle=f"round,pad=0.02,rounding_size={corner_radius}",
            linewidth=1.6,
            edgecolor='#D1D5DB',
            facecolor='white',
            zorder=1
        )
        ax.add_patch(card)

        # ---- Left Color Strip ----
        strip = FancyBboxPatch(
            (card_x, y),
            strip_width,
            card_height,
            boxstyle=f"round,pad=0.02,rounding_size={corner_radius}",
            linewidth=0,
            facecolor=strip_color,
            zorder=2
        )
        ax.add_patch(strip)

        # ---- Goal Text ----
        ax.text(
            card_x + strip_width + 0.45,
            y + card_height / 2,
            goal_name,
            ha='left',
            va='center',
            fontsize=font_size,
            fontweight='bold',
            color='#111827',
            zorder=3
        )

    # ---- Export ----
    buf = io.BytesIO()
    plt.tight_layout(pad=0.3)
    plt.savefig(buf, format='png', dpi=180, bbox_inches='tight', transparent=True)
    buf.seek(0)
    plt.close()

    img_buf = io.BytesIO()
    Image.open(buf).save(img_buf, format='PNG')
    img_buf.seek(0)

    # ---- PPT Placement ----
    if position is None:
        left, top = Inches(4.0), Inches(3.1)
        width, height = Inches(5.8), Inches(4.0)
    else:
        left = position.get('left', Inches(9.5))
        top = position.get('top', Inches(4.8))
        width = position.get('width', Inches(5.0))
        height = position.get('height', Inches(1.2))

    slide.shapes.add_picture(img_buf, left, top, width, height)

    buf.close()
    img_buf.close()

    return presentation

from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.oxml.xmlchemy import OxmlElement

def _sub_element(parent, tagname, **kwargs):
           element = OxmlElement(tagname)
           element.attrib.update(kwargs)
           parent.append(element)
           return element

        
def _set_cell_border(cell, border_color="000000", border_width="12700"):
            tc = cell._tc
            tcPr = tc.get_or_add_tcPr()
            for line_tag in ["a:lnL", "a:lnR", "a:lnT", "a:lnB"]:
              ln = _sub_element(tcPr, line_tag, w=border_width, cap="flat", cmpd="sng", algn="ctr")
              solidFill = _sub_element(ln, "a:solidFill")
              _sub_element(solidFill, "a:srgbClr", val=border_color)
              _sub_element(ln, "a:prstDash", val="solid")
              _sub_element(ln, "a:round")
              _sub_element(ln, "a:headEnd", type="none", w="med", len="med")
              _sub_element(ln, "a:tailEnd", type="none", w="med", len="med")
        
def add_table_to_slide(slide,data,left=Inches(1),top=Inches(1),
    width=Inches(6),
    height=Inches(1.5),font_size=Pt(11),
    font_color=RGBColor(0, 0, 0),header_fill=RGBColor(200, 200, 200),
    body_fill=RGBColor(255, 255, 255),border_color="000000",
    border_width="12700", header_font_color=None):
            """Add a bordered table to a slide.

            Args:
            slide: The pptx slide object to add the table to.
            data: 2D list where data[0] is the header row.
            left, top: Position of the table on the slide.
            width, height: Dimensions of the table.
            font_size: Font size for all cells.
            font_color: Font color for all cells.
            header_fill: Background color for the header row.
            body_fill: Background color for body rows.
            border_color: Hex color string for borders (without #).
            border_width: Border width in EMUs (12700 = 1pt).

            Returns:
             The table shape object.
            """
            rows = len(data)
            cols = len(data[0])

            shape = slide.shapes.add_table(rows, cols, left, top, width, height)
            table = shape.table
            table.table_style_id = ""

            for r_idx, row in enumerate(data):
              for c_idx, value in enumerate(row):
               cell = table.cell(r_idx, c_idx)

               # 1. Borders
               _set_cell_border(cell, border_color=border_color, border_width=border_width)

               # 2. Content
               cell.text = str(value)
               font = cell.text_frame.paragraphs[0].font
               font.size = font_size
               font.color.rgb = font_color

               # 3. Fill
               cell.fill.solid()
               if r_idx == 0:
                  cell.fill.fore_color.rgb = header_fill
                  font.bold = True
                  font.color.rgb = header_font_color
               else:
                cell.fill.fore_color.rgb = body_fill

            return shape
