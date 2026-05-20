from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.oxml.xmlchemy import OxmlElement

# --- ENHANCED XML BORDER LOGIC ---
def SubElement(parent, tagname, **kwargs):
    element = OxmlElement(tagname)
    element.attrib.update(kwargs)
    parent.append(element)
    return element

def _set_cell_border(cell, border_color="000000", border_width='12700'):
    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()
    for lines in ['a:lnL','a:lnR','a:lnT','a:lnB']:
        ln = SubElement(tcPr, lines, w=border_width, cap='flat', cmpd='sng', algn='ctr')
        solidFill = SubElement(ln, 'a:solidFill')
        SubElement(solidFill, 'a:srgbClr', val=border_color)
        SubElement(ln, 'a:prstDash', val='solid')
        SubElement(ln, 'a:round')
        SubElement(ln, 'a:headEnd', type='none', w='med', len='med')
        SubElement(ln, 'a:tailEnd', type='none', w='med', len='med')

# 1. Setup
prs = Presentation()
slide = prs.slides.add_slide(prs.slide_layouts[6])

# 2. Table Creation
data = [
    ["Fund Name", "Weight", "Market Cap"],
    ["Canara Robeco Mid Cap", "100%", "Mid Cap"],
    ["HDFC Bank Ltd", "2.72%", "Large Cap"],
    ["ICICI Bank","2.34%","Small Cap"]
]

shape = slide.shapes.add_table(len(data), len(data[0]), Inches(1), Inches(1), Inches(6), Inches(1.5))
table = shape.table
table.table_style_id = "" # Reset style first

# 3. Reordered Loop: Borders -> Content -> Fill
for r_idx, row in enumerate(data):
    for c_idx, value in enumerate(row):
        cell = table.cell(r_idx, c_idx)

        _set_cell_border(cell, border_color="000000", border_width='12700')

        cell.text = str(value)
        font = cell.text_frame.paragraphs[0].font
        font.size = Pt(11)
        font.color.rgb = RGBColor(0, 0, 0)
 
        cell.fill.solid()
        if r_idx == 0:
            cell.fill.fore_color.rgb = RGBColor(200, 200, 200) # Light Grey Header
            font.bold = True
        else:
            cell.fill.fore_color.rgb = RGBColor(255, 255, 255) # White Body

# 4. Save
prs.save('Borders_First_Table.pptx')
print("Table generated: Borders applied before Fill.")